"""Smoke-test the bundled document scripts against generated fixtures.

Proves the scripts in THIS repo actually run and behave as documented (right exit
codes, right verdicts) — reproducibly, in CI or by hand. Generates real .docx/.pptx
fixtures and asserts each outcome. Focused on the Document Template Suite: the
templatize→build→fill→validate engine (incl. the document-property/cover pass and
image/logo slots), the docx/pptx producer validators, and the extractors.

Run from the repo root:
    python skill-builder/scripts/smoke_test_scripts.py

Requires: python-docx, python-pptx, pillow. Exits non-zero on any failure.
"""
import json
import os
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
results = []


def run(script_rel, *args, env=None):
    """Run a bundled script; return (exit_code, parsed_json_or_None)."""
    script = ROOT / script_rel
    if not script.exists():
        return None, None
    run_env = None
    if env:
        run_env = os.environ.copy()
        run_env.update({k: str(v) for k, v in env.items()})
    proc = subprocess.run([sys.executable, str(script), *[str(a) for a in args]],
                          capture_output=True, text=True, env=run_env)
    data = None
    try:
        data = json.loads(proc.stdout)
    except (ValueError, json.JSONDecodeError):
        pass
    return proc.returncode, data


def check(name, cond):
    results.append((name, bool(cond)))
    print(f"  {'PASS' if cond else 'FAIL'}  {name}")


def make_fixtures(tmp):
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image as PILImage

    # docx: bad (leftover tag) and good (heading + para)
    b = Document(); b.add_paragraph("Hi {{ name }}"); b.save(tmp / "bad.docx")
    g = Document(); g.add_heading("Title", 1); g.add_paragraph("ok"); g.save(tmp / "good.docx")

    # docx to templatize: a Label:value line + a bullet (text + list fields)
    t = Document(); t.add_heading("Weekly Update", 1)
    lp = t.add_paragraph(); lp.add_run("Client: ").bold = True; lp.add_run("Globex")
    t.add_paragraph("First achievement", style="List Bullet")
    t.save(tmp / "tmpl_src.docx")

    # docx whose title lives ONLY in a document property (data-bound cover analogue)
    tp = Document(); tp.add_heading("Body Heading", 1)
    tp.core_properties.title = "PROPCODE Report"
    tp.save(tmp / "tmpl_props.docx")

    # docx with an embedded image -> an image SLOT to swap
    PILImage.new("RGB", (120, 60), (10, 20, 200)).save(tmp / "orig_logo.png")
    PILImage.new("RGB", (120, 60), (200, 20, 10)).save(tmp / "new_logo.png")
    ti = Document(); ti.add_heading("Doc With Image", 1)
    ti.add_picture(str(tmp / "orig_logo.png")); ti.save(tmp / "tmpl_img.docx")

    # pptx: bad (lorem) and good
    pb = Presentation(); s = pb.slides.add_slide(pb.slide_layouts[0])
    s.shapes.title.text = "Lorem ipsum"; pb.save(tmp / "bad.pptx")
    pg = Presentation(); sg = pg.slides.add_slide(pg.slide_layouts[0])
    sg.shapes.title.text = "Real Title"; pg.save(tmp / "good.pptx")

    # pptx to templatize: title + TWO independent bullet lists (no interleave)
    pd = Presentation(); sd = pd.slides.add_slide(pd.slide_layouts[1])
    sd.shapes.title.text = "Globex"
    body = sd.placeholders[1].text_frame
    body.text = "Alpha one"; body.add_paragraph().text = "Alpha two"
    box = sd.shapes.add_textbox(Inches(1), Inches(4), Inches(5), Inches(2)).text_frame
    box.text = "Beta one"; box.add_paragraph().text = "Beta two"
    pd.save(tmp / "tmpl_deck.pptx")


def main():
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        make_fixtures(tmp)

        print("processing-word-documents/extract_docx.py")
        rc, d = run("processing-word-documents/scripts/extract_docx.py", tmp / "good.docx")
        check("docx extractor returns markdown", rc == 0 and d and "Title" in d["markdown"])

        print("authoring-word-documents/validate_docx.py")
        rc, d = run("authoring-word-documents/scripts/validate_docx.py", tmp / "bad.docx")
        check("docx validator fails on unfilled tag", rc == 1 and d and d["status"] == "FAILED")
        rc, d = run("authoring-word-documents/scripts/validate_docx.py", tmp / "good.docx")
        check("docx validator passes clean doc", rc == 0 and d and d["status"] == "OK")

        print("processing-powerpoint-files/extract_pptx.py")
        rc, d = run("processing-powerpoint-files/scripts/extract_pptx.py", tmp / "good.pptx")
        check("pptx extractor returns slides", rc == 0 and d and d["slides"][0]["title"] == "Real Title")

        print("building-powerpoint-decks/validate_pptx.py")
        rc, d = run("building-powerpoint-decks/scripts/validate_pptx.py", tmp / "bad.pptx")
        check("pptx validator fails on lorem", rc == 1 and d and d["status"] == "FAILED")
        rc, d = run("building-powerpoint-decks/scripts/validate_pptx.py", tmp / "good.pptx")
        check("pptx validator passes clean deck", rc == 0 and d and d["status"] == "OK")

        print("processing-documents/detect_type.py")
        rc, d = run("processing-documents/scripts/detect_type.py", tmp / "good.docx")
        check("type detector identifies docx", rc == 0 and d and d["detected_type"] == "docx")

        print("building-document-templates: templatize -> build -> fill -> validate")
        tdir = "building-document-templates/scripts"
        reg = tmp / "registry"
        env = {"TEMPLATE_REGISTRY": reg}
        rc, _ = run(f"{tdir}/templatize.py", "propose", "--file", tmp / "tmpl_src.docx",
                    "--out", tmp / "proposal.json")
        prop = json.loads((tmp / "proposal.json").read_text(encoding="utf-8")) if rc == 0 else {}
        names = {c["current_text"] for c in prop.get("candidates", [])}
        check("templatize propose extracts candidates", rc == 0 and "Globex" in names)
        (tmp / "reviewed.json").write_text(json.dumps({
            "format": "docx", "source_file": "tmpl_src.docx", "candidates": [
                {"current_text": "Globex", "suggest_name": "client_name",
                 "suggest_type": "text", "keep": "variable"},
                {"current_text": "First achievement", "suggest_name": "achievements",
                 "suggest_type": "list", "keep": "variable"},
            ]}), encoding="utf-8")
        rc, _ = run(f"{tdir}/templatize.py", "build", "--file", tmp / "tmpl_src.docx",
                    "--fields", tmp / "reviewed.json", "--client", "acme",
                    "--doc-type", "weekly", "--created", "2026-07-13", env=env)
        man = reg / "acme" / "weekly" / "manifest.json"
        check("templatize build registers template + manifest", rc == 0 and man.exists())
        (tmp / "content.json").write_text(json.dumps({
            "client_name": "Initech", "achievements": ["Alpha", "Beta", "Gamma"]}),
            encoding="utf-8")
        out_docx = tmp / "filled.docx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "weekly",
                    "--data", tmp / "content.json", "--out", out_docx, env=env)
        check("fill produces output document", rc == 0 and out_docx.exists())
        rc, d = run(f"{tdir}/validate.py", out_docx)
        check("validate passes a fully filled document", rc == 0 and d and d["status"] == "OK")
        rc, d = run(f"{tdir}/validate.py", tmp / "bad.docx")
        check("validate fails on a leftover template tag", rc == 1 and d and d["status"] == "FAIL")
        from docx import Document as _Doc
        filled_texts = [p.text for p in _Doc(str(out_docx)).paragraphs]
        check("fill expands a list field into separate items",
              all(x in filled_texts for x in ("Alpha", "Beta", "Gamma")))

        # missing REQUIRED field -> tag stays, fill exits non-zero, validate FAILs
        (tmp / "partial.json").write_text(json.dumps({"achievements": ["only this"]}),
                                          encoding="utf-8")
        out_partial = tmp / "partial.docx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "weekly",
                    "--data", tmp / "partial.json", "--out", out_partial, env=env)
        check("fill exits non-zero when a required field is missing",
              rc != 0 and out_partial.exists())
        rc, d = run(f"{tdir}/validate.py", out_partial)
        check("validate fails when a required field was left unfilled",
              rc == 1 and d and d["status"] == "FAIL")

        # document-property / data-bound-cover pass
        rc, _ = run(f"{tdir}/templatize.py", "propose", "--file", tmp / "tmpl_props.docx",
                    "--out", tmp / "props_prop.json")
        pp = json.loads((tmp / "props_prop.json").read_text(encoding="utf-8")) if rc == 0 else {}
        has_prop = any(c.get("source") == "property" and c["current_text"] == "PROPCODE Report"
                       for c in pp.get("candidates", []))
        check("propose surfaces a document-property (cover) candidate", rc == 0 and has_prop)
        (tmp / "props_reviewed.json").write_text(json.dumps({
            "format": "docx", "source_file": "tmpl_props.docx", "candidates": [
                {"current_text": "PROPCODE Report", "suggest_name": "doc_title",
                 "suggest_type": "text", "keep": "variable", "source": "property"},
            ]}), encoding="utf-8")
        rc, _ = run(f"{tdir}/templatize.py", "build", "--file", tmp / "tmpl_props.docx",
                    "--fields", tmp / "props_reviewed.json", "--client", "acme",
                    "--doc-type", "propdoc", "--created", "2026-07-13", env=env)
        tmpl_core = zipfile.ZipFile(reg / "acme" / "propdoc" / "template.docx").read("docProps/core.xml").decode()
        check("build injects the tag into docProps/core.xml", rc == 0 and "{{ doc_title }}" in tmpl_core)
        (tmp / "props_content.json").write_text(json.dumps({"doc_title": "NEWCODE Report"}),
                                                encoding="utf-8")
        out_props = tmp / "props_filled.docx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "propdoc",
                    "--data", tmp / "props_content.json", "--out", out_props, env=env)
        filled_core = zipfile.ZipFile(out_props).read("docProps/core.xml").decode()
        check("fill replaces the property tag with the value",
              rc == 0 and "NEWCODE Report" in filled_core and "{{ doc_title }}" not in filled_core)

        # image/logo slot: templatize an embedded image, fill with a replacement
        rc, _ = run(f"{tdir}/templatize.py", "propose", "--file", tmp / "tmpl_img.docx",
                    "--out", tmp / "img_prop.json")
        ip = json.loads((tmp / "img_prop.json").read_text(encoding="utf-8")) if rc == 0 else {}
        imgc = [c for c in ip.get("candidates", []) if c.get("source") == "image"]
        check("propose surfaces an image slot", rc == 0 and len(imgc) >= 1)
        if imgc:
            media = imgc[0]["media_part"]
            (tmp / "img_reviewed.json").write_text(json.dumps({
                "format": "docx", "source_file": "tmpl_img.docx", "candidates": [
                    {"current_text": imgc[0]["current_text"], "suggest_name": "logo",
                     "suggest_type": "image", "keep": "variable", "source": "image",
                     "media_part": media},
                ]}), encoding="utf-8")
            rc, _ = run(f"{tdir}/templatize.py", "build", "--file", tmp / "tmpl_img.docx",
                        "--fields", tmp / "img_reviewed.json", "--client", "acme",
                        "--doc-type", "imgdoc", "--created", "2026-07-13", env=env)
            orig = zipfile.ZipFile(reg / "acme" / "imgdoc" / "template.docx").read(media)
            (tmp / "img_content.json").write_text(
                json.dumps({"logo": str(tmp / "new_logo.png")}), encoding="utf-8")
            out_img = tmp / "img_filled.docx"
            rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "imgdoc",
                        "--data", tmp / "img_content.json", "--out", out_img, env=env)
            swapped = zipfile.ZipFile(out_img).read(media)
            check("fill swaps the image-slot media bytes",
                  rc == 0 and swapped != orig and len(swapped) > 0)

        # PPTX round-trip with TWO independent list fields
        rc, _ = run(f"{tdir}/templatize.py", "propose", "--file", tmp / "tmpl_deck.pptx",
                    "--out", tmp / "deck_prop.json")
        (tmp / "deck_reviewed.json").write_text(json.dumps({
            "format": "pptx", "source_file": "tmpl_deck.pptx", "candidates": [
                {"current_text": "Globex", "suggest_name": "client_name",
                 "suggest_type": "text", "keep": "variable"},
                {"current_text": "Alpha one", "suggest_name": "list_a",
                 "suggest_type": "list", "keep": "variable"},
                {"current_text": "Alpha two", "keep": "remove"},
                {"current_text": "Beta one", "suggest_name": "list_b",
                 "suggest_type": "list", "keep": "variable"},
                {"current_text": "Beta two", "keep": "remove"},
            ]}), encoding="utf-8")
        rc, _ = run(f"{tdir}/templatize.py", "build", "--file", tmp / "tmpl_deck.pptx",
                    "--fields", tmp / "deck_reviewed.json", "--client", "acme",
                    "--doc-type", "deck", "--created", "2026-07-13", env=env)
        deck_man = reg / "acme" / "deck" / "manifest.json"
        check("pptx templatize build registers template + manifest",
              rc == 0 and deck_man.exists())
        (tmp / "deck_content.json").write_text(json.dumps({
            "client_name": "Initech", "list_a": ["A1", "A2", "A3"],
            "list_b": ["B1", "B2"]}), encoding="utf-8")
        out_pptx = tmp / "deck_filled.pptx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "deck",
                    "--data", tmp / "deck_content.json", "--out", out_pptx, env=env)
        check("pptx fill produces output deck", rc == 0 and out_pptx.exists())
        rc, d = run(f"{tdir}/validate.py", out_pptx)
        check("pptx validate passes a fully filled deck",
              rc == 0 and d and d["status"] == "OK")
        from pptx import Presentation as _Prs
        deck_texts = [p.text for sl in _Prs(str(out_pptx)).slides
                      for sh in sl.shapes if sh.has_text_frame
                      for p in sh.text_frame.paragraphs]
        check("pptx fill expands two independent list fields",
              all(x in deck_texts for x in ("A1", "A2", "A3", "B1", "B2")))

        print("presenting-with-html: build_html -> validate")
        hdir = "presenting-with-html/scripts"
        built_deck = tmp / "built_deck.html"
        rc, _ = run(f"{hdir}/build_html.py", "--content",
                    ROOT / "presenting-with-html" / "examples" / "deck-content.json",
                    "--out", built_deck)
        check("build_html builds the deck fixture", rc == 0 and built_deck.exists())
        rc, d = run(f"{hdir}/validate_html.py", built_deck)
        check("built deck passes the validator", rc == 0 and d and d["status"] == "OK")
        built_report = tmp / "built_report.html"
        rc, _ = run(f"{hdir}/build_html.py", "--content",
                    ROOT / "presenting-with-html" / "examples" / "report-content.json",
                    "--out", built_report)
        check("build_html builds the report fixture", rc == 0 and built_report.exists())
        rc, d = run(f"{hdir}/validate_html.py", built_report)
        check("built report passes the validator", rc == 0 and d and d["status"] == "OK")

        # second style preset builds and validates
        built_clean = tmp / "built_clean.html"
        rc, _ = run(f"{hdir}/build_html.py", "--content",
                    ROOT / "presenting-with-html" / "examples" / "deck-content.json",
                    "--out", built_clean, "--style", "clean")
        clean_html = built_clean.read_text(encoding="utf-8") if built_clean.exists() else ""
        check("build_html builds the clean style preset",
              rc == 0 and 'data-style="clean"' in clean_html and 'data-theme="light"' in clean_html)
        rc, d = run(f"{hdir}/validate_html.py", built_clean)
        check("clean-style deck passes the validator", rc == 0 and d and d["status"] == "OK")

        # bad content -> precise, non-zero failure (not a broken build)
        (tmp / "bad_content.json").write_text(json.dumps({
            "format": "deck", "meta": {"title": "X"},
            "slides": [{"type": "nonsense"}]}), encoding="utf-8")
        rc, _ = run(f"{hdir}/build_html.py", "--content", tmp / "bad_content.json",
                    "--out", tmp / "never.html")
        check("build_html rejects an unknown block type", rc != 0)

        # pre-build content check + did-you-mean on wrong field names
        rc, _ = run(f"{hdir}/build_html.py", "--content",
                    ROOT / "presenting-with-html" / "examples" / "deck-content.json",
                    "--validate-only")
        check("build_html --validate-only passes the deck fixture", rc == 0)
        (tmp / "alias_content.json").write_text(json.dumps({
            "format": "deck", "meta": {"title": "X"},
            "slides": [{"type": "bullets", "title": "H", "bullets": ["a"]}]}),
            encoding="utf-8")
        proc = subprocess.run(
            [sys.executable, str(ROOT / hdir / "build_html.py"),
             "--content", str(tmp / "alias_content.json"), "--validate-only"],
            capture_output=True, text=True)
        check("build_html suggests exact field names (did-you-mean)",
              proc.returncode != 0 and 'did you mean "items"' in proc.stdout
              and 'did you mean "heading"' in proc.stdout)

        # partial checkout: the skill folder alone (no repo-root brands/) still builds
        import shutil
        iso = tmp / "iso" / "presenting-with-html"
        shutil.copytree(ROOT / "presenting-with-html" / "scripts", iso / "scripts")
        shutil.copytree(ROOT / "presenting-with-html" / "assets", iso / "assets")
        proc = subprocess.run(
            [sys.executable, str(iso / "scripts" / "build_html.py"),
             "--content", str(ROOT / "presenting-with-html" / "examples" / "deck-content.json"),
             "--out", str(tmp / "iso_deck.html")],
            capture_output=True, text=True)
        check("build_html falls back to embedded brand without repo brands/",
              proc.returncode == 0 and (tmp / "iso_deck.html").exists()
              and "embedded" in proc.stdout)

        # QA hooks ship in the built output (?theme= both formats, ?slide= deck)
        deck_text = built_deck.read_text(encoding="utf-8")
        check("built output carries the ?theme= QA hook", "URLSearchParams" in deck_text)
        check("built deck carries the ?slide= QA hook", "get('slide')" in deck_text)

        # validator slide count is token-exact: slide-inner / slide-count must not
        # inflate it (field test: 8-slide deck reported as 17)
        fixture = json.loads((ROOT / "presenting-with-html" / "examples" /
                              "deck-content.json").read_text(encoding="utf-8"))
        rc, d = run(f"{hdir}/validate_html.py", built_deck)
        check("validator n_slides == content slides + auto title slide",
              rc == 0 and d and d["checks"]["n_slides"] == len(fixture["slides"]) + 1)
        # and the deck fixture exercises the status block (RAG pills)
        check("deck fixture includes a status block",
              any(s.get("type") == "status" for s in fixture["slides"]))
        # build prints an explicit success line
        proc = subprocess.run(
            [sys.executable, str(ROOT / hdir / "build_html.py"),
             "--content", str(ROOT / "presenting-with-html" / "examples" / "deck-content.json"),
             "--out", str(tmp / "ok_msg.html"), "--lite"],
            capture_output=True, text=True)
        check("build_html prints an explicit BUILD OK line",
              proc.returncode == 0 and "BUILD OK" in proc.stdout)

        # scaffold emits a parseable skeleton with exact field names
        rc, scaffold = run(f"{hdir}/build_html.py", "--scaffold", "report")
        check("build_html --scaffold emits a valid report skeleton",
              rc == 0 and scaffold and scaffold.get("format") == "report"
              and isinstance(scaffold.get("sections"), list)
              and any(b.get("type") == "kpi" for b in scaffold["sections"]))

        # safe mini-markup: whitelist converts, everything else stays escaped
        (tmp / "rich_content.json").write_text(json.dumps({
            "format": "report", "meta": {"title": "Rich"},
            "sections": [
                {"type": "text", "heading": "A", "paragraphs":
                    ["**bold** and *em* and `code` and [ok](https://x.y)"]},
                {"type": "text", "heading": "B", "paragraphs":
                    ["<script>alert(1)</script>", "[bad](javascript:alert(1))"]},
            ]}), encoding="utf-8")
        rich_out = tmp / "rich.html"
        rc, _ = run(f"{hdir}/build_html.py", "--content", tmp / "rich_content.json",
                    "--out", rich_out)
        rich_html = rich_out.read_text(encoding="utf-8") if rich_out.exists() else ""
        check("mini-markup converts the four safe patterns",
              rc == 0 and "<strong>bold</strong>" in rich_html
              and "<em>em</em>" in rich_html and "<code>code</code>" in rich_html
              and '<a href="https://x.y">ok</a>' in rich_html)
        check("mini-markup keeps raw HTML escaped and rejects unsafe URLs",
              "&lt;script&gt;alert(1)&lt;/script&gt;" in rich_html
              and "[bad](javascript:alert(1))" in rich_html
              and "<script>alert(1)" not in rich_html)
        rc, d = run(f"{hdir}/validate_html.py", rich_out)
        check("mini-markup output passes the validator",
              rc == 0 and d and d["status"] == "OK")

        # office components: all 8 new blocks build and validate in one showcase
        built_office = tmp / "office.html"
        rc, _ = run(f"{hdir}/build_html.py", "--content",
                    ROOT / "presenting-with-html" / "examples" / "office-components.json",
                    "--out", built_office)
        office_html = built_office.read_text(encoding="utf-8") if built_office.exists() else ""
        rc2, d = run(f"{hdir}/validate_html.py", built_office)
        check("office blocks (agenda/callout/team/status/contact/steps/feature/definitions) build + validate",
              rc == 0 and rc2 == 0 and d and d["status"] == "OK"
              and not d["checks"].get("undefined_classes")
              and all(c in office_html for c in
                      ('class="agenda"', 'class="callout recommendation"', 'class="team-grid"',
                       'class="rag amber"', 'class="steps"', 'class="defs"')))

        # self-correcting content errors + warnings
        (tmp / "warn_content.json").write_text(json.dumps({
            "format": "deck", "meta": {"title": "T", "title_accent": "Nope"},
            "blocks": [{"type": "bullets", "heading": "H", "items": ["a"]}]}), encoding="utf-8")
        proc = subprocess.run(
            [sys.executable, str(ROOT / hdir / "build_html.py"),
             "--content", str(tmp / "warn_content.json"), "--validate-only"],
            capture_output=True, text=True)
        check('"blocks" suggests the right list name for the format',
              proc.returncode != 0 and 'did you mean "slides"' in proc.stdout)
        (tmp / "warn_content.json").write_text(json.dumps({
            "format": "deck", "meta": {"title": "T", "title_accent": "Nope"},
            "slides": [{"type": "bullets", "heading": "H",
                        "items": [f"item {i}" for i in range(12)]}]}), encoding="utf-8")
        proc = subprocess.run(
            [sys.executable, str(ROOT / hdir / "build_html.py"),
             "--content", str(tmp / "warn_content.json"), "--validate-only"],
            capture_output=True, text=True)
        check("title_accent + overflow warnings fire without failing the build",
              proc.returncode == 0 and "accent will not render" in proc.stdout
              and "split this slide" in proc.stdout)

        # minimal scaffold + passthrough-chart warning
        rc, mini = run(f"{hdir}/build_html.py", "--scaffold", "report", "--minimal")
        check("--scaffold --minimal emits the smallest skeleton",
              rc == 0 and mini and len(mini.get("sections", [])) == 2)
        (tmp / "pt_content.json").write_text(json.dumps({
            "format": "report", "meta": {"title": "T"},
            "sections": [{"type": "chart", "heading": "H",
                          "chart": {"plotly": {"data": [{"type": "bar", "x": [1], "y": [2]}]}}},
                         {"type": "text", "heading": "X", "paragraphs": ["p"]}]}), encoding="utf-8")
        proc = subprocess.run(
            [sys.executable, str(ROOT / hdir / "build_html.py"),
             "--content", str(tmp / "pt_content.json"), "--validate-only"],
            capture_output=True, text=True)
        check("raw plotly passthrough triggers the no-theme-restyle warning",
              proc.returncode == 0 and "NOT restyle" in proc.stdout)

        # third style preset builds and validates
        built_exec = tmp / "built_exec.html"
        rc, _ = run(f"{hdir}/build_html.py", "--content",
                    ROOT / "presenting-with-html" / "examples" / "report-content.json",
                    "--out", built_exec, "--style", "executive")
        exec_html = built_exec.read_text(encoding="utf-8") if built_exec.exists() else ""
        check("build_html builds the executive style preset",
              rc == 0 and 'data-style="executive"' in exec_html
              and 'data-theme="light"' in exec_html)
        rc, d = run(f"{hdir}/validate_html.py", built_exec)
        check("executive-style report passes the validator",
              rc == 0 and d and d["status"] == "OK")

        print("presenting-with-html/validate_html.py (integrity hardening)")
        # The bespoke boilerplates (which carried data-sample markers) were retired;
        # assert the data-sample residue gate directly by injecting the marker into a
        # real built deck.
        sample_src = built_deck.read_text(encoding="utf-8").replace(
            '<div class="stage"', '<div data-sample class="stage"', 1)
        sample_file = tmp / "sample.html"
        sample_file.write_text(sample_src, encoding="utf-8")
        rc, d = run(f"{hdir}/validate_html.py", sample_file)
        check("validator FAILS a page with data-sample residue",
              rc == 1 and d and d["status"] == "FAIL" and d["checks"].get("sample_blocks", 0) > 0)
        # Regression: the observed small-model failure — a second document appended
        # after </html> with duplicate IDs must hard-fail.
        deck_html = built_deck.read_text(encoding="utf-8")
        concat = deck_html + ('<div id="deck"><section class="slide-header">'
                              '<h1 class="slide-title">Orphan</h1></section>'
                              '<button id="prevBtn"></button></div>')
        (tmp / "concat.html").write_text(concat, encoding="utf-8")
        rc, d = run(f"{hdir}/validate_html.py", tmp / "concat.html")
        check("validator fails content appended after </html> + duplicate IDs",
              rc == 1 and d and not d["checks"]["nothing_after_html"]
              and d["checks"]["duplicate_ids"])
        bad_html = tmp / "bad_report.html"
        bad_html.write_text("<html><body><h1>Report {{ title }}</h1>"
                            "<p>lorem ipsum</p></body></html>", encoding="utf-8")
        rc, d = run("presenting-with-html/scripts/validate_html.py", bad_html)
        check("html validator fails a non-deck / placeholder page",
              rc == 1 and d and d["status"] == "FAIL")

        # false-positive regressions (field-tested): TODO/FIXME inside a vendored-size
        # script must not trip the placeholder scan, and overflow:hidden on a non-body
        # element must not raise the scroll-lock warning.
        report_text = built_report.read_text(encoding="utf-8")
        vendored = report_text.replace(
            "</body>", "<script>/* FIXME TODO */" + "x" * 60000 + "</script></body>", 1)
        (tmp / "vendored.html").write_text(vendored, encoding="utf-8")
        rc, d = run(f"{hdir}/validate_html.py", tmp / "vendored.html")
        check("validator ignores TODO/FIXME inside vendored-size scripts",
              rc == 0 and d and d["status"] == "OK"
              and d["checks"].get("vendored_scripts_skipped") == 1)
        check("overflow warning scoped: clipped figures don't flag",
              d and d["checks"].get("no_horizontal_overflow") is True)
        locked = report_text.replace("</style>", "body{overflow:hidden}</style>", 1)
        (tmp / "locked.html").write_text(locked, encoding="utf-8")
        rc, d = run(f"{hdir}/validate_html.py", tmp / "locked.html")
        check("validator still warns when html/body itself is scroll-locked",
              d and d["checks"].get("no_horizontal_overflow") is False)

        print("built-in template libraries: generate -> scaffold -> fill -> validate")
        rc, _ = run("building-powerpoint-decks/scripts/build_template_library.py",
                    "--only", "exec_update", "--registry", reg, "--created", "2026-07-13")
        pman = reg / "_builtin" / "exec_update" / "manifest.json"
        check("pptx library builder registers exec_update", rc == 0 and pman.exists())
        rc, _ = run("authoring-word-documents/scripts/build_doc_library.py",
                    "--only", "memo", "--registry", reg, "--created", "2026-07-13")
        dman = reg / "_builtin" / "memo" / "manifest.json"
        check("docx library builder registers memo", rc == 0 and dman.exists())

        # slide groups (the family/client-template engine): a dedicated mini
        # fixture — one repeatable topic slide, one optional slide that drops.
        from pptx.util import Inches as _In
        gp = _Prs()
        gp.slide_width, gp.slide_height = 12192000, 6858000
        titles = ("{{ deck_title }}", "{{ t_head }}", "{{ o_head }}", "{{ closing }}")
        for t in titles:
            s = gp.slides.add_slide(gp.slide_layouts[5])
            s.shapes.title.text = t
        gp.slides[1].shapes.add_textbox(_In(1), _In(2), _In(8), _In(2)) \
            .text_frame.text = "{{ t_points }}"
        gdir = reg / "_builtin" / "mini_groups"
        gdir.mkdir(parents=True, exist_ok=True)
        gp.save(gdir / "template.pptx")
        gman = {"template_id": "_builtin/mini_groups", "client": "_builtin",
                "doc_type": "mini_groups", "format": "pptx",
                "template_file": "template.pptx", "version": "1.0.0",
                "description": "slide-group engine fixture",
                "fields": [{"name": "deck_title", "type": "text", "example": "T",
                            "guidance": "", "required": True},
                           {"name": "closing", "type": "text", "example": "C",
                            "guidance": "", "required": True}],
                "row_groups": [], "source_terms": [],
                "slide_groups": [
                    {"name": "topic_slides", "slide_index": 1, "min": 1, "max": 4,
                     "purpose": "one topic per slide",
                     "fields": [{"name": "t_head", "type": "text", "example": "H",
                                 "guidance": "", "required": True},
                                {"name": "t_points", "type": "list", "example": "P",
                                 "guidance": "", "required": True}]},
                    {"name": "opt_slides", "slide_index": 2, "min": 0, "max": 2,
                     "purpose": "optional slide",
                     "fields": [{"name": "o_head", "type": "text", "example": "O",
                                 "guidance": "", "required": True}]}]}
        (gdir / "manifest.json").write_text(json.dumps(gman), encoding="utf-8")
        rc, _ = run(f"{tdir}/registry.py", "scaffold", "--builtin", "mini_groups",
                    "--out", tmp / "sg_content.json", "--with-examples", env=env)
        sg = json.loads((tmp / "sg_content.json").read_text(encoding="utf-8"))
        check("scaffold emits slide-group entry lists",
              rc == 0 and isinstance(sg.get("topic_slides"), list)
              and isinstance(sg.get("opt_slides"), list))
        sg = {"deck_title": "Initech", "closing": "Done", "opt_slides": [],
              "topic_slides": [{"t_head": f"Topic {i}", "t_points": [f"P{i}a", f"P{i}b"]}
                               for i in (1, 2, 3)]}
        (tmp / "sg_filled.json").write_text(json.dumps(sg), encoding="utf-8")
        out_sg = tmp / "sg_deck.pptx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "_builtin", "--doc-type", "mini_groups",
                    "--data", tmp / "sg_filled.json", "--out", out_sg, env=env)
        rc2, d = run(f"{tdir}/validate.py", out_sg, "--template",
                     gdir / "template.pptx", "--manifest", gdir / "manifest.json")
        sg_texts = [p.text for sl in _Prs(str(out_sg)).slides for sh in sl.shapes
                    if sh.has_text_frame for p in sh.text_frame.paragraphs]
        check("slide group expands per entry and drops the optional slide",
              rc == 0 and rc2 == 0 and d and d["status"] == "OK"
              and d["checks"]["n_slides"] == 5      # title + 3 topics + closing
              and all(f"Topic {i}" in sg_texts for i in (1, 2, 3)))
        # delete+clone must not collide on slide partnames (duplicate zip entries)
        from collections import Counter
        sg_names = Counter(zipfile.ZipFile(out_sg).namelist())
        check("slide clones get unique part names (no duplicate zip entries)",
              not [n for n, c in sg_names.items() if c > 1])
        # library bullets carry a hanging indent — without marL/indent the glyph
        # renders flush against the text ("•Like this"; field-test finding)
        with zipfile.ZipFile(reg / "_builtin" / "exec_update" / "template.pptx") as ztpl:
            slide_xml = b"".join(ztpl.read(n) for n in ztpl.namelist()
                                 if n.startswith("ppt/slides/slide"))
        check("library bullets set a hanging indent (marL/indent)",
              b"buChar" in slide_xml and b'marL="228600"' in slide_xml
              and b'indent="-228600"' in slide_xml)
        # scaffold warns against inventing facts (field test: Haiku fabricated
        # owners/dates/tool names to fill required fields)
        proc = subprocess.run(
            [sys.executable, str(ROOT / tdir / "registry.py"), "scaffold",
             "--builtin", "exec_update", "--out", str(tmp / "sg_rule.json")],
            capture_output=True, text=True,
            env={**os.environ, "TEMPLATE_REGISTRY": str(reg)})
        check("scaffold prints the no-invented-facts rule",
              proc.returncode == 0 and "ONLY with facts the user supplied" in proc.stdout)
        sg_missing = {k: v for k, v in sg.items() if k != "topic_slides"}
        (tmp / "sg_missing.json").write_text(json.dumps(sg_missing), encoding="utf-8")
        rc, _ = run(f"{tdir}/fill.py", "--client", "_builtin", "--doc-type", "mini_groups",
                    "--data", tmp / "sg_missing.json", "--out", tmp / "sg_missing.pptx", env=env)
        check("missing required slide group fails the fill", rc != 0)

        # flex_deck: composable typed body — any mix/order, native charts+tables,
        # item-row cloning, placeholder-visual gate.
        rc, _ = run("building-powerpoint-decks/scripts/build_template_library.py",
                    "--only", "flex_deck", "--registry", reg, "--created", "2026-07-16")
        fman = reg / "_builtin" / "flex_deck" / "manifest.json"
        check("library builder registers flex_deck (composable body)", rc == 0 and fman.exists())
        rc, _ = run(f"{tdir}/registry.py", "scaffold", "--builtin", "flex_deck",
                    "--out", tmp / "flex_scaffold.json", "--with-examples", env=env)
        fs = json.loads((tmp / "flex_scaffold.json").read_text(encoding="utf-8"))
        check("flex scaffold emits an ordered typed body list",
              rc == 0 and isinstance(fs.get("body"), list) and
              all("type" in e for e in fs["body"]))
        # next_steps deliberately OMITTED (optional empty list must leave a valid
        # txBody) and an image entry included (swap must drop the placeholder rel).
        flex_data = {
            "deck_eyebrow": "REVIEW", "deck_title": "Initech Q3", "deck_subtitle": "Sub.",
            "author_line": "Ops", "closing_statement": "Approve the plan",
            "body": [
                {"type": "agenda", "heading": "Cover",
                 "items": [{"title": "One", "text": "A"}, {"title": "Two"}]},
                {"type": "chart", "heading": "Grew",
                 "chart": {"chart_type": "column", "categories": ["A", "B"],
                           "series": [{"name": "S", "values": [1, 2]}]}},
                {"type": "table", "heading": "Costs",
                 "table": {"columns": ["Site", "Cost"], "rows": [["X", "1"], ["Y", "2"]]}},
                {"type": "team", "heading": "Owners",
                 "items": [{"initials": "AB", "name": "Alpha", "role": "Lead"}]},
                {"type": "image", "heading": "Pic", "image": str(tmp / "orig_logo.png"),
                 "caption": "cap"},
                {"type": "quote", "quote": "It works."},
            ],
        }
        (tmp / "flex_data.json").write_text(json.dumps(flex_data), encoding="utf-8")
        out_flex = tmp / "flex.pptx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "_builtin", "--doc-type", "flex_deck",
                    "--data", tmp / "flex_data.json", "--out", out_flex, env=env)
        rc2, d = run(f"{tdir}/validate.py", out_flex, "--template",
                     reg / "_builtin" / "flex_deck" / "template.pptx", "--manifest", fman)
        flex_names = zipfile.ZipFile(out_flex).namelist()
        check("flex body fills 6 typed slides + native chart + table + image, validates",
              rc == 0 and rc2 == 0 and d and d["status"] == "OK"
              and d["checks"]["n_slides"] == 8          # cover + 6 body + closing
              and not d["checks"]["placeholder_visuals"]   # swapped image drops old rel
              and any(n.startswith("ppt/charts/chart") for n in flex_names))
        # Empty optional list must never leave a paragraph-less <p:txBody> —
        # schema-invalid; PowerPoint refuses to open the file (field-test find).
        import re as _re
        with zipfile.ZipFile(out_flex) as zf:
            txbody_ok = all(
                "<a:p" in seg
                for n in flex_names if _re.fullmatch(r"ppt/slides/slide\d+\.xml", n)
                for seg in _re.findall(r"<p:txBody>.*?</p:txBody>",
                                       zf.read(n).decode("utf-8", "ignore"), _re.S))
        check("empty optional list leaves a schema-valid text body", txbody_ok)
        # unswapped placeholder visual must FAIL validation (warehouse-audit bug)
        rc, d = run(f"{tdir}/validate.py", reg / "_builtin" / "flex_deck" / "template.pptx",
                    "--template", reg / "_builtin" / "flex_deck" / "template.pptx",
                    "--manifest", fman)
        check("placeholder visual left in a deck fails validation",
              rc == 1 and d and d["checks"].get("placeholder_visuals"))
        # image body entry without a file -> fill refuses
        bad = dict(flex_data)
        bad["body"] = [{"type": "image", "heading": "Pic", "caption": "c"}]
        (tmp / "flex_bad.json").write_text(json.dumps(bad), encoding="utf-8")
        rc, _ = run(f"{tdir}/fill.py", "--client", "_builtin", "--doc-type", "flex_deck",
                    "--data", tmp / "flex_bad.json", "--out", tmp / "flex_bad.pptx", env=env)
        check("image body slide without an image fails the fill", rc != 0)

        rc, _ = run(f"{tdir}/registry.py", "scaffold", "--builtin", "memo",
                    "--out", tmp / "memo_content.json", "--with-examples", env=env)
        check("registry scaffold emits a content file", rc == 0
              and (tmp / "memo_content.json").exists())
        # Unedited examples must FAIL the source-residue gate...
        out_memo = tmp / "memo_unedited.docx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "_builtin", "--doc-type", "memo",
                    "--data", tmp / "memo_content.json", "--out", out_memo, env=env)
        rc, d = run(f"{tdir}/validate.py", out_memo, "--template",
                    reg / "_builtin" / "memo" / "template.docx",
                    "--manifest", dman)
        check("unedited example fill fails the source-residue gate",
              rc == 1 and d and d["status"] == "FAIL")
        # ...and an edited fill must pass.
        memo_content = json.loads((tmp / "memo_content.json").read_text(encoding="utf-8"))
        edited = json.loads(json.dumps(memo_content)
                            .replace("Acme Mining", "Initech")
                            .replace("Jane Mokoena", "Pat Lee")
                            .replace("Sipho Dlamini", "Sam Cole"))
        (tmp / "memo_edited.json").write_text(json.dumps(edited), encoding="utf-8")
        out_memo2 = tmp / "memo_edited.docx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "_builtin", "--doc-type", "memo",
                    "--data", tmp / "memo_edited.json", "--out", out_memo2, env=env)
        rc, d = run(f"{tdir}/validate.py", out_memo2, "--template",
                    reg / "_builtin" / "memo" / "template.docx",
                    "--manifest", dman)
        check("edited builtin fill passes validation",
              rc == 0 and d and d["status"] == "OK")

    passed = sum(1 for _, ok in results if ok)
    print(f"\n{passed}/{len(results)} checks passed")
    sys.exit(0 if passed == len(results) else 1)


if __name__ == "__main__":
    main()
