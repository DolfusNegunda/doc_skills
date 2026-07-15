"""Generate the starter deck the skill adapts — the .pptx analog of deck-template.html.

The presenting-with-html skill's spine is "adapt the boilerplate, don't reinvent it":
it ships a complete, styled deck you copy and refill. PowerPoint decks deserve the
same. This builds assets/starter-template.pptx: a 16:9 deck with a neutral, premium,
BRAND-SWAPPABLE look and one example of each core layout (title, section divider,
content, two-column, chart/evidence, closing). Copy it, replace the content, keep the
geometry, palette, and type scale.

Brand swap (the org's "branding is data, not code" rule): change PALETTE and FONT
below and re-run — every slide re-skins from these values, exactly like swapping the
--accent tokens in the HTML template. Or drop the org's real .potx over the theme.

    python scripts/build_starter_template.py            # -> assets/starter-template.pptx
    python scripts/build_starter_template.py --out x.pptx

Type scale respects deck-anatomy minimums (title >=32, body >=24, never <18).
Requires python-pptx.
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- Brand tokens (swap these to re-skin the whole deck) ----
PALETTE = {
    "bg":      RGBColor(0xFF, 0xFF, 0xFF),  # slide background
    "ink":     RGBColor(0x0F, 0x17, 0x2A),  # primary text (deep navy)
    "muted":   RGBColor(0x54, 0x61, 0x7A),  # secondary text
    "accent":  RGBColor(0x25, 0x63, 0xEB),  # primary accent (blue)
    "accent2": RGBColor(0x08, 0x91, 0xB2),  # secondary accent (cyan)
    "accent3": RGBColor(0x7C, 0x3A, 0xED),  # tertiary accent (purple)
    "hair":    RGBColor(0xE2, 0xE8, 0xF0),  # hairlines / grey series
    "panel":   RGBColor(0xF1, 0xF5, 0xF9),  # light panel fill
}
FONT = "Segoe UI"  # safe on Windows/Office; swap for the brand font

# 16:9 canvas + a shared grid
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)
CONTENT_W = SW - 2 * MARGIN


def _txt(slide, left, top, width, height, text, size, *, color, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return box


def _bullets(slide, left, top, width, height, items, size=24):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = PALETTE["ink"]
    return box


def _rect(slide, left, top, width, height, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# Layout 5 = "Title Only": a real title placeholder (so the deck is navigable and the
# validator sees titles) with full freedom below it. We restyle the title per slide but
# it still inherits from the layout — master-based, not orphan formatting.
TITLE_ONLY = 5


def _slide(prs, *, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PALETTE["ink"] if dark else PALETTE["bg"]
    return s


def _title(slide, text, *, left, top, width, height, size, color, bold=True, spacing=1.05):
    ph = slide.shapes.title
    ph.left, ph.top, ph.width, ph.height = int(left), int(top), int(width), int(height)
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = spacing
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = color
    return ph


def _accent_bar(slide, top=MARGIN, color=PALETTE["accent"]):
    _rect(slide, MARGIN, top, Inches(0.9), Inches(0.14), color)


def slide_title(prs):
    s = _slide(prs)
    _rect(s, 0, 0, Inches(0.28), SH, PALETTE["accent"])          # left spine
    _txt(s, MARGIN, Inches(1.7), CONTENT_W, Inches(0.5),
         "QUARTERLY REVIEW · Q1 2026", 18, color=PALETTE["accent"], bold=True)
    _title(s, "AI Platform Spend & Adoption",
           left=MARGIN, top=Inches(2.2), width=CONTENT_W, height=Inches(2.0),
           size=44, color=PALETTE["ink"], spacing=1.0)
    _txt(s, MARGIN, Inches(4.3), CONTENT_W, Inches(1.0),
         "Where spend concentrated, how it moved, and where to act.", 24, color=PALETTE["muted"])
    _txt(s, MARGIN, Inches(6.4), CONTENT_W, Inches(0.5),
         "Prepared by Platform Finance  ·  15 Jul 2026", 18, color=PALETTE["muted"])
    return s


def slide_section(prs):
    s = _slide(prs, dark=True)
    _txt(s, MARGIN, Inches(2.2), Inches(3.0), Inches(2.0),
         "01", 96, color=PALETTE["accent"], bold=True)
    _title(s, "Where the spend went",
           left=MARGIN, top=Inches(4.3), width=CONTENT_W, height=Inches(1.2),
           size=40, color=PALETTE["bg"])
    return s


def slide_content(prs):
    s = _slide(prs)
    _accent_bar(s)
    _title(s, "Spend rose every month and concentrated in the top three companies",
           left=MARGIN, top=Inches(1.15), width=CONTENT_W, height=Inches(1.3),
           size=32, color=PALETTE["ink"])
    _bullets(s, MARGIN, Inches(2.8), CONTENT_W, Inches(4.0), [
        "Q1 total reached $486k, up 12% on the prior quarter",
        "March was the largest month at $186k (+18%)",
        "Top 3 companies drove 68% of all spend",
        "214 distinct users were active across the platform",
    ], size=24)
    return s


def slide_two_col(prs):
    s = _slide(prs)
    _accent_bar(s)
    _title(s, "Northwind's growth is broad-based, not one heavy user",
           left=MARGIN, top=Inches(1.15), width=CONTENT_W, height=Inches(1.3),
           size=32, color=PALETTE["ink"])
    half = (CONTENT_W - Inches(0.6)) / 2
    _bullets(s, MARGIN, Inches(2.8), half, Inches(3.8), [
        "18 active users this quarter",
        "Spend rose steadily Jan → Mar",
        "Top user A. Chen at $8.4k in March",
        "No single user above 15% of the total",
    ], size=24)
    _rect(s, int(MARGIN + half + Inches(0.6)), int(Inches(2.8)), int(half), int(Inches(3.4)), PALETTE["panel"])
    _txt(s, MARGIN + half + Inches(0.6), Inches(4.2), half, Inches(0.8),
         "[ visual / chart ]", 20, color=PALETTE["muted"], align=PP_ALIGN.CENTER)
    return s


def slide_chart(prs):
    s = _slide(prs)
    _accent_bar(s)
    _title(s, "March drove the quarter's growth across all three companies",
           left=MARGIN, top=Inches(1.15), width=CONTENT_W, height=Inches(1.3),
           size=32, color=PALETTE["ink"])
    data = CategoryChartData()
    data.categories = ["Jan", "Feb", "Mar"]
    data.add_series("Northwind", (41, 47, 58))
    data.add_series("Acme", (33, 36, 40))
    data.add_series("Globex", (28, 30, 34))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            MARGIN, Inches(2.6), CONTENT_W, Inches(4.2), data)
    chart = gf.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    for series, col in zip(chart.plots[0].series,
                           (PALETTE["accent"], PALETTE["accent2"], PALETTE["accent3"])):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = col
    return s


def slide_closing(prs):
    s = _slide(prs, dark=True)
    _rect(s, 0, 0, Inches(0.28), SH, PALETTE["accent"])
    _txt(s, MARGIN, Inches(1.5), CONTENT_W, Inches(0.5),
         "RECOMMENDATION", 18, color=PALETTE["accent"], bold=True)
    _title(s, "Tie spend to outcomes and review the heaviest users monthly",
           left=MARGIN, top=Inches(2.1), width=CONTENT_W, height=Inches(1.8),
           size=34, color=PALETTE["bg"])
    _bullets_light(s, MARGIN, Inches(4.3), CONTENT_W, Inches(2.4), [
        "Monitor the top decile of users each month",
        "Attribute cost to shipped results, not raw activity",
        "Set guardrails on premium models and tool-heavy flows",
    ])
    return s


def _bullets_light(slide, left, top, width, height, items, size=24):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10); p.line_spacing = 1.15
        r = p.add_run(); r.text = "•  " + item
        r.font.size = Pt(size); r.font.name = FONT
        r.font.color.rgb = PALETTE["panel"]
    return box


def build(out: Path):
    prs = Presentation()
    prs.slide_width = int(SW)
    prs.slide_height = int(SH)
    slide_title(prs)
    slide_section(prs)
    slide_content(prs)
    slide_two_col(prs)
    slide_chart(prs)
    slide_closing(prs)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Wrote {out} — {len(prs.slides._sldIdLst)} slides, 16:9, brand-swappable via PALETTE/FONT.")


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    default = Path(__file__).resolve().parent.parent / "assets" / "starter-template.pptx"
    ap.add_argument("--out", default=str(default), help="output .pptx path")
    args = ap.parse_args()
    build(Path(args.out))


if __name__ == "__main__":
    main()
