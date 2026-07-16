"""Validate a .pptx deck before it ships.

Deterministic quality gate for a produced deck. Catches the failure modes that make
a deck look unfinished, prints a machine-readable JSON report, and exits non-zero on
ERRORS so the caller can fix and re-run (produce -> validate -> fix -> re-validate).

This gate reads the deck's *markup*. It is the structural half of a two-part check;
the visual half is render_pptx.py + a human/vision Read of each slide. Some defects
(text autofit-shrink or clipping, shape overlap, off-brand color, a chart that
renders wrong) are only visible in the render — this script cannot see them, and it
CANNOT resolve font sizes inherited from the master/layout (python-pptx leaves those
as None). That blind spot is exactly why the vision pass is mandatory, not optional.

Checks:
  ERROR   - Leftover placeholder text (lorem/ipsum, TBD, TODO, FIXME, XXX,
            PLACEHOLDER, and {{ }} / {% %} template tags) on any slide or in notes.
  WARNING - Empty slides (no title and no text/table content).
  WARNING - Slides with no title (harder to navigate; check intent).
  WARNING - Very text-dense slides (> ~60 words), a wall-of-text smell.
  WARNING - Shapes that overflow the slide bounds (clipped / bled off-canvas).
  WARNING - Explicitly-set fonts below the 18pt floor (inherited sizes not seen).
  WARNING - Raster images below ~96 DPI at their placed size (will look pixelated).

Usage:
    python scripts/validate_pptx.py path/to/deck.pptx
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PLACEHOLDER = re.compile(
    r"\bTBD\b|\bTODO\b|\bFIXME\b|\bXXX+\b|\bPLACEHOLDER\b|lorem\s*ipsum|{{.*?}}|{%.*?%}",
    re.IGNORECASE,
)
DENSE_WORDS = 60
MIN_FONT_PT = 18       # deck-anatomy: never below 18pt on any slide
MIN_IMAGE_DPI = 96     # below this a raster image looks pixelated even on screen
OVERFLOW_TOL_EMU = 91440  # ~0.1 inch slack before a shape counts as overflowing


def slide_texts(slide):
    """(title, [body strings], has_table, notes) for one slide."""
    title = slide.shapes.title.text.strip() if slide.shapes.title else ""
    body, has_table = [], False
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.has_table:
            has_table = True
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        body.append(cell.text.strip())
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        elif shape.has_text_frame and shape.text_frame.text.strip():
            body.append(shape.text_frame.text.strip())
    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    return title, body, has_table, notes


def check_geometry(slide, idx, sw, sh, warnings):
    """Warn on shapes whose box falls outside the slide (clipped / off-canvas).
    Top-level shapes only — a child that bleeds off-slide inside a group is not
    recursed into, and text that autofit-shrinks or overflows its (in-bounds) box
    is invisible here; both are the vision pass's job."""
    for shape in slide.shapes:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if None in (l, t, w, h):
            continue  # inherited/auto geometry — can't judge
        if (l < -OVERFLOW_TOL_EMU or t < -OVERFLOW_TOL_EMU
                or l + w > sw + OVERFLOW_TOL_EMU or t + h > sh + OVERFLOW_TOL_EMU):
            name = (shape.name or "shape").strip()
            warnings.append(f"slide {idx}: '{name}' overflows the slide bounds (clipped/off-canvas)")


def check_fonts(slide, idx, warnings):
    """Warn on EXPLICIT font sizes below the floor. Inherited (None) sizes are invisible
    to python-pptx — those are the vision pass's job, not this gate's."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                if sz is not None and sz.pt < MIN_FONT_PT:
                    snippet = (run.text or "").strip()[:30]
                    warnings.append(
                        f"slide {idx}: {sz.pt:g}pt text below {MIN_FONT_PT}pt floor — {snippet!r}"
                    )
                    return  # one report per slide is enough


def check_images(slide, idx, warnings):
    """Warn on raster images placed larger than their pixel data supports (low DPI)."""
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            img = shape.image
            px_w, px_h = img.size
            disp_w_in = shape.width / 914400.0
            disp_h_in = shape.height / 914400.0
        except Exception:
            continue
        if disp_w_in <= 0 or disp_h_in <= 0:
            continue
        dpi = min(px_w / disp_w_in, px_h / disp_h_in)
        if dpi < MIN_IMAGE_DPI:
            name = (shape.name or "image").strip()
            warnings.append(f"slide {idx}: '{name}' is ~{dpi:.0f} DPI at its size (<{MIN_IMAGE_DPI}) — will look pixelated")


def main():
    if len(sys.argv) != 2:
        sys.exit("usage: python scripts/validate_pptx.py <deck.pptx>")
    path = Path(sys.argv[1])
    if not path.exists():
        sys.exit(f"file not found: {path}")

    prs = Presentation(str(path))
    errors, warnings = [], []
    placeholders = []
    sw, sh = prs.slide_width, prs.slide_height

    for idx, slide in enumerate(prs.slides, start=1):
        title, body, has_table, notes = slide_texts(slide)
        for chunk in [title, *body, notes]:
            m = PLACEHOLDER.search(chunk or "")
            if m:
                placeholders.append(f"slide {idx}: {m.group(0)!r}")
        if not title and not body and not has_table:
            warnings.append(f"slide {idx}: empty (no title or content)")
        elif not title:
            warnings.append(f"slide {idx}: no title")
        words = sum(len(b.split()) for b in body)
        if words > DENSE_WORDS:
            warnings.append(f"slide {idx}: dense (~{words} words) — wall-of-text risk")
        check_geometry(slide, idx, sw, sh, warnings)
        check_fonts(slide, idx, warnings)
        check_images(slide, idx, warnings)

    if placeholders:
        errors.append(f"Leftover placeholder text: {placeholders}")

    report = {
        "file": str(path),
        "status": "OK" if not errors else "FAILED",
        "errors": errors,
        "warnings": warnings,
        "n_slides": len(prs.slides._sldIdLst),
        "note": "Structural gate only. Also run render_pptx.py and Read every slide — "
                "overflow, autofit shrink, overlap, and off-brand color are not visible here. "
                "If this deck was FILLED from a registry template, this is the WRONG gate: run "
                "building-document-templates/scripts/validate.py <deck> --client <c> --doc-type <t> "
                "(the placeholder/slide-count/residue checks live there).",
    }
    print(json.dumps(report, indent=2, ensure_ascii=False))
    sys.exit(0 if not errors else 1)


if __name__ == "__main__":
    main()
