"""Derive client-agnostic assets from a designer layout-gallery deck.

Input: a professionally designed PowerPoint "template gallery" (one showcase slide
per layout, content in empty placeholders, branding baked into masters/layouts).
Output, in two passes:

  1. DE-BRAND -> assets/layout-gallery.pptx
     Logo images are replaced with neutral CLIENT-LOGO placeholder chips (same
     geometry, same media parts — so a fill can swap the real client logo in one
     image-slot operation), legal/footer/address text is neutralized. The result
     carries no organization identity and is safe to commit; capable models copy
     it and prune slides for bespoke decks.

  2. CURATE -> registry/_builtin/<name>_visual/{template.pptx, manifest.json}
     Fill-ready templates per report type: a subset of gallery slides in a fixed
     narrative order, every empty text placeholder tagged `{{ field }}`, unused
     slides/layouts/masters pruned (file size), and a manifest authored in the
     same run — the standard scaffold -> fill -> validate -> render loop applies.

    python scripts/derive_gallery_templates.py --source "path/to/gallery.pptx"
    python scripts/derive_gallery_templates.py --source gallery.pptx --only exec_update_visual

Requires python-pptx + pillow.
"""
from __future__ import annotations

import argparse
import io
import json
import os
import re
import shutil
import sys
import zipfile
from datetime import date
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

SKILL_ROOT = Path(__file__).resolve().parent.parent
REPO_ROOT = SKILL_ROOT.parent
DEFAULT_REGISTRY = REPO_ROOT / "building-document-templates" / "registry"
GALLERY_OUT = SKILL_ROOT / "assets" / "layout-gallery.pptx"


def tag(name: str) -> str:
    return "{{ " + name + " }}"


# ─────────────────────────── Pass 1: de-brand ───────────────────────────────
# Media parts carrying organization identity -> replaced with neutral art of the
# SAME pixel size so layout geometry is untouched. The footer/cover logo parts
# become CLIENT-LOGO chips and are exposed as image slots in every curated
# manifest; tagline art becomes blank.
LOGO_PARTS = {
    "ppt/media/image1.png": "chip",     # footer logo (dark bg), used across masters
    "ppt/media/image3.png": "chip",     # cover lockup (on dark art)
    "ppt/media/image15.png": "chip",    # color logo on light bg
}
BLANK_PARTS = {
    "ppt/media/image4.png": "png",      # tagline strip
    "ppt/media/image16.jpeg": "jpeg",   # big tagline art (back cover)
}
# Text neutralization: any <a:t> run containing a key gets the replacement.
TEXT_SWAPS = [
    ("BUSINESS SCIENCE", "CONFIDENTIAL"),
    ("REVENUE SCIENCE", ""),
    ("Rivonia", "Company address  |  Contact  |  Website"),
    ("bscglobal", ""),
    ("+27", ""),                      # phone numbers
]


def _client_logo_chip(w: int, h: int) -> bytes:
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    pad = max(2, h // 12)
    d.rounded_rectangle([pad, pad, w - pad, h - pad], radius=max(4, h // 8),
                        outline=(150, 158, 170, 200), width=max(2, h // 40))
    text = "CLIENT LOGO"
    # default bitmap font ~6x11px; center it
    tw = len(text) * 6
    d.text(((w - tw) / 2, h / 2 - 6), text, fill=(150, 158, 170, 220))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _blank(kind: str, w: int, h: int) -> bytes:
    buf = io.BytesIO()
    if kind == "jpeg":
        # match the slide background so the blanked art disappears
        Image.new("RGB", (w, h), (245, 245, 245)).save(buf, format="JPEG", quality=80)
    else:
        Image.new("RGBA", (w, h), (0, 0, 0, 0)).save(buf, format="PNG")
    return buf.getvalue()


def _neutralize_xml(xml: str) -> str:
    def fix(m):
        text = m.group(1)
        for key, repl in TEXT_SWAPS:
            if key.lower() in text.lower():
                return f"<a:t>{repl}</a:t>"
        return m.group(0)
    return re.sub(r"<a:t>([^<]*)</a:t>", fix, xml)


# Collaboration/authorship parts carry names+emails and serve no template purpose.
DROP_PARTS = re.compile(r"ppt/(changesInfos/|authors\.xml|revisionInfo\.xml)")


def debrand(source: Path, out: Path) -> None:
    src = zipfile.ZipFile(source)
    dropped = {n for n in src.namelist() if DROP_PARTS.match(n)}
    out.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            if item.filename in dropped:
                continue
            data = src.read(item.filename)
            if item.filename in LOGO_PARTS:
                with Image.open(io.BytesIO(data)) as im:
                    data = _client_logo_chip(im.width, im.height)
            elif item.filename in BLANK_PARTS:
                with Image.open(io.BytesIO(data)) as im:
                    data = _blank(BLANK_PARTS[item.filename], im.width, im.height)
            elif (re.match(r"ppt/media/.*\.(png|jpe?g)$", item.filename, re.I)
                    and len(data) > 800_000):
                # decorative photography — slim it so the gallery and every curated
                # template stay committable (same part name/format, smaller pixels)
                with Image.open(io.BytesIO(data)) as im:
                    im = im.convert("RGBA") if item.filename.lower().endswith("png") else im.convert("RGB")
                    if im.width > 1920:
                        im = im.resize((1920, round(im.height * 1920 / im.width)), Image.LANCZOS)
                    buf = io.BytesIO()
                    if item.filename.lower().endswith("png"):
                        im.save(buf, format="PNG", optimize=True)
                    else:
                        im.save(buf, format="JPEG", quality=82, optimize=True)
                    if buf.tell() < len(data):
                        data = buf.getvalue()
            elif re.match(r"ppt/(slideMasters|slideLayouts|slides)/[^_].*\.xml$", item.filename):
                data = _neutralize_xml(data.decode("utf-8")).encode("utf-8")
            elif item.filename == "[Content_Types].xml":
                xml = data.decode("utf-8")
                xml = re.sub(r'<Override[^>]*PartName="/(?:ppt/(?:changesInfos/[^"]*|authors\.xml|revisionInfo\.xml))"[^>]*/>', "", xml)
                data = xml.encode("utf-8")
            elif item.filename.endswith(".rels"):
                xml = data.decode("utf-8")
                xml = re.sub(r'<Relationship[^>]*Target="[^"]*(?:changesInfos/[^"]*|authors\.xml|revisionInfo\.xml)"[^>]*/>', "", xml)
                data = xml.encode("utf-8")
            elif item.filename == "docProps/core.xml":
                xml = data.decode("utf-8")
                xml = re.sub(r"(<dc:creator>)[^<]*(</dc:creator>)", r"\1\2", xml)
                xml = re.sub(r"(<cp:lastModifiedBy>)[^<]*(</cp:lastModifiedBy>)", r"\1\2", xml)
                data = xml.encode("utf-8")
            dst.writestr(item, data)
    # prove nothing survived
    z = zipfile.ZipFile(out)
    leftovers = []
    for name in z.namelist():
        if name.endswith(".xml"):
            xml = z.read(name).decode("utf-8", errors="ignore")
            for pat in ("BUSINESS SCIENCE", "REVENUE SCIENCE", "Rivonia", "bscglobal"):
                if pat.lower() in xml.lower():
                    leftovers.append((name, pat))
    if leftovers:
        sys.exit(f"de-brand incomplete, identity text survived: {leftovers[:6]}")
    print(f"De-branded gallery -> {out} ({out.stat().st_size / 1e6:.1f} MB)")


# ─────────────────────────── Pass 2: curate ─────────────────────────────────

def _keep_and_order_slides(prs: Presentation, keep: list[int]) -> None:
    """Prune to `keep` (1-based source indexes) and reorder to match."""
    id_lst = prs.slides._sldIdLst
    entries = list(id_lst)
    by_source = {i: e for i, e in enumerate(entries, 1)}
    for i, entry in by_source.items():
        if i not in keep:
            id_lst.remove(entry)
            prs.part.drop_rel(entry.get(qn("r:id")))
    for entry in [by_source[i] for i in keep]:      # reorder to spec order
        id_lst.remove(entry)
        id_lst.append(entry)


def _prune_layouts_and_masters(prs: Presentation) -> None:
    """Drop layouts no kept slide uses, then masters with no layouts left —
    unreachable parts (incl. their media) are not serialized on save."""
    used = {s.slide_layout.part for s in prs.slides}
    for master in prs.slide_masters:
        id_lst = master.element.find(qn("p:sldLayoutIdLst"))
        if id_lst is None:
            continue
        for lid in list(id_lst):
            r_id = lid.get(qn("r:id"))
            if master.part.rels[r_id].target_part not in used:
                id_lst.remove(lid)
                master.part.drop_rel(r_id)
    m_lst = prs.element.find(qn("p:sldMasterIdLst"))
    for mid in list(m_lst):
        r_id = mid.get(qn("r:id"))
        mpart = prs.part.rels[r_id].target_part
        lay_lst = mpart._element.find(qn("p:sldLayoutIdLst"))
        if lay_lst is None or len(lay_lst) == 0:
            m_lst.remove(mid)
            prs.part.drop_rel(r_id)


def _max_chars(ph) -> int:
    """Rough capacity hint from the placeholder box (no autofit in the engine)."""
    try:
        w_in, h_in = Emu(ph.width).inches, Emu(ph.height).inches
    except TypeError:
        return 80
    lines = max(1, int(h_in / 0.32))
    return max(20, int(w_in * 11) * lines)


def F(name, example, guidance, *, type="text", required=True, media_part=""):
    return {"name": name, "type": type, "example": example, "guidance": guidance,
            "required": required, "media_part": media_part}


def tag_slide(slide, spec: dict) -> list[dict]:
    """Set every empty text placeholder to a `{{ tag }}` and return its fields.

    spec: {prefix, title (example), items_example, mode: items|columns,
           col_names: (left,right) when mode=columns}
    """
    fields: list[dict] = []
    prefix = spec["prefix"]
    title_ph, bodies = None, []
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        if t == 1 or ph.placeholder_format.idx == 0:            # TITLE
            title_ph = ph
        elif ph.has_text_frame:                                  # BODY/other text
            bodies.append(ph)
        # picture placeholders are left empty (render blank; user art optional)

    if title_ph is not None:
        name = f"{prefix}_title"
        title_ph.text_frame.text = tag(name)
        fields.append(F(name, spec["title"],
                        f"Headline for the {spec['label']} slide — state the takeaway."))

    if spec.get("mode") == "columns" and bodies:
        # a wide placeholder near the top is the slide's note/lead, not a column entry
        slide_w = 12192000  # 16:9 EMU
        wide = sorted((p for p in bodies if p.width and p.width > slide_w * 0.45),
                      key=lambda p: p.top)
        for j, ph in enumerate(wide, 1):
            name = f"{prefix}_note" if len(wide) == 1 else f"{prefix}_note{j}"
            ph.text_frame.text = tag(name)
            fields.append(F(name, "One muted line under the heading.",
                            "Short framing line for the slide.", required=False))
            bodies.remove(ph)
        mid = sum(p.left for p in bodies) / len(bodies)
        left = sorted((p for p in bodies if p.left <= mid), key=lambda p: (p.top, p.left))
        right = sorted((p for p in bodies if p.left > mid), key=lambda p: (p.top, p.left))
        for col_name, col in zip(spec["col_names"], (left, right)):
            for i, ph in enumerate(col, 1):
                name = f"{col_name}_{i}"
                ph.text_frame.text = tag(name)
                fields.append(F(name, f"{spec['items_example']} ({col_name} {i})",
                                f"{col_name.replace('_', ' ').title()} entry {i} "
                                f"(≈{_max_chars(ph)} chars fit).", required=False))
    else:
        bodies.sort(key=lambda p: (p.top, p.left))
        for i, ph in enumerate(bodies, 1):
            name = f"{prefix}_text{i}" if len(bodies) > 1 else f"{prefix}_body"
            ph.text_frame.text = tag(name)
            fields.append(F(name, spec["items_example"],
                            f"{spec['label']} slot {i} of {len(bodies)}, top-to-bottom "
                            f"left-to-right (≈{_max_chars(ph)} chars fit).",
                            required=False))
    return fields


# ─────────────────────────── Gallery index (the picker's map) ───────────────
# Human-reviewed category per gallery slide (vision pass over the full render).
# Ranges keep this maintainable; the --index command merges these with the
# objective per-slide facts (text/picture slot counts) into a browsable catalog.
CATEGORIES: list[tuple[range, str]] = [
    (range(1, 3), "cover — photographic"),
    (range(3, 4), "cover — dark, icon hexagons"),
    (range(4, 16), "agenda / numbered list (chevrons, bars, pills; 4–8 items)"),
    (range(16, 21), "numbered list with side panel / card grid"),
    (range(21, 26), "banner arrows / paired circles"),
    (range(26, 31), "plain title + content / card grids"),
    (range(31, 33), "two-panel content / icon feature panel"),
    (range(33, 38), "challenges — staircase bars or numbered list (3–6)"),
    (range(38, 41), "objectives — target diagrams with items"),
    (range(41, 48), "concept diagrams — gears, bulb, head/mind (3–5 parts)"),
    (range(48, 58), "content panels — numbered rows, drop shapes, wide panels"),
    (range(58, 64), "team / people — circle cards (3–5)"),
    (range(64, 72), "process steps — numbered panels, squares, loops (3–5)"),
    (range(72, 82), "lists with markers / layered bars / callout rows"),
    (range(82, 93), "journeys & roadmaps — roads, winding paths, rising-arrow milestones"),
    (range(93, 98), "timelines — horizontal markers, tag panels, dotted paths"),
    (range(98, 106), "sequence chains — links, circle chains, people arcs"),
    (range(106, 113), "quadrants / matrix — panels, petals, puzzles, pinwheel"),
    (range(113, 118), "layered arrows / mountain climb (A–E)"),
    (range(118, 126), "hierarchies — pyramids, funnels, rings (3–5 levels)"),
    (range(126, 132), "cycles & pros/cons — orbit, gear arc, PROS/CONS, T-chart"),
    (range(132, 137), "comparisons — branches, side-by-side, VS panels"),
    (range(137, 138), "back cover — contact/close"),
    (range(138, 148), "data-flavored — SA maps, KPI donuts/bars (fixed art: numbers must match!)"),
]


def category_for(n: int) -> str:
    for rng, label in CATEGORIES:
        if n in rng:
            return label
    return "diagram"


def build_index(gallery: Path, out_md: Path) -> None:
    """Emit the picker's catalog: per slide, what it is and what it can hold —
    so a model shortlists by reading, then renders ONLY the shortlist to confirm."""
    prs = Presentation(str(gallery))
    lines = [
        "# Layout gallery index",
        "",
        "Generated by `derive_gallery_templates.py --index` from `layout-gallery.pptx`.",
        "Pick candidate slides here, then RENDER YOUR SHORTLIST and look before using",
        "(`render_pptx.py`) — the category tells you what a slide is for; only vision",
        "tells you it fits. Slides marked *fixed art* encode quantities in the drawing;",
        "use them only when your numbers match the art.",
        "",
        "| # | Category | Text slots | Pic slots | Largest slot (chars ≈) |",
        "|---|----------|-----------:|----------:|-----------------------:|",
    ]
    for i, slide in enumerate(prs.slides, 1):
        texts, pics = [], 0
        for ph in slide.placeholders:
            if ph.placeholder_format.type == 18:
                pics += 1
            elif ph.has_text_frame:
                texts.append(_max_chars(ph))
        biggest = max(texts) if texts else 0
        lines.append(f"| {i} | {category_for(i)} | {len(texts)} | {pics} | {biggest} |")
    out_md.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"Gallery index -> {out_md} ({len(prs.slides._sldIdLst)} slides)")


# ─────────────────────────── Template specs ─────────────────────────────────
# Slide numbers are 1-based positions in the SOURCE gallery. Every spec ends on
# the neutral back-cover (137, no fields). Slides with quantity-encoding art
# (fixed donuts/bars) are deliberately excluded — their geometry would lie.

def S(n, prefix, label, title, items_example, mode="items", col_names=()):
    return {"n": n, "prefix": prefix, "label": label, "title": title,
            "items_example": items_example, "mode": mode, "col_names": col_names}


TEMPLATES: dict[str, dict] = {
    "exec_update_visual": {
        "description": "Executive/quarterly update on the designer layout system: cover, agenda, summary, challenges, roadmap, risks vs mitigations.",
        "slides": [
            S(2, "cover", "cover", "Q3 2026 Business Update", "Quarterly performance, drivers, and decisions requested."),
            S(8, "agenda", "agenda", "Agenda", "Performance against target"),
            S(50, "summary", "executive summary", "Executive summary", "The quarter closed 4% above plan on the strength of two enterprise renewals; margin improved 1.2 points."),
            S(36, "challenge", "key challenges", "What we are solving", "Cycle times vary 2.3x between shifts"),
            S(89, "milestone", "roadmap", "The road ahead", "Q4 hiring approved and posted"),
            S(130, "tradeoff", "pros and cons", "The trade-offs of the recommended plan", "Locks in the renewal upside", "columns", ("pro", "con")),
        ],
    },
    "project_kickoff_visual": {
        "description": "Project kickoff on the designer layout system: cover, agenda, objectives, team, journey, timeline, next steps.",
        "slides": [
            S(3, "cover", "cover", "Project Falcon — Kickoff", "Simulating the fleet transition across three sites."),
            S(7, "agenda", "agenda", "Agenda", "Project introduction and objectives"),
            S(39, "objective", "objectives", "Definition of victory", "A validated model of current operations"),
            S(58, "team", "team", "Meet the team", "Jane Mokoena — Engagement Lead"),
            S(84, "phase", "journey", "How we get there", "Phase 1: data audit"),
            S(93, "timeline", "timeline", "Timeline", "Sprint 1 — Weeks 1–2"),
            S(51, "next", "next steps", "Next steps", "Data extract of 12 months' records — by Friday"),
        ],
    },
    "proposal_visual": {
        "description": "Client proposal on the designer layout system: cover, context, challenges, approach, scope, team, investment.",
        "slides": [
            S(1, "cover", "cover", "Fleet Optimisation Study", "Proposal prepared for Acme Mining."),
            S(52, "context", "context", "Context and challenge", "Haulage costs rose 18% while utilisation fell — the operating model is the constraint."),
            S(37, "pain", "pain points", "What is holding performance back", "Cycle times vary 2.3x between shifts"),
            S(64, "approach", "approach", "Our approach", "Phase 1: baseline the operation"),
            S(14, "scope", "scope", "Scope and deliverables", "Validated simulation model"),
            S(59, "team", "team", "The team", "Jane Mokoena — Engagement Lead"),
            S(53, "investment", "investment", "Timeline and investment", "Eight weeks, three consultants, fixed fee."),
        ],
    },
    "report_out_visual": {
        "description": "Results report-out on the designer layout system: cover, summary, findings, before/after comparison, recommendations.",
        "slides": [
            S(2, "cover", "cover", "Haulage Study — Results", "What we found, what it means, what to do next."),
            S(50, "summary", "executive summary", "Executive summary", "The phased transition cuts cost per tonne 14% with payback inside 30 months."),
            S(15, "finding", "findings", "Key findings", "Cost per tonne falls 14% in the phased scenario"),
            S(135, "compare", "before/after comparison", "Today vs the proposed model", "Current state", "columns", ("before", "after")),
            S(51, "recommendation", "recommendations", "Recommendations", "Commit to the phased transition from Q2"),
        ],
    },
    "project_status_visual": {
        "description": "Project status/progress update on the designer layout system: cover, progress milestones, workstreams, risks, next steps.",
        "slides": [
            S(2, "cover", "cover", "Project Falcon — Status Week 6", "Progress, risks, and the decisions we need."),
            S(90, "progress", "progress milestones", "Where we are", "Data audit complete"),
            S(78, "workstream", "workstreams", "Workstream status", "Modelling — on track"),
            S(129, "tradeoff", "pros and cons", "Trade-offs in the current plan", "Sprint pace ahead of schedule", "columns", ("pro", "con")),
            S(51, "next", "next steps", "Next steps", "Scenario list sign-off — Thursday"),
        ],
    },
}
BACK_COVER = 137     # neutral closing slide, appended to every template (no fields)


def curate(gallery: Path, name: str, spec: dict, registry: Path, created: str) -> None:
    prs = Presentation(str(gallery))
    keep = [s["n"] for s in spec["slides"]] + [BACK_COVER]
    _keep_and_order_slides(prs, keep)
    _prune_layouts_and_masters(prs)

    fields: list[dict] = []
    for slide, s_spec in zip(prs.slides, spec["slides"]):
        fields += tag_slide(slide, s_spec)

    dest = registry / "_builtin" / name
    dest.mkdir(parents=True, exist_ok=True)
    tpl = dest / "template.pptx"
    prs.save(str(tpl))

    # logo image slots: expose whichever placeholder-chip parts this subset kept
    with zipfile.ZipFile(tpl) as z:
        kept_media = set(z.namelist())
    slot_specs = [("client_logo", "ppt/media/image1.png", "Footer/master logo"),
                  ("client_logo_cover", "ppt/media/image3.png", "Cover logo lockup"),
                  ("client_logo_light", "ppt/media/image15.png", "Logo on light backgrounds")]
    for fname, part, label in slot_specs:
        if part in kept_media:
            fields.append(F(fname, "", f"{label}: PNG path for the client's logo — swapped "
                            "everywhere this slot appears. Leave unset to keep the neutral "
                            "CLIENT LOGO chip.", type="image", required=False, media_part=part))

    manifest = {
        "template_id": f"_builtin/{name}",
        "client": "_builtin",
        "doc_type": name,
        "format": "pptx",
        "template_file": "template.pptx",
        "source_file": "curated from assets/layout-gallery.pptx by derive_gallery_templates.py",
        "version": "1.0.0",
        "owner": "",
        "created": created,
        "changelog": [f"{created}: curated from the de-branded layout gallery"],
        "description": spec["description"],
        "fields": fields,
        "row_groups": [],
        "source_terms": ["Acme Mining", "Jane Mokoena"],
    }
    (dest / "manifest.json").write_text(json.dumps(manifest, indent=2, ensure_ascii=False),
                                        encoding="utf-8")
    size = tpl.stat().st_size / 1e6
    print(f"Registered _builtin/{name}: {len(prs.slides._sldIdLst)} slides, "
          f"{len(fields)} fields, {size:.1f} MB -> {dest}")


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--source", help="The branded gallery .pptx (only needed to (re)build "
                                     "assets/layout-gallery.pptx; curation reuses the de-branded copy)")
    ap.add_argument("--only", choices=sorted(TEMPLATES), help="Curate a single template")
    ap.add_argument("--index", action="store_true",
                    help="(Re)generate assets/layout-gallery-index.md and exit")
    ap.add_argument("--skip-debrand", action="store_true",
                    help="Reuse the existing assets/layout-gallery.pptx")
    ap.add_argument("--registry", default=os.environ.get("TEMPLATE_REGISTRY", str(DEFAULT_REGISTRY)))
    ap.add_argument("--created", default=date.today().isoformat())
    args = ap.parse_args()

    if args.index:
        if not GALLERY_OUT.exists():
            sys.exit(f"{GALLERY_OUT} missing — run once with --source")
        build_index(GALLERY_OUT, GALLERY_OUT.with_name("layout-gallery-index.md"))
        return

    if not args.skip_debrand:
        if not args.source:
            sys.exit("--source <branded gallery.pptx> is required unless --skip-debrand")
        debrand(Path(args.source), GALLERY_OUT)
    elif not GALLERY_OUT.exists():
        sys.exit(f"{GALLERY_OUT} missing — run once with --source")

    registry = Path(args.registry)
    for name in ([args.only] if args.only else sorted(TEMPLATES)):
        curate(GALLERY_OUT, name, TEMPLATES[name], registry, args.created)
    print("\nDiscover: python ../building-document-templates/scripts/registry.py list")


if __name__ == "__main__":
    main()
