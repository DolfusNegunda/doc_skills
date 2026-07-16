"""Generate the built-in PPTX template library — professional decks a small model FILLS.

Each template is a complete, branded, 16:9 deck whose variable text is already
`{{ tagged }}`, written straight into the document-template registry together with a
hand-authored manifest (fields with examples + guidance). Template and manifest come
from the same code, so they can never drift. The fill path is the standard engine:

    python ../building-document-templates/scripts/registry.py show --client _builtin --doc-type exec_update
    python ../building-document-templates/scripts/registry.py scaffold --builtin exec_update --out content.json
    python ../building-document-templates/scripts/fill.py --client _builtin --doc-type exec_update \
        --data content.json --out out.pptx
    python ../building-document-templates/scripts/validate.py out.pptx --template <registry>/_builtin/exec_update/template.pptx

Branding is data, not code: every color, font, logo, and footer string comes from a
brand pack (see ../../brands/README.md). Re-run with a different --brand to re-skin
the whole library.

    python scripts/build_template_library.py                    # all templates, default brand
    python scripts/build_template_library.py --brand path/to/client-pack   # re-skin for a client
    python scripts/build_template_library.py --only exec_update # one template

Templates: exec_update (QBR/quarterly review), project_kickoff, proposal, report_out.
Requires python-pptx + pillow.
"""
from __future__ import annotations

import argparse
import copy
import io
import json
import os
import sys
import zipfile
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SKILL_ROOT = Path(__file__).resolve().parent.parent
REPO_ROOT = SKILL_ROOT.parent
BRANDS = REPO_ROOT / "brands"
DEFAULT_REGISTRY = REPO_ROOT / "building-document-templates" / "registry"

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)
CONTENT_W = SW - 2 * MARGIN


def tag(name: str) -> str:
    return "{{ " + name + " }}"


# ---------------- Brand pack ----------------

def _deep_merge(base: dict, over: dict) -> dict:
    out = dict(base)
    for k, v in over.items():
        out[k] = _deep_merge(base[k], v) if isinstance(v, dict) and isinstance(base.get(k), dict) else v
    return out


def load_brand(name_or_path: str | None) -> dict:
    default = json.loads((BRANDS / "default" / "brand.json").read_text(encoding="utf-8"))
    default["_dir"] = BRANDS / "default"
    if not name_or_path or name_or_path == "default":
        return default
    p = Path(name_or_path)
    if p.suffix == ".json" and p.exists():
        pack_file, pack_dir = p, p.parent
    elif p.is_dir() and (p / "brand.json").exists():
        pack_file, pack_dir = p / "brand.json", p
    elif (BRANDS / name_or_path / "brand.json").exists():
        pack_dir = BRANDS / name_or_path
        pack_file = pack_dir / "brand.json"
    else:
        known = sorted(d.name for d in BRANDS.iterdir() if (d / "brand.json").exists())
        sys.exit(f"Unknown brand '{name_or_path}'. Available: {', '.join(known)}")
    pack = _deep_merge(default, json.loads(pack_file.read_text(encoding="utf-8")))
    pack["_dir"] = pack_dir
    return pack


def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


class Style:
    """Brand pack resolved into the handful of drawing tokens the slides use."""

    def __init__(self, brand: dict):
        c = brand["colors"]
        self.bg = rgb(c["bg"])
        self.ink = rgb(c["ink"])
        self.muted = rgb(c["muted"])
        self.accent = rgb(c["primary"])
        self.accent2 = rgb(c["accent"])
        self.dark = rgb(c["dark"])
        self.panel = rgb(c["panel"])
        self.hair = rgb(c["hairline"])
        self.light = rgb("#F5F7FA")
        self.font = brand["fonts"]["body"]
        self.font_h = brand["fonts"]["heading"]
        year = date.today().year
        f = brand.get("footer", {})
        cop = (f.get("copyright") or "").replace("{year}", str(year)) \
                                        .replace("{company}", brand.get("display_name", ""))
        conf = f.get("confidentiality", "")
        self.footer = "   ·   ".join(x for x in (cop, conf) if x)
        logo_rel = brand.get("logo")
        self.logo_bytes = (brand["_dir"] / logo_rel).read_bytes() if logo_rel else None


# ---------------- Drawing helpers ----------------

def _txt(slide, st, left, top, width, height, text, size, *, color=None, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=None, spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font or st.font
    r.font.color.rgb = color if color is not None else st.ink
    return box


def _bullet_para(p, char="•"):
    """Give a paragraph a real PowerPoint bullet (survives fill.py's list expansion,
    which deep-copies the paragraph XML per item). Without a hanging indent the
    glyph renders flush against the text ("•Like this"), so set marL/indent too."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", "228600")      # 0.25" left margin...
    pPr.set("indent", "-228600")   # ...hanging, so the bullet sits in the gutter
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial", "pitchFamily": "34", "charset": "0"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(buFont)
    pPr.append(buChar)


def _list_box(slide, st, left, top, width, height, list_tag, size=20, color=None):
    """A textbox whose single paragraph is exactly `{{ tag }}` — fill.py clones the
    (bulleted) paragraph once per list item, giving real repeating bullets."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    p.space_after = Pt(10)
    _bullet_para(p)
    r = p.add_run()
    r.text = list_tag
    r.font.size = Pt(size)
    r.font.name = st.font
    r.font.color.rgb = color if color is not None else st.ink
    return box


def _rect(slide, left, top, width, height, color, *, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, int(left), int(top), int(width), int(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


TITLE_ONLY = 5  # layout with a real title placeholder; everything else is free


def _slide(prs, st, *, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = st.dark if dark else st.bg
    return s


def _title(slide, st, text, *, left=MARGIN, top=Inches(1.15), width=CONTENT_W,
           height=Inches(1.3), size=32, color=None, bold=True, spacing=1.05):
    ph = slide.shapes.title
    ph.left, ph.top, ph.width, ph.height = int(left), int(top), int(width), int(height)
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = spacing
    p.text = text
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = st.font_h
    r.font.color.rgb = color if color is not None else st.ink
    return ph


def _accent_bar(slide, st, top=MARGIN):
    _rect(slide, MARGIN, top, Inches(0.9), Inches(0.14), st.accent)


def _footer(slide, st, *, dark=False):
    if st.footer:
        _txt(slide, st, MARGIN, SH - Inches(0.42), CONTENT_W, Inches(0.3), st.footer, 8.5,
             color=st.light if dark else st.muted)


def _logo(slide, st, *, top=Inches(0.55), left=None, height=Inches(0.6)):
    """Logo on a white chip so it reads on light AND dark slides."""
    if not st.logo_bytes:
        return
    from PIL import Image
    img = Image.open(io.BytesIO(st.logo_bytes))
    ratio = img.width / img.height
    w = int(height * ratio)
    x = int(left) if left is not None else int(SW - MARGIN - w)
    pad = Inches(0.09)
    _rect(slide, x - pad, top - pad, w + 2 * pad, height + 2 * pad, rgb("#FFFFFF"), rounded=True)
    slide.shapes.add_picture(io.BytesIO(st.logo_bytes), x, int(top), height=int(height))


def _eyebrow(slide, st, text, top=Inches(1.7), *, color=None):
    _txt(slide, st, MARGIN, top, CONTENT_W, Inches(0.45), text, 15,
         color=color if color is not None else st.accent, bold=True)


# ---------------- Slide constructors ----------------

def title_slide(prs, st, eyebrow, title, subtitle, footnote):
    s = _slide(prs, st)
    _rect(s, 0, 0, Inches(0.28), SH, st.accent)
    _logo(s, st)
    _eyebrow(s, st, eyebrow)
    _title(s, st, title, top=Inches(2.2), height=Inches(1.9), size=42, spacing=1.0)
    _txt(s, st, MARGIN, Inches(4.25), CONTENT_W, Inches(1.0), subtitle, 22, color=st.muted)
    _txt(s, st, MARGIN, Inches(6.35), CONTENT_W, Inches(0.5), footnote, 15, color=st.muted)
    _footer(s, st)
    return s


def divider_slide(prs, st, number, heading, note):
    s = _slide(prs, st, dark=True)
    _txt(s, st, MARGIN, Inches(2.0), Inches(3.0), Inches(1.8), number, 88,
         color=st.accent2, bold=True)
    _title(s, st, heading, top=Inches(4.1), height=Inches(1.2), size=38, color=st.bg)
    _txt(s, st, MARGIN, Inches(5.4), CONTENT_W, Inches(0.9), note, 18, color=st.light)
    _footer(s, st, dark=True)
    return s


def bullets_slide(prs, st, heading, list_tag, *, intro=None, size=20):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, heading)
    top = Inches(2.6)
    if intro:
        _txt(s, st, MARGIN, top, CONTENT_W, Inches(0.9), intro, 17, color=st.muted, spacing=1.25)
        top = Inches(3.5)
    _list_box(s, st, MARGIN, top, CONTENT_W, SH - top - Inches(0.7), list_tag, size=size)
    _footer(s, st)
    return s


def text_slide(prs, st, heading, body_tag):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, heading)
    _txt(s, st, MARGIN, Inches(2.6), CONTENT_W, Inches(4.0), body_tag, 19,
         color=st.ink, spacing=1.35)
    _footer(s, st)
    return s


def kpi_slide(prs, st, heading, kpis):
    """kpis: list of (label_tag, value_tag, delta_tag)."""
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, heading)
    n = len(kpis)
    gap = Inches(0.3)
    w = (CONTENT_W - gap * (n - 1)) / n
    top, h = Inches(2.9), Inches(2.6)
    for i, (label, value, delta) in enumerate(kpis):
        x = MARGIN + i * (w + gap)
        _rect(s, x, top, w, h, st.panel, rounded=True)
        pad = Inches(0.25)
        _txt(s, st, x + pad, top + Inches(0.3), w - 2 * pad, Inches(0.4), label, 12,
             color=st.muted, bold=True)
        _txt(s, st, x + pad, top + Inches(0.85), w - 2 * pad, Inches(0.9), value, 30,
             color=st.ink, bold=True)
        _txt(s, st, x + pad, top + Inches(1.9), w - 2 * pad, Inches(0.45), delta, 13,
             color=st.accent)
    _footer(s, st)
    return s


def two_list_slide(prs, st, heading, left_label, left_tag, right_label, right_tag):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, heading)
    half = (CONTENT_W - Inches(0.6)) / 2
    top = Inches(2.7)
    for label, ltag, x in ((left_label, left_tag, MARGIN),
                           (right_label, right_tag, MARGIN + half + Inches(0.6))):
        _txt(s, st, x, top, half, Inches(0.4), label, 13, color=st.accent, bold=True)
        _list_box(s, st, x, top + Inches(0.55), half, Inches(3.6), ltag, size=17)
    _footer(s, st)
    return s


def visual_slide(prs, st, heading, caption_tag, placeholder_png):
    """Image slot: the picture is swapped by fill.py through the manifest's
    image-type field (media_part). Geometry is preserved on swap."""
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, heading)
    img_top, img_h = Inches(2.55), Inches(3.7)
    s.shapes.add_picture(io.BytesIO(placeholder_png), int(MARGIN), int(img_top),
                         width=int(CONTENT_W), height=int(img_h))
    _txt(s, st, MARGIN, img_top + img_h + Inches(0.15), CONTENT_W, Inches(0.5),
         caption_tag, 14, color=st.muted)
    _footer(s, st)
    return s


def closing_slide(prs, st, eyebrow, statement_tag, list_tag):
    s = _slide(prs, st, dark=True)
    _rect(s, 0, 0, Inches(0.28), SH, st.accent)
    _logo(s, st, top=Inches(0.55))
    _eyebrow(s, st, eyebrow, top=Inches(1.4), color=st.accent2)
    _title(s, st, statement_tag, top=Inches(1.95), height=Inches(1.7), size=32, color=st.bg)
    _list_box(s, st, MARGIN, Inches(3.9), CONTENT_W, Inches(2.6), list_tag, size=19, color=st.light)
    _footer(s, st, dark=True)
    return s


# ---------------- Placeholder visual (image slot content) ----------------

def placeholder_chart_png(st) -> bytes:
    """A branded 'replace me' visual for image slots — obvious in any vision pass."""
    from PIL import Image, ImageDraw
    W, H = 1600, 480
    img = Image.new("RGB", (W, H), tuple(st.panel))
    d = ImageDraw.Draw(img)
    heights = [0.45, 0.62, 0.5, 0.74, 0.66, 0.88, 0.8]
    n = len(heights)
    bw = W // (n * 2)
    base = H - 70
    for i, hgt in enumerate(heights):
        x0 = int(W * 0.06) + i * 2 * bw
        col = tuple(st.accent) if i % 2 == 0 else tuple(st.accent2)
        d.rectangle([x0, base - int(hgt * (H - 150)), x0 + bw, base], fill=col)
    d.rectangle([0, 0, W - 1, H - 1], outline=tuple(st.muted), width=2)
    d.text((int(W * 0.06), 18), "SAMPLE VISUAL - swap via the manifest's image field "
           "(fill.py) or replace before shipping", fill=tuple(st.ink))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ---------------- Field helper ----------------

def F(name, example, guidance, *, type="text", required=True, media_part=""):
    return {"name": name, "type": type, "example": example, "guidance": guidance,
            "required": required, "media_part": media_part}


# ---------------- Template definitions ----------------
# Examples deliberately use the fictional client "Acme Mining" and fictional people —
# they are also the manifest's source_terms, so a fill that ships the examples
# verbatim FAILS validate.py's source-residue check.

SOURCE_TERMS = ["Acme Mining", "Jane Mokoena", "Sipho Dlamini"]


# ---------------- Deck presets (all built on the composable flex body) ----------------
# One visual system for every deck builtin: cover + 12 designed body types + closing
# (see _flex_spec below). Each preset differs only in its cover/closing EXAMPLES and
# its recommended DEFAULT body sequence — the scaffold emits that sequence, and the
# filler may add/remove/reorder/swap types freely. No fixed decorative imagery:
# the image type is a required-swap slot, charts/tables are built from data.

def _flex_preset(prs, st, png, *, default, eyebrow, title, subtitle, closing):
    spec = _flex_spec(prs, st, png)
    spec["body"]["default"] = default
    ex = {"deck_eyebrow": eyebrow, "deck_title": title, "deck_subtitle": subtitle,
          "closing_statement": closing}
    for f in spec["fields"]:
        if f["name"] in ex:
            f["example"] = ex[f["name"]]
    return spec


def build_exec_update(prs, st, png):
    return _flex_preset(
        prs, st, png,
        default=["stats", "section", "chart", "bullets", "bullets", "two_col"],
        eyebrow="EXECUTIVE REVIEW · Q3 2026", title="Q3 2026 Business Update",
        subtitle="Performance against target, the drivers behind it, and the decisions we need.",
        closing="Approve the Q4 plan")


def build_project_kickoff(prs, st, png):
    return _flex_preset(
        prs, st, png,
        default=["agenda", "bullets", "numbered", "team", "timeline", "two_col"],
        eyebrow="PROJECT KICKOFF · 3 NOVEMBER 2026", title="Haulage Transition Simulation Study",
        subtitle="Simulating the fleet transition to trolley-assist across three sites.",
        closing="Data access and site contacts unlock sprint 1")


def build_proposal(prs, st, png):
    return _flex_preset(
        prs, st, png,
        default=["bullets", "numbered", "bullets", "team", "table", "quote"],
        eyebrow="PROPOSAL · 20 NOVEMBER 2026", title="Fleet Optimisation Study",
        subtitle="Prepared for the operations leadership team.",
        closing="Approve the engagement to start in January")


def build_report_out(prs, st, png):
    return _flex_preset(
        prs, st, png,
        default=["bullets", "section", "bullets", "bullets", "numbered", "quote"],
        eyebrow="PROJECT REPORT · 12 DECEMBER 2026", title="Haulage Transition Study — Results",
        subtitle="What we found, what it means, and what to do next.",
        closing="Board decision on phase 1 within 90 days")


# ---------------- flex_deck: composable body slides (Rev Sci-style motifs) ----------------
# The whole middle of the deck is a typed, ordered "body" list: the fill engine
# clones one designed source slide per entry (any mix, any order) and deletes the
# unused sources. Item rows inside a slide are cloned shape GROUPS (circle rows,
# chevrons, stat cards, timeline nodes) — entry count = row count, colors walk a
# brand ramp. No fixed decorative image anywhere; logo/footer come from the brand pack.

def _mix(c1, c2, t):
    return RGBColor(*(round(a + (b - a) * t) for a, b in zip(c1, c2)))


def _ramp(st, n=6):
    """Brand color ramp for numbered/row motifs (dark -> primary -> accent -> tint)."""
    stops = [st.dark, st.accent, st.accent2, _mix(st.accent2, RGBColor(0xFF, 0xFF, 0xFF), 0.45)]
    out = []
    for i in range(n):
        pos = i * (len(stops) - 1) / (n - 1)
        lo = min(int(pos), len(stops) - 2)
        out.append(_mix(stops[lo], stops[lo + 1], pos - lo))
    return out


def _grp(slide):
    g = slide.shapes.add_group_shape()
    return g


def _item_spec(field, shape, *, dx=0, dy=0, mx=5, ramp=None, subfields=(),
               center=False):
    spec = {"field": field, "shape": shape, "dx": int(dx), "dy": int(dy), "max": mx,
            "subfields": list(subfields)}
    if ramp:
        spec["ramp"] = [str(c) for c in ramp]
    if center:
        spec.update(center=True, span=int(CONTENT_W), span_left=int(MARGIN))
    return spec


def _sub(name, example, guidance, *, required=True):
    return {"name": name, "example": example, "guidance": guidance, "required": required}


def flex_agenda_slide(prs, st, ramp):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, tag("heading"))
    g = _grp(s)
    g.name = "ITEM"
    y = Inches(2.45)
    d = Inches(0.62)
    circle = g.shapes.add_shape(MSO_SHAPE.OVAL, int(MARGIN), int(y), int(d), int(d))
    circle.fill.solid(); circle.fill.fore_color.rgb = ramp[0]
    circle.line.fill.background(); circle.shadow.inherit = False
    tf = circle.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag("item._n")
    r.font.size, r.font.bold, r.font.name = Pt(18), True, st.font
    r.font.color.rgb = rgb("#FFFFFF")
    tx = MARGIN + d + Inches(0.35)
    _txt(g, st, tx, y - Inches(0.02), CONTENT_W - d - Inches(0.35), Inches(0.4),
         tag("item.title"), 17, bold=True)
    _txt(g, st, tx, y + Inches(0.32), CONTENT_W - d - Inches(0.35), Inches(0.35),
         tag("item.text"), 12.5, color=st.muted)
    _footer(s, st)
    return s


def flex_numbered_slide(prs, st, ramp):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, tag("heading"))
    g = _grp(s)
    g.name = "ITEM"
    y = Inches(2.45)
    ch_w, ch_h = Inches(0.85), Inches(0.7)
    chev = g.shapes.add_shape(MSO_SHAPE.CHEVRON, int(MARGIN), int(y), int(ch_w), int(ch_h))
    chev.fill.solid(); chev.fill.fore_color.rgb = ramp[0]
    chev.line.fill.background(); chev.shadow.inherit = False
    tf = chev.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag("item._n")
    r.font.size, r.font.bold, r.font.name = Pt(16), True, st.font
    r.font.color.rgb = rgb("#FFFFFF")
    tx = MARGIN + ch_w + Inches(0.35)
    _txt(g, st, tx, y - Inches(0.02), CONTENT_W - ch_w - Inches(0.35), Inches(0.4),
         tag("item.title"), 16.5, bold=True)
    _txt(g, st, tx, y + Inches(0.34), CONTENT_W - ch_w - Inches(0.35), Inches(0.35),
         tag("item.text"), 12.5, color=st.muted)
    _footer(s, st)
    return s


def flex_stats_slide(prs, st, ramp):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, tag("heading"))
    g = _grp(s)
    g.name = "ITEM"
    w, h = Inches(2.75), Inches(2.5)
    x, y = MARGIN, Inches(2.9)
    card = g.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
    card.fill.solid(); card.fill.fore_color.rgb = st.panel
    card.line.fill.background(); card.shadow.inherit = False
    strip = g.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x + Inches(0.25)),
                               int(y + Inches(0.25)), int(Inches(0.6)), int(Inches(0.12)))
    strip.fill.solid(); strip.fill.fore_color.rgb = ramp[0]
    strip.line.fill.background(); strip.shadow.inherit = False
    pad = Inches(0.25)
    _txt(g, st, x + pad, y + Inches(0.55), w - 2 * pad, Inches(0.4), tag("item.label"),
         12, color=st.muted, bold=True)
    _txt(g, st, x + pad, y + Inches(1.0), w - 2 * pad, Inches(0.8), tag("item.value"),
         28, bold=True)
    _txt(g, st, x + pad, y + Inches(1.9), w - 2 * pad, Inches(0.45), tag("item.note"),
         12.5, color=st.accent)
    _footer(s, st)
    return s


def flex_team_slide(prs, st, ramp):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, tag("heading"))
    g = _grp(s)
    g.name = "ITEM"
    y = Inches(2.4)
    d = Inches(0.72)
    circle = g.shapes.add_shape(MSO_SHAPE.OVAL, int(MARGIN), int(y), int(d), int(d))
    circle.fill.solid(); circle.fill.fore_color.rgb = ramp[0]
    circle.line.fill.background(); circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag("item.initials")
    r.font.size, r.font.bold, r.font.name = Pt(15), True, st.font
    r.font.color.rgb = rgb("#FFFFFF")
    tx = MARGIN + d + Inches(0.35)
    _txt(g, st, tx, y, Inches(4.6), Inches(0.4), tag("item.name"), 16, bold=True)
    _txt(g, st, tx, y + Inches(0.38), Inches(4.6), Inches(0.35), tag("item.role"),
         12.5, color=st.accent, bold=True)
    _txt(g, st, MARGIN + Inches(6.2), y + Inches(0.1), CONTENT_W - Inches(6.2),
         Inches(0.6), tag("item.note"), 12.5, color=st.muted)
    _footer(s, st)
    return s


def flex_timeline_slide(prs, st, ramp):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, tag("heading"))
    _rect(s, MARGIN, Inches(4.04), CONTENT_W, Inches(0.03), st.hair)
    g = _grp(s)
    g.name = "ITEM"
    x = MARGIN
    d = Inches(0.34)
    dot = g.shapes.add_shape(MSO_SHAPE.OVAL, int(x + Inches(0.9)), int(Inches(3.9)),
                             int(d), int(d))
    dot.fill.solid(); dot.fill.fore_color.rgb = ramp[0]
    dot.line.color.rgb = rgb("#FFFFFF"); dot.line.width = Pt(2)
    dot.shadow.inherit = False
    _txt(g, st, x, Inches(2.95), Inches(2.15), Inches(0.4), tag("item.label"), 14,
         bold=True, align=PP_ALIGN.CENTER)
    _txt(g, st, x, Inches(3.42), Inches(2.15), Inches(0.35), tag("item.date"), 12,
         color=st.accent, bold=True, align=PP_ALIGN.CENTER)
    _txt(g, st, x, Inches(4.5), Inches(2.15), Inches(1.5), tag("item.text"), 11.5,
         color=st.muted, align=PP_ALIGN.CENTER)
    _footer(s, st)
    return s


def flex_chart_slide(prs, st):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, tag("heading"))
    slot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(MARGIN), int(Inches(2.5)),
                              int(CONTENT_W), int(Inches(3.9)))
    slot.name = "CHART_SLOT"
    slot.fill.solid(); slot.fill.fore_color.rgb = st.panel
    slot.line.color.rgb = st.hair; slot.shadow.inherit = False
    tfp = slot.text_frame.paragraphs[0]
    r = tfp.add_run(); r.text = "Native chart renders here (fill.py builds it from the entry's data)"
    r.font.size, r.font.name, r.font.color.rgb = Pt(12), st.font, st.muted
    _txt(s, st, MARGIN, Inches(6.55), CONTENT_W, Inches(0.45), tag("takeaway"), 13,
         color=st.accent, bold=True)
    _footer(s, st)
    return s


def flex_table_slide(prs, st):
    s = _slide(prs, st)
    _accent_bar(s, st)
    _title(s, st, tag("heading"))
    slot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(MARGIN), int(Inches(2.5)),
                              int(CONTENT_W), int(Inches(3.9)))
    slot.name = "TABLE_SLOT"
    slot.fill.solid(); slot.fill.fore_color.rgb = st.panel
    slot.line.color.rgb = st.hair; slot.shadow.inherit = False
    tfp = slot.text_frame.paragraphs[0]
    r = tfp.add_run(); r.text = "Native table renders here (fill.py builds it from the entry's columns/rows)"
    r.font.size, r.font.name, r.font.color.rgb = Pt(12), st.font, st.muted
    _txt(s, st, MARGIN, Inches(6.55), CONTENT_W, Inches(0.45), tag("note"), 12,
         color=st.muted)
    _footer(s, st)
    return s


def flex_quote_slide(prs, st):
    s = _slide(prs, st, dark=True)
    _rect(s, 0, 0, Inches(0.28), SH, st.accent)
    _txt(s, st, MARGIN, Inches(1.6), Inches(1.2), Inches(1.4), "“", 110,
         color=st.accent2, bold=True)
    _txt(s, st, MARGIN + Inches(0.1), Inches(2.7), CONTENT_W - Inches(0.4), Inches(2.4),
         tag("quote"), 28, color=st.bg, spacing=1.25)
    _txt(s, st, MARGIN + Inches(0.1), Inches(5.6), CONTENT_W, Inches(0.5),
         tag("attribution"), 15, color=st.light)
    _footer(s, st, dark=True)
    return s


def _flex_spec(prs, st, png):
    ramp = _ramp(st)
    ramp_hex = ["#" + str(c) for c in ramp]

    title_slide(prs, st, tag("deck_eyebrow"), tag("deck_title"),
                tag("deck_subtitle"), tag("author_line"))                       # 0
    flex_agenda_slide(prs, st, ramp)                                            # 1  agenda
    bullets_slide(prs, st, tag("heading"), tag("points"), intro=tag("intro"))   # 2  bullets
    flex_numbered_slide(prs, st, ramp)                                          # 3  numbered
    flex_stats_slide(prs, st, ramp)                                             # 4  stats
    two_list_slide(prs, st, tag("heading"), tag("left_label"), tag("left_items"),
                   tag("right_label"), tag("right_items"))                      # 5  two_col
    flex_team_slide(prs, st, ramp)                                              # 6  team
    flex_timeline_slide(prs, st, ramp)                                          # 7  timeline
    flex_chart_slide(prs, st)                                                   # 8  chart
    flex_table_slide(prs, st)                                                   # 9  table
    visual_slide(prs, st, tag("heading"), tag("caption"), png)                  # 10 image
    flex_quote_slide(prs, st)                                                   # 11 quote
    divider_slide(prs, st, tag("marker"), tag("heading"), tag("note"))          # 12 section
    closing_slide(prs, st, "CLOSING", tag("closing_statement"), tag("next_steps"))  # 13

    fields = [
        F("deck_eyebrow", "OPERATIONS REVIEW · Q3 2026", "Small uppercase kicker on the cover: deck kind + period."),
        F("deck_title", "Q3 2026 Operations Review", "The deck title. Keep under ~8 words."),
        F("deck_subtitle", "What moved, what's at risk, and the decisions we need.", "One-sentence framing under the title. Omit if the user gave nothing suitable.", required=False),
        F("author_line", "Prepared by Operations · 15 October 2026", "Author/team and date. Omit if unknown — never invent one.", required=False),
        F("closing_statement", "Approve the Q4 operations plan", "The single closing ask, one sentence."),
        F("next_steps", "Confirm peak-season staffing plan — owner TBC — by 31 October", "Closing bullets ONLY from user-supplied actions/owners/dates. Omit entirely if none were given.", type="list", required=False),
    ]

    H = lambda ex: F("heading", ex, "This slide's headline — state the takeaway, not the topic.")
    body_types = {
        "agenda": {"slide_index": 1, "purpose": "Numbered agenda/overview rows (2–5).",
                   "fields": [H("What we'll cover")],
                   "items": _item_spec("items", "ITEM", dy=Inches(0.98), mx=5, ramp=ramp_hex,
                                       subfields=[_sub("title", "Financial performance", "Row title (3–6 words)."),
                                                  _sub("text", "Revenue, margin, and the fuel-cost squeeze", "One-line description.", required=False)])},
        "bullets": {"slide_index": 2, "purpose": "Classic message slide: headline + 3–6 bullets (+ optional intro line).",
                    "fields": [H("On-time delivery is approaching target"),
                               F("intro", "", "Optional one-line setup above the bullets.", required=False),
                               F("points", "94.1% on-time, 0.9 points below the 95% target", "3–6 bullets, one message each.", type="list")]},
        "numbered": {"slide_index": 3, "purpose": "Sequential steps/priorities as chevron rows (2–5).",
                     "fields": [H("How we get to launch")],
                     "items": _item_spec("items", "ITEM", dy=Inches(0.95), mx=5, ramp=ramp_hex,
                                         subfields=[_sub("title", "Complete discovery", "Step title."),
                                                    _sub("text", "Data audit and stakeholder interviews, weeks 1–2", "One-line detail.", required=False)])},
        "stats": {"slide_index": 4, "purpose": "2–4 KPI stat cards.",
                  "fields": [H("The quarter at a glance")],
                  "items": _item_spec("items", "ITEM", dx=Inches(3.05), mx=4, ramp=ramp_hex, center=True,
                                      subfields=[_sub("label", "REVENUE", "Short uppercase label."),
                                                 _sub("value", "$48.2M", "The big number."),
                                                 _sub("note", "+6% QoQ", "Movement/context line.", required=False)])},
        "two_col": {"slide_index": 5, "purpose": "Two labelled columns (risks/mitigations, before/after, pros/cons).",
                    "fields": [H("Top risks and how we contain them"),
                               F("left_label", "RISKS", "Left column label (uppercase)."),
                               F("left_items", "Key-person dependency on night shift", "3–5 left bullets.", type="list"),
                               F("right_label", "MITIGATIONS", "Right column label (uppercase)."),
                               F("right_items", "Cross-training rotation from November", "One right bullet per left bullet.", type="list")]},
        "team": {"slide_index": 6, "purpose": "Team/owner rows (2–5): initials avatar, name, role, note.",
                 "fields": [H("Who is in the room")],
                 "items": _item_spec("items", "ITEM", dy=Inches(0.95), mx=5, ramp=ramp_hex,
                                     subfields=[_sub("initials", "JM", "1–3 letter initials for the avatar circle.", required=False),
                                                _sub("name", "Jane Mokoena", "Person or team name."),
                                                _sub("role", "Engagement Lead", "Role/responsibility."),
                                                _sub("note", "Decision owner for scope changes", "Optional context.", required=False)])},
        "timeline": {"slide_index": 7, "purpose": "Horizontal milestones (2–5): label, date, note.",
                     "fields": [H("The road to cutover")],
                     "items": _item_spec("items", "ITEM", dx=Inches(2.35), mx=5, ramp=ramp_hex, center=True,
                                         subfields=[_sub("label", "Pilot live", "Milestone name."),
                                                    _sub("date", "End September", "Timing."),
                                                    _sub("text", "First depot on the new app", "Optional one-liner.", required=False)])},
        "chart": {"slide_index": 8, "purpose": "Native editable chart built from data (column/bar/line/pie/area/doughnut).",
                  "fields": [H("Throughput grew every month"),
                             F("chart", {"chart_type": "column", "categories": ["Jul", "Aug", "Sep"],
                                         "series": [{"name": "Actual", "values": [82, 91, 103]}]},
                               'Chart spec: {"chart_type", "categories", "series":[{"name","values"}]}. Values numeric, one per category.', type="chart"),
                             F("takeaway", "Sept exceeded plan by 8% — capacity holds through peak", "One-line takeaway under the chart.", required=False)]},
        "table": {"slide_index": 9, "purpose": "Native table from columns/rows.",
                  "fields": [H("Cost per order by site"),
                             F("table", {"columns": ["Site", "Cost/order", "Trend"],
                                         "rows": [["Columbus", "$2.41", "▼ 3%"]]},
                               'Table spec: {"columns": [...], "rows": [[...], ...]} — every row matches columns length.', type="table"),
                             F("note", "Source: finance close, Oct 2026.", "Optional source/footnote.", required=False)]},
        "image": {"slide_index": 10, "purpose": "Full-width image with caption — REQUIRES an image path; omit this slide if you have no image.",
                  "fields": [H("The loading bay at shift change"),
                             F("image", "", "PNG/JPG path. REQUIRED — never ship this slide without a real image.", type="image"),
                             F("caption", "Source: site walkthrough, July 2026.", "Small caption: what it shows + source.")]},
        "quote": {"slide_index": 11, "purpose": "Big statement/quote on a dark slide.",
                  "fields": [F("quote", "This quarter proved the network can absorb peak volume without adding sites.", "The statement (1–2 sentences)."),
                             F("attribution", "Chief Operating Officer, September review", "Who said it / where it's from.", required=False)]},
        "section": {"slide_index": 12, "purpose": "Dark chapter divider.",
                    "fields": [F("marker", "01", "Big marker, e.g. 01 / 02 / A. Omit for no number.", required=False),
                               H("What the audit showed"),
                               F("note", "Three findings drive the recommendations.", "One muted line under the heading.", required=False)]},
    }
    body = {
        "anchor_index": 0,
        "min": 1, "max": 16,
        "chart_style": {"colors": ramp_hex, "font": st.font,
                        "table_head": "#" + str(st.dark), "table_band": "#" + str(st.panel),
                        "ink": "#" + str(st.ink)},
        "types": body_types,
    }
    slides = [
        {"index": 0, "name": "Cover", "purpose": "Eyebrow, title, subtitle, author line."},
        *[{"index": d["slide_index"], "name": t, "body_type": t, "purpose": d["purpose"]}
          for t, d in sorted(body_types.items(), key=lambda kv: kv[1]["slide_index"])],
        {"index": 13, "name": "Closing", "purpose": "Single ask + next steps."},
    ]
    return {"fields": fields, "slide_groups": [], "slides": slides, "body": body}


def build_flex_deck(prs, st, png):
    return _flex_spec(prs, st, png)


TEMPLATES = {
    "exec_update": (build_exec_update, "Executive/quarterly business update (QBR). Flex preset — scaffold suggests stat cards, section, chart, topic bullets, risks two-col; any body type can be added/reordered."),
    "flex_deck": (build_flex_deck, "Universal composable deck: pick ANY mix/order of body slides — agenda, bullets, numbered steps, stat cards, two-col, team, timeline, native charts, tables, image evidence, quote, section dividers. Use when no preset narrative fits."),
    "project_kickoff": (build_project_kickoff, "Project kickoff. Flex preset — scaffold suggests agenda, objectives, numbered approach, team, milestone timeline, comms two-col; any body type can be added/reordered."),
    "proposal": (build_proposal, "Client proposal. Flex preset — scaffold suggests problem bullets, numbered approach, scope, team, investment table, quote; any body type can be added/reordered."),
    "report_out": (build_report_out, "Project results report-out. Flex preset — scaffold suggests summary, findings section+bullets, numbered recommendations, quote; add image evidence slides when the user supplies files."),
}


# ---------------- Build + register ----------------

def media_part_for(pptx_path: Path, image_bytes: bytes) -> str:
    """Locate the package media part holding exactly these bytes (the image slot)."""
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if name.startswith("ppt/media/") and z.read(name) == image_bytes:
                return name
    return ""


def build_one(name: str, brand: dict, registry: Path, owner: str, created: str) -> Path:
    builder, description = TEMPLATES[name]
    st = Style(brand)
    png = placeholder_chart_png(st)

    prs = Presentation()
    prs.slide_width = int(SW)
    prs.slide_height = int(SH)
    spec = builder(prs, st, png)
    if isinstance(spec, list):                       # older builders return just fields
        spec = {"fields": spec, "slide_groups": [], "slides": []}
    fields = spec["fields"]
    slide_groups = spec.get("slide_groups", [])
    slides_meta = spec.get("slides", [])

    dest = registry / "_builtin" / name
    dest.mkdir(parents=True, exist_ok=True)
    tpl = dest / "template.pptx"
    prs.save(str(tpl))

    # Wire image-slot fields (global, per-group AND per-body-type) to their media part.
    slot = media_part_for(tpl, png)
    body = spec.get("body")
    body_fields = [bf for d in (body or {}).get("types", {}).values()
                   for bf in d.get("fields", [])]
    for f in fields + [gf for g in slide_groups for gf in g["fields"]] + body_fields:
        if f["type"] == "image" and not f["media_part"]:
            f["media_part"] = slot

    import hashlib
    manifest = {
        "template_id": f"_builtin/{name}",
        "client": "_builtin",
        "doc_type": name,
        "format": "pptx",
        "template_file": "template.pptx",
        "source_file": f"generated by building-powerpoint-decks/scripts/build_template_library.py (brand '{brand['name']}')",
        "version": "1.0.0",
        "owner": owner,
        "created": created,
        "changelog": [f"{created}: generated with brand pack '{brand['name']}'"],
        "description": description,
        "fields": fields,
        "row_groups": [],
        "slide_groups": slide_groups,
        "slides": slides_meta,
        "source_terms": SOURCE_TERMS,
        # validate.py fails a delivered deck whose media still contains this
        # placeholder image — an image slot that was never swapped must not ship.
        "placeholder_media_sha1": [hashlib.sha1(png).hexdigest()],
    }
    if body:
        manifest["body"] = body
    (dest / "manifest.json").write_text(json.dumps(manifest, indent=2, ensure_ascii=False),
                                        encoding="utf-8")
    n_slides = len(prs.slides._sldIdLst)
    n_fields = len(fields) + sum(len(g["fields"]) for g in slide_groups) + len(body_fields)
    groups_note = (" (+" + ", ".join(f"{g['name']} ×{g.get('min', 1)}–{g.get('max') or '∞'}"
                                     for g in slide_groups) + ")") if slide_groups else ""
    if body:
        groups_note += f" (body: {len(body['types'])} composable types)"
    print(f"Registered _builtin/{name}: {n_slides} slides{groups_note}, {n_fields} fields -> {dest}")
    return dest


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--brand", default="default", help="Brand pack name under brands/ (or a path)")
    ap.add_argument("--only", choices=sorted(TEMPLATES), help="Build a single template")
    ap.add_argument("--registry", default=os.environ.get("TEMPLATE_REGISTRY", str(DEFAULT_REGISTRY)),
                    help="Registry root (default: building-document-templates/registry or $TEMPLATE_REGISTRY)")
    ap.add_argument("--owner", default="", help="Recorded in each manifest")
    ap.add_argument("--created", default=date.today().isoformat())
    args = ap.parse_args()

    brand = load_brand(args.brand)
    registry = Path(args.registry)
    names = [args.only] if args.only else sorted(TEMPLATES)
    for name in names:
        build_one(name, brand, registry, args.owner, args.created)
    print(f"\nDiscover:  python ../building-document-templates/scripts/registry.py list")
    print(f"Fill:      scaffold -> fill.py -> validate.py -> render_pages.py (see the SKILL)")


if __name__ == "__main__":
    main()
