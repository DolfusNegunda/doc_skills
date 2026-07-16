"""Gate a filled document before it ships.

Fails loudly (non-zero exit) with specific messages so problems are fixable:
  * No leftover ``{{ tag }}`` / ``{% ... %}`` anywhere (body, tables, headers,
    footers, slides, notes) — the classic "shipped with a placeholder in it" bug.
  * The output still has structure — non-empty, and (vs. the template, if given)
    the same slide count / no lost sections.

Prints a JSON report; exit 0 only when status == "OK".

Usage:
    python scripts/validate.py out/acme-q4.docx
    python scripts/validate.py out/deck.pptx --template registry/acme/board-deck/template.pptx
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import common as C


def _residue_errors(texts, source_terms):
    """Flag source-exemplar content that should have been replaced by the fill.

    'No leftover placeholders' is NOT the same as 'successfully reused': a fill can
    keep the source's client/project/dates and still have zero {{ tags }}. Given the
    terms the canonical template was learned from, fail if any survive in the output.
    Case-insensitive substring match; short/ambiguous terms (<4 chars) are skipped."""
    if not source_terms:
        return [], []
    blob = "\n".join(texts).lower()
    found = sorted({t for t in source_terms
                    if len(t.strip()) >= 4 and t.strip().lower() in blob})
    if found:
        return [f"Source-exemplar content still present (not replaced): {found}"], found
    return [], []


def docx_report(path: Path, template: Path | None, source_terms):
    from docx import Document
    doc = Document(str(path))
    texts = [C.para_text(p) for p in C.iter_docx_paragraphs(doc)]
    texts += C.property_texts(path)   # cover-page / data-bound property leaves
    texts += [d["text"] for d in C.smartart_texts(path)]   # SmartArt diagram text
    errors, checks = [], {}

    leftover = sorted({m.group(0) for t in texts for m in C.ANY_TAG_RE.finditer(t)})
    if leftover:
        errors.append(f"Unfilled template tags remain: {leftover}")
    checks["nonempty_paragraphs"] = sum(1 for t in texts if t.strip())
    if checks["nonempty_paragraphs"] == 0:
        errors.append("Document has no text content.")

    res_err, res_found = _residue_errors(texts, source_terms)
    errors += res_err
    checks["source_residue"] = res_found
    # Informational: SmartArt/chart text the fill can't touch — QA must eyeball these.
    checks["unsupported_objects"] = [f"{u['kind']}:{u['part']}" for u in C.iter_unsupported_objects(path)]

    if template and template.exists():
        tdoc = Document(str(template))
        t_sections = len(tdoc.sections)
        checks["sections"] = len(doc.sections)
        if len(doc.sections) != t_sections:
            errors.append(f"Section count changed: template {t_sections} -> output {len(doc.sections)}")
    elif template:
        errors.append(f"--template not found: {template} (structure check could not run)")
    return errors, checks


def pptx_report(path: Path, template: Path | None, source_terms, slide_groups=None):
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = [C.para_text(p) for p in C.iter_pptx_paragraphs(prs)]
    texts += C.property_texts(path)   # cover-page / data-bound property leaves
    texts += [d["text"] for d in C.smartart_texts(path)]   # SmartArt diagram text
    errors, checks = [], {}

    leftover = sorted({m.group(0) for t in texts for m in C.ANY_TAG_RE.finditer(t)})
    if leftover:
        errors.append(f"Unfilled template tags remain: {leftover}")
    checks["n_slides"] = len(prs.slides)
    if checks["n_slides"] == 0:
        errors.append("Deck has no slides.")

    res_err, res_found = _residue_errors(texts, source_terms)
    errors += res_err
    checks["source_residue"] = res_found
    # Informational: SmartArt/chart text the fill can't touch — QA must eyeball these.
    checks["unsupported_objects"] = [f"{u['kind']}:{u['part']}" for u in C.iter_unsupported_objects(path)]

    if template and template.exists():
        tprs = Presentation(str(template))
        t_slides = len(tprs.slides)
        groups = slide_groups or []
        if groups:
            # Repeatable/optional slides: each group's ONE template slide may
            # legitimately become min..max slides in the output.
            lo = t_slides + sum(g.get("min", 1) - 1 for g in groups)
            hi = t_slides + sum((g.get("max") or 99) - 1 for g in groups)
            checks["expected_slides"] = f"{lo}-{hi}"
            if not (lo <= checks["n_slides"] <= hi):
                errors.append(f"Slide count {checks['n_slides']} outside the template's "
                              f"expected range {lo}-{hi} (base {t_slides} + slide_groups).")
        elif checks["n_slides"] != t_slides:
            errors.append(f"Slide count changed: template {t_slides} -> output {checks['n_slides']} "
                          "(if this template has repeatable slide_groups, pass --manifest).")
    elif template:
        errors.append(f"--template not found: {template} (structure check could not run)")
    return errors, checks


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("file")
    ap.add_argument("--template", default=None, help="Compare structure against this template")
    ap.add_argument("--manifest", default=None,
                    help="Manifest to read source_terms from (source-residue check)")
    ap.add_argument("--source-terms", default=None,
                    help="Comma-separated source terms that must NOT survive the fill")
    args = ap.parse_args()

    path = Path(args.file)
    if not path.exists():
        sys.exit(f"file not found: {path}")
    template = Path(args.template) if args.template else None
    fmt = C.detect_format(path)

    # Source-residue terms + slide_groups: from --manifest and/or --source-terms.
    source_terms, slide_groups = [], []
    if args.manifest and Path(args.manifest).exists():
        man = C.load_manifest(Path(args.manifest))
        source_terms += man.get("source_terms", [])
        slide_groups = man.get("slide_groups", [])
    if args.source_terms:
        source_terms += [t.strip() for t in args.source_terms.split(",") if t.strip()]

    if fmt == "docx":
        errors, checks = docx_report(path, template, source_terms)
    else:
        errors, checks = pptx_report(path, template, source_terms, slide_groups)
    report = {
        "file": str(path),
        "format": fmt,
        "status": "OK" if not errors else "FAIL",
        "checks": checks,
        "errors": errors,
    }
    print(json.dumps(report, indent=2, ensure_ascii=False))
    sys.exit(0 if not errors else 1)


if __name__ == "__main__":
    main()
