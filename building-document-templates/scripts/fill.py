"""Fill a registered template with content and emit a finished document.

Resolves a template from the gallery by client + doc-type (or an explicit path),
reads a data JSON keyed by the manifest's field names, and renders:

  * text  fields -> the value is dropped in wherever ``{{ field }}`` appears.
  * list  fields -> if the placeholder is alone on a paragraph (a bullet/row), that
                    paragraph is *duplicated* once per item so you get real bullets,
                    not one line with commas. Otherwise the items are joined inline.

The template's layout, styles, logos and masters are never touched — only the
placeholder text changes. Optionally exports the result to PDF.

Usage:
    python scripts/fill.py --client acme --doc-type quarterly-review \
        --data content.json --out out/acme-q4.docx [--export-pdf]
    python scripts/fill.py --template path/to/template.pptx --manifest path/to/manifest.json \
        --data content.json --out out/deck.pptx
"""
from __future__ import annotations

import argparse
import copy
import json
import subprocess
import sys
from pathlib import Path

import common as C


def _wrap_docx(elem, parent):
    from docx.text.paragraph import Paragraph
    return Paragraph(elem, parent)


def _wrap_pptx(elem, parent):
    from pptx.text.text import _Paragraph
    return _Paragraph(elem, parent)


def expand_list(paragraph, tag, items, wrap):
    """Duplicate a bullet/row paragraph once per list item, preserving order.

    The original paragraph becomes item[0]; each further item is a deep copy of the
    pristine (still-tagged) paragraph inserted right after the previous one.
    """
    if not items:
        paragraph._p.getparent().remove(paragraph._p)
        return
    template_p = copy.deepcopy(paragraph._p)   # pristine copy (still holds the tag)
    parent = paragraph._parent
    C.replace_in_paragraph(paragraph, tag, str(items[0]))
    anchor = paragraph._p
    for item in items[1:]:
        clone = copy.deepcopy(template_p)
        anchor.addnext(clone)                  # insert immediately after anchor...
        anchor = clone                         # ...advance so the next stays in order
        C.replace_in_paragraph(wrap(clone, parent), tag, str(item))


def _fill_row_element(tr, columns, item):
    """Replace each column field's {{ tag }} inside one <w:tr> clone from `item`
    (a dict keyed by the column field names). Preserves each cell's formatting."""
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph
    for p in tr.iter(qn("w:p")):
        para = Paragraph(p, None)
        for col in columns:
            val = item.get(col, "") if isinstance(item, dict) else ""
            C.replace_in_paragraph(para, C.placeholder(col), "" if val is None else str(val))


def expand_row_groups(doc, row_groups, data):
    """Clone a table's template ROW once per data item — real repeating rows.

    A row_group is {name, columns:[field,...]}. The template row is the <w:tr>
    whose cells hold those column tags. `data[name]` is a list of dicts keyed by
    the column field names. Empty/missing -> the template row is removed (header
    stays). This is what paragraph-level list expansion cannot do (it stacks
    paragraphs inside one cell instead of adding rows)."""
    import copy
    for g in row_groups:
        name, columns = g["name"], g["columns"]
        items = data.get(name) or []
        if isinstance(items, dict):
            items = [items]
        tag0 = C.placeholder(columns[0])
        template_tr = None
        for tbl in C.all_docx_tables(doc):
            for row in tbl.rows:
                if any(tag0 in c.text for c in row.cells):
                    template_tr = row._tr
                    break
            if template_tr is not None:
                break
        if template_tr is None:
            continue
        if not items:
            template_tr.getparent().remove(template_tr)
            continue
        pristine = copy.deepcopy(template_tr)     # still-tagged copy for the clones
        _fill_row_element(template_tr, columns, items[0])
        anchor = template_tr
        for item in items[1:]:
            clone = copy.deepcopy(pristine)
            anchor.addnext(clone)
            anchor = clone
            _fill_row_element(clone, columns, item)


def _is_empty(value):
    """Absent/None/blank/empty-list all count as 'no value supplied'."""
    return value is None or value == "" or value == []


# ---------------- PPTX slide groups (repeatable / optional slides) ----------------
# A manifest slide_group marks ONE slide as a repeating unit:
#   {"name", "slide_index", "min", "max", "purpose", "fields": [field-specs]}
# data[name] is a LIST of dicts (one slide per entry). The engine clones the
# designer slide byte-for-byte per entry — same layout, background, art, theme —
# and fills each clone from its entry. Zero entries: min 0 removes the slide;
# a required group keeps its tags so validate.py fails loudly.

def _slide_paragraphs(slide):
    """Paragraphs of ONE slide (text frames, grouped shapes, table cells)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(sh.shapes)
            elif getattr(sh, "has_text_frame", False):
                yield from sh.text_frame.paragraphs
            elif getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    for cell in row.cells:
                        yield from cell.text_frame.paragraphs

    yield from walk(slide.shapes)


def _sld_el(prs, slide):
    from pptx.oxml.ns import qn  # noqa: F401  (namespace map side effect)
    for el in prs.slides._sldIdLst:
        if int(el.get("id")) == slide.slide_id:
            return el
    raise ValueError(f"slide id {slide.slide_id} not in sldIdLst")


def delete_pptx_slide(prs, slide):
    from pptx.oxml.ns import qn
    el = _sld_el(prs, slide)
    prs.part.drop_rel(el.get(qn("r:id")))
    el.getparent().remove(el)


def _move_slide_after(prs, dest, anchor):
    lst = prs.slides._sldIdLst
    d_el, a_el = _sld_el(prs, dest), _sld_el(prs, anchor)
    lst.remove(d_el)
    a_el.addnext(d_el)


def clone_pptx_slide(prs, source):
    """Byte-faithful duplicate of a slide: same layout, background, art, images.

    python-pptx has no public duplicate API. Approach: add_slide() on the SAME
    layout (correct part scaffolding + theme inheritance), swap in a deep copy of
    the source's entire <p:cSld> (background + every shape), then re-create the
    source's relationships on the new part and remap the rIds inside the copied
    XML (two passes via placeholder tokens so colliding ids can't chain-rewrite).
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.oxml.ns import qn

    dest = prs.slides.add_slide(source.slide_layout)
    # python-pptx names new slide parts len(slides)+1, which COLLIDES with an
    # existing part after a slide deletion (duplicate zip entries corrupt the
    # deck) — rename the clone to the first genuinely free slide partname.
    from pptx.opc.packuri import PackURI
    taken = {str(p.partname) for p in prs.part.package.iter_parts()} - {str(dest.part.partname)}
    n = 1
    while f"/ppt/slides/slide{n}.xml" in taken:
        n += 1
    dest.part.partname = PackURI(f"/ppt/slides/slide{n}.xml")
    d_el, s_el = dest._element, source._element
    d_el.remove(d_el.find(qn("p:cSld")))
    d_el.insert(0, copy.deepcopy(s_el.find(qn("p:cSld"))))
    # add_slide() already touched dest.shapes, caching a wrapper around the
    # spTree we just replaced — drop the caches so fills hit the copied shapes.
    dest.__dict__.pop("shapes", None)
    dest.__dict__.pop("placeholders", None)

    mapping = {}
    for rid, rel in source.part.rels.items():
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        new_rid = (dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                   if rel.is_external
                   else dest.part.relate_to(rel.target_part, rel.reltype))
        mapping[rid] = new_rid
    tokens = {old: f"__RIDTMP{i}__" for i, old in enumerate(mapping)}
    for el in d_el.iter():
        for attr, val in el.attrib.items():
            if val in tokens:
                el.set(attr, tokens[val])
    inverse = {tokens[old]: new for old, new in mapping.items()}
    for el in d_el.iter():
        for attr, val in el.attrib.items():
            if val in inverse:
                el.set(attr, inverse[val])
    return dest


def _swap_instance_image(slide, media_part, image_path):
    """Give THIS slide instance its own image: add a new image part and repoint
    the picture that references `media_part` (clones share the original part, so
    a package-level media swap would change every clone at once)."""
    from pptx.oxml.ns import qn
    p = Path(image_path)
    if not p.exists():
        print(f"WARNING: image not found: {p} — slide keeps its placeholder visual.")
        return False
    want = "/" + str(media_part).lstrip("/")
    target_rids = {rid for rid, rel in slide.part.rels.items()
                   if not rel.is_external and str(rel.target_part.partname) == want}
    if not target_rids:
        return False
    image_part, new_rid = slide.part.get_or_add_image_part(str(p))
    hit = False
    for blip in slide._element.iter(qn("a:blip")):
        if blip.get(qn("r:embed")) in target_rids:
            blip.set(qn("r:embed"), new_rid)
            hit = True
    return hit


def expand_slide_groups(prs, groups, data):
    """Clone/drop whole slides per the manifest's slide_groups; fill each clone
    from its entry. Returns the set of missing required group/field names."""
    missing = set()
    # Descending template index so earlier group indices stay valid while we
    # insert/delete around later ones.
    for g in sorted(groups, key=lambda g: g["slide_index"], reverse=True):
        name = g["name"]
        entries = data.get(name)
        if isinstance(entries, dict):
            entries = [entries]
        entries = entries or []
        mn, mx = g.get("min", 1), g.get("max")
        slides = list(prs.slides)
        idx = g["slide_index"]
        if idx >= len(slides):
            print(f"WARNING: slide group '{name}' points at slide {idx + 1} but the deck "
                  f"has {len(slides)} slides — group skipped.")
            continue
        source = slides[idx]

        if not entries:
            if mn == 0:
                delete_pptx_slide(prs, source)   # optional slide, no content -> gone
            else:
                missing.add(name)                # tags stay -> validate.py fails loudly
            continue
        if mx and len(entries) > mx:
            print(f"WARNING: slide group '{name}': {len(entries)} entries exceeds max {mx} — "
                  "building them all; check the rendered deck for pacing.")

        instances, anchor = [source], source
        for _ in entries[1:]:
            clone = clone_pptx_slide(prs, source)
            _move_slide_after(prs, clone, anchor)
            instances.append(clone)
            anchor = clone

        for slide, entry in zip(instances, entries):
            text_fields = [f for f in g["fields"] if f.get("type", "text") != "image"]
            paras = list(_slide_paragraphs(slide))
            miss, types = fill_paragraphs(paras, entry, text_fields, _wrap_pptx)
            missing |= {f"{name}.{m}" for m in miss}
            paras = list(_slide_paragraphs(slide))   # re-collect after list expansion
            for f in text_fields:
                if types[f["name"]] == "list" or f["name"] in miss:
                    continue
                val = entry.get(f["name"], "")
                tag_ = C.placeholder(f["name"])
                for p in paras:
                    C.replace_in_paragraph(p, tag_, "" if val is None else str(val))
            for f in g["fields"]:
                if f.get("type") != "image":
                    continue
                val = entry.get(f["name"])
                if val and not _swap_instance_image(slide, f.get("media_part", ""), val):
                    print(f"WARNING: image slot for '{name}.{f['name']}' not found on its slide.")
    return missing


# ---------------- Composable deck body (manifest key: "body") ----------------
# A body-enabled template ships ONE source slide per body TYPE (bullets, chart,
# team, ...). The data's "body" is an ordered list of typed entries; the engine
# clones the matching source per entry (any order, any mix), fills it, and
# finally deletes every source slide — so the deck contains exactly the slides
# the content asked for, all byte-faithful to the designer originals.

def _shape_by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def _first_off(el):
    from pptx.oxml.ns import qn
    return next(el.iter(qn("a:off")), None)


def _insert_native_chart(slide, spec, style, field_name):
    """Replace the CHART_SLOT marker with a real, editable PowerPoint chart."""
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Pt

    marker = _shape_by_name(slide, "CHART_SLOT")
    if marker is None:
        print(f"WARNING: no CHART_SLOT marker for '{field_name}' — chart skipped.")
        return False
    kinds = {"bar": XL_CHART_TYPE.BAR_CLUSTERED, "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
             "line": XL_CHART_TYPE.LINE, "pie": XL_CHART_TYPE.PIE,
             "area": XL_CHART_TYPE.AREA, "doughnut": XL_CHART_TYPE.DOUGHNUT}
    kind = str(spec.get("chart_type", "column")).lower()
    if kind not in kinds:
        sys.exit(f"body chart '{field_name}': unknown chart_type {spec.get('chart_type')!r} "
                 f"— valid: {', '.join(sorted(kinds))}")
    cats = spec.get("categories") or []
    series = spec.get("series") or []
    if not cats or not series:
        sys.exit(f"body chart '{field_name}': needs 'categories' and 'series' "
                 '[{"name", "values"}] with values matching categories length.')
    cd = CategoryChartData()
    cd.categories = [str(c) for c in cats]
    for s in series:
        vals = s.get("values") or []
        if len(vals) != len(cats):
            sys.exit(f"body chart '{field_name}': series '{s.get('name', '')}' has "
                     f"{len(vals)} values for {len(cats)} categories.")
        cd.add_series(str(s.get("name", "")), tuple(float(v) for v in vals))

    x, y, cx, cy = marker.left, marker.top, marker.width, marker.height
    marker._element.getparent().remove(marker._element)
    chart = slide.shapes.add_chart(kinds[kind], x, y, cx, cy, cd).chart

    colors = [RGBColor.from_string(c.lstrip("#")) for c in (style or {}).get("colors", [])] \
        or [RGBColor.from_string(c) for c in ("2563EB", "0E7490", "38BDF8", "64748B", "94A3B8")]
    chart.has_title = False
    chart.font.size = Pt(12)
    if (style or {}).get("font"):
        chart.font.name = style["font"]
    if kind in ("pie", "doughnut"):
        pts = chart.plots[0].series[0].points
        for i, pt in enumerate(pts):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        chart.has_legend = True
    else:
        for i, ser in enumerate(chart.series):
            if kind in ("line",):
                ser.format.line.color.rgb = colors[i % len(colors)]
                ser.format.line.width = Pt(2.5)
                ser.smooth = False
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = colors[i % len(colors)]
        chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    return True


def _insert_native_table(slide, spec, style, field_name):
    """Replace the TABLE_SLOT marker with a real table styled from the brand."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    marker = _shape_by_name(slide, "TABLE_SLOT")
    if marker is None:
        print(f"WARNING: no TABLE_SLOT marker for '{field_name}' — table skipped.")
        return False
    cols = spec.get("columns") or []
    rows = spec.get("rows") or []
    if not cols or not rows:
        sys.exit(f"body table '{field_name}': needs 'columns' [..] and 'rows' [[..]].")
    for i, r in enumerate(rows):
        if len(r) != len(cols):
            sys.exit(f"body table '{field_name}': row {i + 1} has {len(r)} cells "
                     f"for {len(cols)} columns.")
    x, y, cx, cy = marker.left, marker.top, marker.width, marker.height
    marker._element.getparent().remove(marker._element)
    tbl = slide.shapes.add_table(len(rows) + 1, len(cols), x, y, cx, cy).table
    # Compact rows (add_table stretches rows to fill the slot's height).
    from pptx.util import Emu
    tbl.rows[0].height = Emu(457200)          # 0.5" header
    for r in list(tbl.rows)[1:]:
        r.height = Emu(402336)                # 0.44" body rows
    s = style or {}
    head = RGBColor.from_string(s.get("table_head", "1E293B").lstrip("#"))
    band = RGBColor.from_string(s.get("table_band", "F1F5F9").lstrip("#"))
    ink = RGBColor.from_string(s.get("ink", "0F172A").lstrip("#"))
    font = s.get("font")
    for c, name in enumerate(cols):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = head
        cell.text = str(name)
        for p in cell.text_frame.paragraphs:
            for r_ in p.runs:
                r_.font.size, r_.font.bold = Pt(12), True
                r_.font.color.rgb = RGBColor.from_string("FFFFFF")
                if font:
                    r_.font.name = font
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = band if ri % 2 == 0 else RGBColor.from_string("FFFFFF")
            cell.text = "" if val is None else str(val)
            for p in cell.text_frame.paragraphs:
                for r_ in p.runs:
                    r_.font.size = Pt(11.5)
                    r_.font.color.rgb = ink
                    if font:
                        r_.font.name = font
    return True


def _expand_items(slide, entry, item_spec, missing, label):
    """Clone the ITEM template shape/group once per item — the Rev Sci pattern:
    each row/card is a designed shape group; entry count = row count; colors walk
    the brand ramp per instance."""
    from pptx.oxml.ns import qn

    tmpl_shape = _shape_by_name(slide, item_spec.get("shape", "ITEM"))
    field = item_spec.get("field", "items")
    items = entry.get(field) or []
    if isinstance(items, dict):
        items = [items]
    items = [{"title": i} if isinstance(i, str) else i for i in items]
    if tmpl_shape is None:
        if items:
            print(f"WARNING: {label}: no '{item_spec.get('shape', 'ITEM')}' template "
                  "shape on the slide — items skipped.")
        return
    tmpl_el = tmpl_shape._element
    if not items:
        missing.add(f"{label}.{field}")
        return
    mx = item_spec.get("max")
    if mx and len(items) > mx:
        print(f"WARNING: {label}: {len(items)} items exceeds max {mx} — building all; "
              "check the render for overflow.")

    pristine = copy.deepcopy(tmpl_el)
    off0 = _first_off(tmpl_el)
    x0, y0 = int(off0.get("x")), int(off0.get("y"))
    dx, dy = int(item_spec.get("dx", 0)), int(item_spec.get("dy", 0))
    ramp = [c.lstrip("#").upper() for c in item_spec.get("ramp", [])]
    subfields = item_spec.get("subfields", [])
    parent_tree = tmpl_el.getparent()

    instances = [tmpl_el]
    anchor = tmpl_el
    for _ in items[1:]:
        clone = copy.deepcopy(pristine)
        anchor.addnext(clone)
        anchor = clone
        instances.append(clone)

    if item_spec.get("center") and dx:
        total = (len(items) - 1) * dx + tmpl_shape.width
        x0 = max(0, (item_spec.get("span", 12192000) - total) // 2 +
                 int(item_spec.get("span_left", 0)))

    for i, (el, item) in enumerate(zip(instances, items)):
        off = _first_off(el)
        off.set("x", str(x0 + i * dx))
        off.set("y", str(y0 + i * dy))
        if ramp:
            want = ramp[i % len(ramp)]
            for clr in el.iter(qn("a:srgbClr")):
                if clr.get("val", "").upper() == ramp[0]:
                    clr.set("val", want)
        paras = [_wrap_pptx(p, None) for p in el.iter(qn("a:p"))]
        for p in paras:                       # auto row number ({{ item._n }})
            C.replace_in_paragraph(p, C.placeholder("item._n"), str(i + 1))
        for sf in subfields:
            name = sf["name"]
            tag_ = C.placeholder(f"item.{name}")
            val = item.get(name)
            if _is_empty(val) and sf.get("required", True):
                missing.add(f"{label}.{field}[{i}].{name}")
                continue
            for p in paras:
                C.replace_in_paragraph(p, tag_, "" if val is None else str(val))
    _ = parent_tree  # tree mutation done in place


def expand_body(prs, spec, data):
    """Compose the deck from typed body entries (any order, any mix). Returns the
    set of missing required names; leaves source-slide tags in place when the body
    itself is missing so validate.py fails loudly."""
    missing = set()
    types = spec.get("types", {})
    entries = data.get("body")
    if isinstance(entries, dict):
        entries = [entries]
    entries = entries or []
    valid = sorted(types)

    bad = [f'body[{i}]: unknown type {e.get("type")!r}' for i, e in enumerate(entries)
           if e.get("type") not in types]
    if bad:
        sys.exit("\n".join(bad) + f"\nValid body types: {', '.join(valid)}")

    mn, mx = spec.get("min", 1), spec.get("max")
    if len(entries) < mn:
        missing.add("body")
        return missing            # sources (and their tags) stay -> validate fails
    if mx and len(entries) > mx:
        print(f"WARNING: body has {len(entries)} slides, max is {mx} — building all; "
              "consider splitting the deck.")

    slides = list(prs.slides)
    sources = {t: slides[d["slide_index"]] for t, d in types.items()
               if d["slide_index"] < len(slides)}
    anchor = slides[spec["anchor_index"]]
    chart_style = spec.get("chart_style", {})

    for i, entry in enumerate(entries):
        t = entry["type"]
        tdef = types[t]
        label = f"body[{i}]({t})"
        slide = clone_pptx_slide(prs, sources[t])
        _move_slide_after(prs, slide, anchor)
        anchor = slide

        text_fields = [f for f in tdef.get("fields", [])
                       if f.get("type", "text") not in ("image", "chart", "table")]
        paras = list(_slide_paragraphs(slide))
        miss, ftypes = fill_paragraphs(paras, entry, text_fields, _wrap_pptx)
        missing |= {f"{label}.{m}" for m in miss}
        paras = list(_slide_paragraphs(slide))
        for f in text_fields:
            if ftypes[f["name"]] == "list" or f["name"] in miss:
                continue
            val = entry.get(f["name"], "")
            tag_ = C.placeholder(f["name"])
            for p in paras:
                C.replace_in_paragraph(p, tag_, "" if val is None else str(val))

        for f in tdef.get("fields", []):
            name, ftype = f["name"], f.get("type", "text")
            val = entry.get(name)
            if ftype == "image":
                if _is_empty(val):
                    if f.get("required", True):
                        missing.add(f"{label}.{name}")
                elif not _swap_instance_image(slide, f.get("media_part", ""), val):
                    print(f"WARNING: image slot for '{label}.{name}' not found.")
            elif ftype == "chart":
                if _is_empty(val):
                    missing.add(f"{label}.{name}")
                else:
                    _insert_native_chart(slide, val, chart_style, f"{label}.{name}")
            elif ftype == "table":
                if _is_empty(val):
                    missing.add(f"{label}.{name}")
                else:
                    _insert_native_table(slide, val, chart_style, f"{label}.{name}")

        if tdef.get("items"):
            _expand_items(slide, entry, tdef["items"], missing, label)

    for t in sorted(sources, key=lambda t: types[t]["slide_index"], reverse=True):
        delete_pptx_slide(prs, sources[t])
    return missing


def fill_paragraphs(paragraphs, data, fields, wrap):
    """Apply the manifest's fields to every paragraph.

    A REQUIRED field with no value is left as its ``{{ tag }}`` on purpose — it is
    NOT blanked. That way `validate.py`'s leftover-tag check fails the document
    instead of silently shipping a blank cell/bullet. An OPTIONAL field with no
    value fills blank (text) or drops the line (list). Returns (missing_required set,
    types dict).
    """
    types = {f["name"]: f.get("type", "text") for f in fields}
    required = {f["name"]: f.get("required", True) for f in fields}
    missing = {f["name"] for f in fields
               if required[f["name"]] and _is_empty(data.get(f["name"]))}

    # List expansion first (it restructures paragraphs), then plain text replacement.
    for name, ftype in types.items():
        if ftype != "list" or name in missing:   # missing+required -> leave the tag
            continue
        tag = C.placeholder(name)
        items = data.get(name) or []
        if isinstance(items, str):
            items = [items]
        # Snapshot: expansion mutates the tree, so collect target paragraphs first.
        targets = [p for p in paragraphs if C.para_text(p).strip() == tag]
        if targets:
            for p in targets:
                expand_list(p, tag, items, wrap)
        else:
            # Placeholder shares a line with other text -> join inline.
            joined = "; ".join(str(i) for i in items)
            for p in paragraphs:
                C.replace_in_paragraph(p, tag, joined)

    return missing, types


def render(fmt, template_path, data, manifest, out_path):
    fields = manifest["fields"]
    row_groups = manifest.get("row_groups", [])
    if fmt == "docx":
        from docx import Document
        doc = Document(str(template_path))
        # Row-group expansion FIRST — it clones <w:tr> rows, restructuring tables
        # before the paragraph-level passes run over the (new) paragraph set.
        if row_groups:
            expand_row_groups(doc, row_groups, data)
        paras = list(C.iter_docx_paragraphs(doc))
        wrap = _wrap_docx
        saver = doc
    else:
        from pptx import Presentation
        prs = Presentation(str(template_path))
        # Slide-group expansion FIRST — it clones/drops whole slides and fills
        # each instance from its own entry, before the global passes run.
        group_missing = set()
        if manifest.get("body"):
            group_missing |= expand_body(prs, manifest["body"], data)
        if manifest.get("slide_groups"):
            group_missing |= expand_slide_groups(prs, manifest["slide_groups"], data)
        paras = list(C.iter_pptx_paragraphs(prs))
        wrap = _wrap_pptx
        saver = prs

    missing, types = fill_paragraphs(paras, data, fields, wrap)
    if fmt == "pptx":
        missing |= group_missing

    # Text pass (re-collect because list expansion changed the paragraph set).
    if fmt == "docx":
        paras = list(C.iter_docx_paragraphs(saver))
    else:
        paras = list(C.iter_pptx_paragraphs(saver))
    for name, ftype in types.items():
        if ftype == "list" or name in missing:    # missing+required -> leave the tag
            continue
        val = data.get(name, "")
        tag = C.placeholder(name)
        for p in paras:
            C.replace_in_paragraph(p, tag, "" if val is None else str(val))

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    saver.save(str(out))

    # Fill tags that live in cover-page / data-bound property parts (docProps +
    # customXml). These drive the cover title/subtitle/date/author and are not in
    # the body runs, so they need a separate pass over the saved package.
    prop_pairs = []
    for f in fields:
        name = f["name"]
        if name in missing:                     # leave the tag so validate.py fails
            continue
        val = data.get(name)
        if f.get("type", "text") == "list":
            items = val if isinstance(val, list) else ([val] if val else [])
            val = "; ".join(str(i) for i in items)
        prop_pairs.append((C.placeholder(name), "" if val is None else str(val)))
    C.patch_property_parts(out, out, C.ordered_replacer(prop_pairs))

    # SmartArt text tags live in the diagram parts (data + cached drawing), not in body
    # runs — fill them there too (no-op where the template has no SmartArt tags).
    C.patch_smartart_parts(out, out, C.ordered_replacer(prop_pairs))

    # Swap image/logo slots: {media_part: new_image_path} from image-type fields whose
    # data value is a path. Unset image slots keep the original picture.
    image_map = {f["media_part"]: data[f["name"]]
                 for f in fields
                 if f.get("type") == "image" and f.get("media_part")
                 and data.get(f["name"])}
    if image_map:
        C.swap_media_parts(out, out, image_map)
    return missing


def export_pdf(docx_or_pptx: Path):
    """Best-effort PDF export via LibreOffice; prints guidance if unavailable."""
    out_dir = docx_or_pptx.resolve().parent
    for exe in ("libreoffice", "soffice"):
        try:
            subprocess.run([exe, "--headless", "--convert-to", "pdf",
                            "--outdir", str(out_dir), str(docx_or_pptx)],
                           check=True, capture_output=True)
            print(f"PDF exported: {docx_or_pptx.with_suffix('.pdf')}")
            return
        except (FileNotFoundError, subprocess.CalledProcessError):
            continue
    print("NOTE: LibreOffice not found — open the file and Save As PDF, or install "
          "LibreOffice for headless export.")


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--client")
    ap.add_argument("--doc-type")
    ap.add_argument("--template", help="Explicit template path (bypasses the registry)")
    ap.add_argument("--manifest", help="Manifest for an explicit --template")
    ap.add_argument("--data", required=True, help="JSON of field_name -> value")
    ap.add_argument("--out", required=True)
    ap.add_argument("--export-pdf", action="store_true")
    args = ap.parse_args()

    if args.template:
        template_path = Path(args.template)
        man_path = Path(args.manifest) if args.manifest else template_path.parent / "manifest.json"
        manifest = C.load_manifest(man_path)
    elif args.client and args.doc_type:
        template_path, manifest = C.find_template(args.client, args.doc_type)
    else:
        sys.exit("Provide --client/--doc-type, or --template (+--manifest).")

    data = json.loads(Path(args.data).read_text(encoding="utf-8"))
    fmt = manifest["format"]

    # Slides don't reflow: a value much longer than the field's example is the main
    # overflow risk (there is no autofit). Warn early; the render/vision pass decides.
    if fmt == "pptx":
        for f in manifest["fields"]:
            if f.get("type", "text") != "text":
                continue
            ex, val = f.get("example") or "", data.get(f["name"])
            if ex and isinstance(val, str) and len(val) > max(60, int(len(ex) * 1.5)):
                print(f"WARNING: '{f['name']}' is {len(val)} chars (example: {len(ex)}) — "
                      f"slide text does not autofit; check the rendered pages for overflow.")

    missing = render(fmt, template_path, data, manifest, args.out)

    print(f"Filled {manifest['template_id']} -> {args.out}")
    if args.export_pdf:
        export_pdf(Path(args.out))
    if missing:
        # The tags for these fields were left in the document on purpose, so
        # validate.py will also fail. Exit non-zero so an automated pipeline stops.
        print(f"ERROR: no value supplied for required fields {sorted(missing)}; "
              f"their placeholders were left in {args.out} — do not ship it.")
        sys.exit(2)


if __name__ == "__main__":
    main()
