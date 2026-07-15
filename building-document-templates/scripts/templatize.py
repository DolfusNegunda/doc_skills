"""Turn a client's existing .docx/.pptx into a reusable template + manifest.

Templatizing preserves the file — layout, fonts, logos, styles, slide masters all
stay exactly as the client made them — and only swaps the *variable* text for
``{{ placeholder }}`` tags. It is a TWO-STEP, ASSISTED process because a single
example can't tell you which text is boilerplate and which changes each time
("Acme Q3 2026" — is Acme the client? is Q3 the quarter? both?):

  1. propose : read the file, extract candidate variable values with context and a
               heuristic guess (plus IMAGE slots, cover PROPERTY leaves, and ROW-GROUP
               table candidates), and write a proposal JSON for a human/agent to edit.
  2. build   : read the (edited) proposal, inject placeholders for everything marked
               keep="variable", expand confirmed tables into repeating ROW-GROUPS,
               and register the template + manifest in the gallery.

In the FAMILY model (the default), a template registers under
``registry/_families/<family>/`` — the one canonical template a document family
converges to — with variable-count tables as ``row_groups`` and the source
exemplar's identifying terms recorded as ``source_terms`` so a later fill is
checked for stale residue. ``--client``/``--doc-type`` still make a per-client
instance for the rare structural exception.

Usage:
    python scripts/templatize.py propose --file client.docx --out proposal.json
    # ...review proposal.json: set keep, rename fields, set type (text|list); mark
    #    variable-count tables as row-groups (keep='variable' on the candidate)...
    python scripts/templatize.py build --file client.docx --fields proposal.json \
        --family lessons-learned --owner you@co.com --created 2026-07-13 \
        --source-terms "Acme Corp,PRJ-1935,2024/03/10,Fragmentation"
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

import common as C

DATE_RE = re.compile(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b|\b\d{4}-\d{2}-\d{2}\b"
                     r"|\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4}\b",
                     re.I)
MONEY_RE = re.compile(r"[$£€]\s?\d")
LABEL_RE = re.compile(r"^\s*([A-Za-z][A-Za-z /&-]{1,30}?)\s*[:\-–]\s*(.+\S)\s*$")
QUARTER_RE = re.compile(r"\bQ[1-4]\b", re.I)

# Paragraphs longer than this are treated as prose/boilerplate by default (keep=fixed).
LONG_TEXT = 140


def load(path: Path):
    fmt = C.detect_format(path)
    if fmt == "docx":
        from docx import Document
        doc = Document(str(path))
        return fmt, doc, list(C.iter_docx_paragraphs(doc))
    from pptx import Presentation
    prs = Presentation(str(path))
    return fmt, prs, list(C.iter_pptx_paragraphs(prs))


def suggest_name(value: str, label: str | None, seen: set[str]) -> str:
    """Best-effort snake_case field name from a label, else the value's shape."""
    if label:
        base = C.slugify(label)
    elif DATE_RE.search(value):
        base = "date"
    elif MONEY_RE.search(value):
        base = "amount"
    elif QUARTER_RE.search(value) and len(value) <= 12:
        base = "period"
    elif len(value) <= 40:
        base = C.slugify(value)[:30] or "field"
    else:
        base = "field"
    name, i = base, 2
    while name in seen:
        name, i = f"{base}_{i}", i + 1
    seen.add(name)
    return name


def guess_keep_and_type(value: str) -> tuple[str, str]:
    """Heuristic default: is this variable content, and is it text or a list?"""
    if len(value) > LONG_TEXT:
        return "fixed", "text"       # prose/boilerplate — flip to variable if needed
    return "variable", "text"


# ── Row-groups (repeating table rows — the core of the family model) ─────────────
# A row-group makes ONE body row of a table a repeating template row: at fill time the
# engine clones that <w:tr> once per data item so the row count varies per project
# (team members, deliverables, cost lines…). `propose` offers each data-table as a
# candidate; the assisted review confirms which tables repeat, names the group and its
# columns, and which rows are template vs. dropped; `build` tags the template row's
# cells and records the group in the manifest. docx only (see SKILL.md limitations).

def detect_row_groups(doc) -> list[dict]:
    """Offer each multi-row/-column body table as a candidate row-group for review."""
    cands = []
    for idx, tbl in enumerate(C.all_docx_tables(doc)):
        rows = tbl.rows
        if len(rows) < 2:
            continue
        header = [c.text.strip() for c in rows[0].cells]
        if len(header) < 2:
            continue
        # snake_case column name per header cell, made unique
        columns, used = [], set()
        for i, h in enumerate(header):
            base = C.slugify(h) or f"col_{i + 1}"
            name, k = base, 2
            while name in used:
                name, k = f"{base}_{k}", k + 1
            used.add(name)
            columns.append(name)
        cands.append({
            "table_index": idx,
            "header": header,
            "sample_row": [c.text.strip() for c in rows[1].cells],
            "n_rows": len(rows),
            "suggest_name": C.slugify(header[0]) or f"table_{idx}",
            "suggest_columns": columns,
            "template_row_index": 1,          # first body row becomes the repeating row
            "drop_rows": list(range(2, len(rows))),  # surplus example rows to delete
            "keep": "fixed",                  # set 'variable' to activate as a row-group
        })
    return cands


def _diag_idx(part: str) -> str:
    """'ppt/diagrams/data2.xml' -> '2' (pairs a diagram with its drawing cache)."""
    m = re.search(r"(?:data|drawing)(\d*)\.xml$", part or "")
    return m.group(1) if m else ""


def _tag_cell(cell, tag: str) -> None:
    """Replace a cell's content with a single {{ tag }}, preserving the first run's
    formatting (font/size/bold) and REMOVING extra paragraphs — so no empty bullet or
    blank line survives in the repeating row."""
    paras = cell.paragraphs
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    runs = paras[0].runs
    if runs:
        runs[0].text = tag
        for r in runs[1:]:
            r.text = ""
    else:
        paras[0].add_run(tag)


def apply_row_groups(doc, proposal: dict) -> list[dict]:
    """For each confirmed (keep='variable') row-group candidate: tag its template row's
    cells with {{ column }}, delete the surplus example rows, and return the manifest
    `row_groups` list. fill.expand_row_groups later finds the template row by its
    {{ columns[0] }} tag and clones it per data item."""
    tables = C.all_docx_tables(doc)
    groups = []
    for rg in proposal.get("row_group_candidates", []):
        if rg.get("keep") != "variable":
            continue
        ti = rg.get("table_index")
        if ti is None or ti >= len(tables):
            continue
        columns = [C.slugify(c) for c in (rg.get("columns") or rg.get("suggest_columns") or [])]
        if not columns:
            continue
        name = C.slugify(rg.get("name") or rg.get("suggest_name") or f"table_{ti}")
        rows = list(tables[ti].rows)
        tri = rg.get("template_row_index", 1)
        if tri >= len(rows):
            continue
        for i, col in enumerate(columns):
            if i < len(rows[tri].cells):
                _tag_cell(rows[tri].cells[i], C.placeholder(col))
        for di in sorted(set(rg.get("drop_rows", [])), reverse=True):
            if di != tri and 0 <= di < len(rows):
                rows[di]._tr.getparent().remove(rows[di]._tr)
        groups.append({"name": name, "columns": columns})
    return groups


def propose(args):
    path = Path(args.file)
    if not path.exists():
        sys.exit(f"file not found: {path}")
    fmt, _doc, paragraphs = load(path)

    # Dedupe by exact value (replace-all semantics): the same client name in 5
    # places becomes ONE field that fills everywhere.
    by_value: dict[str, dict] = {}
    seen_names: set[str] = set()
    for para in paragraphs:
        text = C.para_text(para).strip()
        if not text:
            continue
        m = LABEL_RE.match(text)
        # For "Label: value" lines, the *value* is the candidate; keep the label fixed.
        value, label = (m.group(2).strip(), m.group(1).strip()) if m else (text, None)
        if not value or C.PLACEHOLDER_RE.search(value):
            continue
        entry = by_value.get(value)
        if entry:
            entry["occurrences"] += 1
            if label and label not in entry["labels"]:
                entry["labels"].append(label)
            continue
        keep, ftype = guess_keep_and_type(value)
        by_value[value] = {
            "current_text": value,
            "occurrences": 1,
            "labels": [label] if label else [],
            "context": text if text != value else "",
            "suggest_name": suggest_name(value, label, seen_names),
            "suggest_type": ftype,
            "keep": keep,
        }

    # Cover-page / data-bound property leaves (docProps + customXml). These render
    # the cover title/subtitle/date/author and are NOT in the body runs, so add them
    # as candidates or the cover silently keeps the source client's values.
    for leaf in C.iter_property_leaves(path):
        value = leaf["current_text"]
        if by_value.get(value.strip()):
            by_value[value.strip()]["labels"].append(leaf["label"])  # also in body
            continue
        name = leaf["suggest_name"]
        while name in seen_names:
            name += "_2"
        seen_names.add(name)
        keep, ftype = guess_keep_and_type(value.strip())
        by_value[f"\x00prop:{leaf['part']}:{leaf['localname']}"] = {
            "current_text": value,
            "occurrences": 1,
            "labels": [leaf["label"]],
            "context": "",
            "suggest_name": name,
            "suggest_type": ftype,
            "keep": keep,
            "source": "property",
        }

    # Image / logo slots (embedded pictures). One entry per media part; the user marks
    # keep='variable' for those a client should be able to swap (cover bg, logo). Default
    # fixed — most images are furniture that stays.
    image_candidates = []
    for slot in C.iter_image_slots(path):
        name = slot["suggest_name"]
        while name in seen_names:
            name += "_2"
        seen_names.add(name)
        dims = f"{slot['width']}x{slot['height']}px" if slot["width"] else "?"
        image_candidates.append({
            "current_text": f"[image {slot['media_part']} {dims}, {slot['ext']}, "
                            f"used {slot['refs']}×]",
            "occurrences": slot["refs"],
            "labels": [],
            "context": "",
            "suggest_name": name,
            "suggest_type": "image",
            "keep": "fixed",
            "source": "image",
            "media_part": slot["media_part"],
        })

    # SmartArt (diagram) TEXT candidates — fillable for FIXED-structure diagrams via a
    # substring replacement in the diagram parts (data+drawing), so they length-sort with
    # the other text. Default fixed; the review activates the ones that vary per project.
    smartart_text_candidates = []
    smartart_by_part: dict[str, list] = {}
    smartart_seen: set[tuple] = set()
    for st in C.smartart_texts(path):
        smartart_by_part.setdefault(st["part"], []).append(st["text"])
        key = (st["part"], st["text"])
        if key in smartart_seen or not st["text"].strip():
            continue
        smartart_seen.add(key)
        smartart_text_candidates.append({
            "current_text": st["text"],
            "occurrences": 1,
            "labels": [],
            "context": f"SmartArt text in {st['part']}",
            "suggest_name": suggest_name(st["text"].strip(), None, seen_names),
            "suggest_type": "text",
            "keep": "fixed",
            "source": "smartart",
            "smartart_part": st["part"],
        })
    # SmartArt -> IMAGE placeholder candidates: one per diagram. Set keep='variable' to
    # abstract the WHOLE graphic to a swappable optional image slot (works regardless of
    # node count; loses the live graphic). Use INSTEAD of the text candidates above for a
    # variable-count diagram (team/deliverables/timeline).
    smartart_image_candidates = []
    for i, (part, texts) in enumerate(sorted(smartart_by_part.items()), 1):
        nm = f"figure_{i}"
        while nm in seen_names:
            nm += "_2"
        seen_names.add(nm)
        sample = "; ".join(t.strip() for t in texts[:4] if t.strip())[:80]
        smartart_image_candidates.append({
            "current_text": f"[SmartArt {part}: {sample}]",
            "occurrences": 1,
            "labels": [],
            "context": "",
            "suggest_name": nm,
            "suggest_type": "image",
            "keep": "fixed",
            "source": "smartart_image",
            "smartart_part": part,
        })

    candidates = sorted([*by_value.values(), *smartart_text_candidates],
                        key=lambda c: (-len(c["current_text"])))
    candidates += image_candidates              # images last; not length-sorted with text
    candidates += smartart_image_candidates     # SmartArt->image conversions last too
    unsupported = C.iter_unsupported_objects(path)
    # Row-group candidates: variable-count table rows (docx only). The review step marks
    # which tables actually repeat; a repeating table should be a row-group, NOT a set of
    # per-cell text fields (mark its cell values keep='remove'/'fixed' to avoid overlap).
    row_group_candidates = detect_row_groups(_doc) if fmt == "docx" else []
    proposal = {
        "format": fmt,
        "source_file": str(path),
        "unsupported_objects": unsupported,
        "instructions": (
            "Review each candidate and set keep to one of: 'variable' (filled each "
            "time), 'fixed' (boilerplate that stays), or 'remove' (delete this line — "
            "use for surplus example bullets/rows a single 'list' field will "
            "regenerate). For a repeating list, mark ONE representative line "
            "keep='variable' + suggest_type='list', and mark the other example lines "
            "keep='remove'. Rename suggest_name to a clean snake_case field. Delete "
            "candidates you don't care about. For a table whose ROW COUNT varies per "
            "project (team, deliverables, costs), set keep='variable' on its "
            "row_group_candidate (rename name/columns; adjust template_row_index & "
            "drop_rows), and set its cell values keep='remove' in candidates above. "
            "Then: templatize.py build --file <file> --fields <this> --family <family> "
            "[--source-terms 'ProjectName,CODE,2024/01/01,Domain Term']."
        ),
        "candidates": candidates,
        "row_group_candidates": row_group_candidates,
    }
    out = Path(args.out)
    out.write_text(json.dumps(proposal, indent=2, ensure_ascii=False), encoding="utf-8")
    n_var = sum(1 for c in candidates if c["keep"] == "variable")
    print(f"Proposed {len(candidates)} candidates ({n_var} variable) -> {out}")
    if row_group_candidates:
        print(f"Detected {len(row_group_candidates)} table(s) as row-group candidates "
              f"(variable row counts) — confirm which repeat in the review step:")
        for rg in row_group_candidates:
            print(f"  - table[{rg['table_index']}] {rg['n_rows']} rows, "
                  f"columns {rg['suggest_columns']}")
    if unsupported:
        n_sa = sum(1 for u in unsupported if u["kind"] == "smartart")
        print(f"WARNING: {len(unsupported)} unsupported object(s). Charts are not fillable. "
              f"For the {n_sa} SmartArt diagram(s) you now have options in the candidates: "
              f"(A) set its TEXT candidates keep='variable' to fill in place — ONLY if its "
              f"node count is fixed across projects; (B) set its smartart_image candidate "
              f"keep='variable' to abstract the whole graphic to a swappable image "
              f"placeholder; or rebuild variable-count diagrams as a native table (row-group):")
        for u in unsupported:
            print(f"  - {u['kind']}: {u['part']} ({u['chars']} chars) e.g. {u['sample']}")
    print("Review/edit it, then run `templatize.py build`.")


def build(args):
    path = Path(args.file)
    if not path.exists():
        sys.exit(f"file not found: {path}")
    proposal = json.loads(Path(args.fields).read_text(encoding="utf-8"))
    variables = [c for c in proposal["candidates"] if c.get("keep") == "variable"]
    removals = [c for c in proposal["candidates"] if c.get("keep") == "remove"]
    if not variables:
        sys.exit("No candidates marked keep='variable' — nothing to templatize.")

    fmt, doc, paragraphs = load(path)
    if fmt != proposal.get("format"):
        sys.exit(f"Proposal format '{proposal.get('format')}' != file format '{fmt}'.")

    # Row-groups FIRST (docx): tag each confirmed table's template row and drop its
    # surplus example rows, so the paragraph passes below run over the FINAL table
    # structure (and stale rows/cells can't linger as leftover text or residue).
    row_groups = []
    if fmt == "docx":
        row_groups = apply_row_groups(doc, proposal)
        if row_groups:
            paragraphs = list(C.iter_docx_paragraphs(doc))

    # Remove extra example paragraphs first (e.g. surplus bullets/rows a `list` field
    # will regenerate). Match on exact stripped text so we only drop intended lines.
    for c in removals:
        target = c["current_text"].strip()
        for para in list(paragraphs):
            if C.para_text(para).strip() == target:
                para._p.getparent().remove(para._p)
    paragraphs = [p for p in paragraphs if p._p.getparent() is not None]

    # Longest first so a longer value ("Q3 2026") is tagged before a shorter substring
    # ("Q3") could clobber part of it.
    variables.sort(key=lambda c: -len(c["current_text"]))

    fields, warnings = [], []
    body_counts: dict[str, int] = {}
    prop_pairs: list[tuple[str, str]] = []
    smartart_pairs: list[tuple[str, str]] = []      # (current_text, tag) — text-fill (A)
    sa_image_targets: dict[str, str] = {}           # diagram_index -> field name (B)
    for c in variables:
        name = C.slugify(c.get("suggest_name") or c["current_text"])
        # Image slots are identified by media part, not by tags — no text replacement.
        if c.get("source") == "image":
            fields.append(C.Field(
                name=name, type="image", example=c["current_text"],
                guidance="Path to a replacement image; re-encoded to the slot's format, "
                         "geometry preserved. Leave unset to keep the original.",
                required=False, media_part=c.get("media_part", ""),
            ))
            continue
        # SmartArt TEXT (A): tag <a:t> in the diagram data+drawing parts, not body runs.
        if c.get("source") == "smartart":
            tag = C.placeholder(name)
            smartart_pairs.append((c["current_text"], tag))
            fields.append(C.Field(
                name=name, type="text", example=c["current_text"],
                guidance="SmartArt text (fixed-structure diagram).", required=True))
            continue
        # SmartArt -> IMAGE placeholder (B): the whole diagram becomes an image slot; the
        # graphicFrame is converted after the paragraph passes (needs the live prs).
        if c.get("source") == "smartart_image":
            sa_image_targets[_diag_idx(c.get("smartart_part", ""))] = name
            continue
        tag = C.placeholder(name)
        replaced = 0
        for para in paragraphs:
            replaced += C.replace_in_paragraph(para, c["current_text"], tag)
        body_counts[tag] = replaced
        prop_pairs.append((c["current_text"], tag))   # also try the property parts
        fields.append(C.Field(
            name=name,
            type=c.get("suggest_type", "text"),
            example=c["current_text"],
            guidance=(f"Fills: {', '.join(c.get('labels', []))}" if c.get("labels") else ""),
            required=True,
        ))

    # Family model (default) vs. per-client instance. A family template lives under
    # registry/_families/<family>/ and is what future files of that family converge to.
    if args.family:
        tdir = C.family_dir(args.family)
        client_id, doc_type_id = "_families", C.slugify(args.family)
        template_id = f"_families/{doc_type_id}"
    elif args.client and args.doc_type:
        tdir = C.template_dir(args.client, args.doc_type)
        client_id, doc_type_id = C.slugify(args.client), C.slugify(args.doc_type)
        template_id = f"{client_id}/{doc_type_id}"
    else:
        sys.exit("Provide --family <name> (the family model), or --client and --doc-type.")
    tdir.mkdir(parents=True, exist_ok=True)
    template_file = f"template.{fmt}"

    # B: abstract flagged SmartArt diagrams to placeholder images BEFORE saving (needs the
    # live prs). Each converted graphicFrame becomes a same-geometry placeholder picture
    # registered as an OPTIONAL image field, so a fill can swap it or leave the box.
    if fmt == "pptx" and sa_image_targets:
        import tempfile as _tf
        _pngdir = Path(_tf.mkdtemp())

        def _png_for(idx, w, h):
            p = _pngdir / f"sa_{idx}.png"
            C.make_placeholder_image(p, sa_image_targets.get(idx, f"figure {idx}"), w, h)
            return str(p)

        for sw in C.smartart_to_placeholder(doc, set(sa_image_targets), _png_for):
            idx = sw["index"]
            fields.append(C.Field(
                name=sa_image_targets[idx], type="image",
                example=f"[SmartArt {idx} abstracted to placeholder]",
                guidance="Optional replacement image; leave unset to keep the placeholder box.",
                required=False, media_part=sw["media_part"]))

    doc.save(str(tdir / template_file))

    # A: fill SmartArt TEXT tags into the diagram data+drawing parts (keeps render + data
    # model in sync). Blank the text of any imaged (now-orphaned) diagram so the source's
    # SmartArt content can't survive the residue check.
    sa_stats: dict[str, int] = {}
    if smartart_pairs:
        C.patch_smartart_parts(tdir / template_file, tdir / template_file,
                               C.ordered_replacer(smartart_pairs, sa_stats))
    if fmt == "pptx" and sa_image_targets:
        C.patch_smartart_parts(tdir / template_file, tdir / template_file,
                               (lambda t: ""), only_parts=set(sa_image_targets))

    # Inject tags into cover-page / data-bound property parts too (docProps +
    # customXml). Longest value first; count matches so we can warn only when a
    # field matched NEITHER the body NOR a property NOR a SmartArt part.
    prop_stats: dict[str, int] = {}
    C.patch_property_parts(tdir / template_file, tdir / template_file,
                           C.ordered_replacer(prop_pairs, prop_stats))
    for c in variables:
        if c.get("source") in ("image", "smartart_image"):
            continue   # keyed by media part, not by text tags
        name = C.slugify(c.get("suggest_name") or c["current_text"])
        tag = C.placeholder(name)
        if body_counts.get(tag, 0) == 0 and prop_stats.get(tag, 0) == 0 and sa_stats.get(tag, 0) == 0:
            warnings.append(f"'{c['current_text'][:40]}' -> {{{{ {name} }}}}: 0 matches "
                            "(text may span formatting boundaries oddly; check manually)")

    source_terms = [t.strip() for t in (args.source_terms or "").split(",") if t.strip()]
    manifest = C.Manifest(
        template_id=template_id,
        client=client_id,
        doc_type=doc_type_id,
        format=fmt,
        template_file=template_file,
        source_file=Path(path).name,
        owner=args.owner,
        created=args.created,
        changelog=[f"{args.created} v1.0.0 templatized from {Path(path).name}"] if args.created else [],
        fields=fields,
        row_groups=row_groups,
        source_terms=source_terms,
    )
    C.save_manifest(manifest, tdir / "manifest.json")

    print(f"Template registered: {manifest.template_id}  ({len(fields)} fields, "
          f"{len(row_groups)} row-group(s), {fmt})")
    print(f"  {tdir / template_file}")
    print(f"  {tdir / 'manifest.json'}")
    if row_groups:
        for g in row_groups:
            print(f"  row-group '{g['name']}' columns {g['columns']}")
    if source_terms:
        print(f"  source_terms (must not survive a fill): {source_terms}")
    for w in warnings:
        print(f"  WARNING: {w}")


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    sub = ap.add_subparsers(dest="cmd", required=True)

    p = sub.add_parser("propose", help="Extract candidate variable fields for review")
    p.add_argument("--file", required=True, help="Client .docx/.pptx to learn from")
    p.add_argument("--out", default="proposal.json")
    p.set_defaults(func=propose)

    b = sub.add_parser("build", help="Inject placeholders and register the template")
    b.add_argument("--file", required=True, help="Same client file as `propose`")
    b.add_argument("--fields", required=True, help="Reviewed proposal JSON")
    b.add_argument("--family", help="Family name -> registry/_families/<family>/ "
                                    "(the family model; preferred over --client/--doc-type)")
    b.add_argument("--client", help="Per-client instance (use --family instead by default)")
    b.add_argument("--doc-type", help="Per-client instance doc type")
    b.add_argument("--source-terms", default="",
                   help="Comma-separated source-exemplar terms that must NOT survive a "
                        "fill (client/project/code/dates/domain terms; not recurring people)")
    b.add_argument("--owner", default="")
    b.add_argument("--created", default="", help="ISO date (scripts have no clock)")
    b.set_defaults(func=build)

    args = ap.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
