"""Shared helpers for the document-template engine.

The engine turns a client's existing file into a reusable *template* (same layout,
fonts, logos, styles — with the variable text swapped for `{{ placeholder }}` tags)
plus a machine-readable *manifest* describing every placeholder. A later fill step
reads the manifest, drops in content, and re-emits a finished document.

Design choices that live here:

* **Placeholder convention** — one simple form everywhere: ``{{ field_name }}``.
  No loop/branch syntax in the template. List and table expansion is driven by the
  *manifest field type* at fill time, not by tags in the document. This keeps
  templatize (which only has ONE example to learn from) from having to synthesize
  control-flow, and makes DOCX and PPTX behave identically.

* **Run-aware, text-based replacement** — we replace the *text* of a value wherever
  it appears, preserving the formatting of the run it lives in. When a value spans
  several runs we consolidate them into the first run (keeping its formatting) so
  the tag always lands in a single run — the one failure mode that otherwise breaks
  in-place fills. Replacing every occurrence is usually what you want (a client name
  appears in many places and should all become ``{{ client_name }}``).

This module is format-agnostic where it can be: it exposes paragraph iterators for
DOCX and PPTX that both yield objects with a ``.runs`` list of runs that each carry
a settable ``.text`` — so the scan/replace/render logic is written once.
"""
from __future__ import annotations

import json
import os
import re
from dataclasses import dataclass, asdict, field as dc_field
from pathlib import Path

# ── Paths ──────────────────────────────────────────────────────────────────────
# scripts/ -> skill root -> registry/ is the version-controlled template gallery.
# Override with $TEMPLATE_REGISTRY to point at a shared gallery outside the repo
# (or to isolate tests).
SKILL_ROOT = Path(__file__).resolve().parents[1]
REGISTRY = Path(os.environ.get("TEMPLATE_REGISTRY", SKILL_ROOT / "registry"))

# ── Placeholder convention ───────────────────────────────────────────────────--
# {{ field_name }} — optional surrounding whitespace; names are [a-z0-9_].
PLACEHOLDER_RE = re.compile(r"\{\{\s*([a-zA-Z0-9_]+)\s*\}\}")
# Any leftover tag (used by the validator to detect an unfilled/partial fill).
ANY_TAG_RE = re.compile(r"\{\{.*?\}\}|\{%.*?%\}")

SUPPORTED_FORMATS = ("docx", "pptx")


def detect_format(path: str | Path) -> str:
    """Map a file path to a supported engine format by extension."""
    ext = Path(path).suffix.lower().lstrip(".")
    if ext not in SUPPORTED_FORMATS:
        raise ValueError(
            f"Unsupported format '.{ext}'. This engine templatizes {SUPPORTED_FORMATS}. "
            "PDF is an export target (fill a DOCX/PPTX then export), not a templatize source."
        )
    return ext


def slugify(text: str) -> str:
    """A filesystem/identifier-safe kebab/snake slug (lowercase, _-separated words)."""
    text = re.sub(r"[^\w\s-]", "", str(text).strip().lower())
    text = re.sub(r"[\s-]+", "_", text)
    return text.strip("_") or "unnamed"


def placeholder(name: str) -> str:
    """Render the canonical tag string for a field name."""
    return "{{ " + name + " }}"


# ── Manifest model ───────────────────────────────────────────────────────────--
@dataclass
class Field:
    """One variable slot in a template."""
    name: str                       # snake_case identifier, e.g. "client_name"
    type: str = "text"              # text | list | image
    example: str = ""               # the value seen in the source (a fill sample)
    guidance: str = ""              # how to fill it, for humans/agents
    required: bool = True
    media_part: str = ""            # for type=image: the package media part to swap

    def to_dict(self) -> dict:
        return asdict(self)


@dataclass
class Manifest:
    """Everything a future fill needs, without re-reading the whole document."""
    template_id: str                # "<client>/<doc_type>"
    client: str
    doc_type: str
    format: str                     # docx | pptx
    template_file: str              # relative name inside the template dir
    source_file: str = ""           # original client file this was learned from
    version: str = "1.0.0"
    owner: str = ""
    created: str = ""               # ISO date; caller stamps it (no clock in scripts)
    changelog: list = dc_field(default_factory=list)
    fields: list = dc_field(default_factory=list)  # list[Field]
    row_groups: list = dc_field(default_factory=list)  # repeating table rows: {name, columns:[field,...]}
    source_terms: list = dc_field(default_factory=list)  # source-exemplar terms that must not survive a fill

    def to_dict(self) -> dict:
        d = asdict(self)
        d["fields"] = [f.to_dict() if isinstance(f, Field) else f for f in self.fields]
        return d

    def field_names(self) -> list[str]:
        return [f["name"] if isinstance(f, dict) else f.name for f in self.fields]


def save_manifest(manifest: Manifest, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(manifest.to_dict(), indent=2, ensure_ascii=False),
                    encoding="utf-8")


def load_manifest(path: Path) -> dict:
    return json.loads(Path(path).read_text(encoding="utf-8"))


# ── Registry (the template gallery) ─────────────────────────────────────────--
def template_dir(client: str, doc_type: str) -> Path:
    """registry/<client>/<doc_type>/ — where a template + its manifest live.

    The reserved namespaces `_builtin` (shipped generic templates) and `_families`
    (governed canonicals) pass through unslugified so `--client _builtin` resolves.
    """
    ns = client if client in ("_builtin", "_families") else slugify(client)
    return REGISTRY / ns / slugify(doc_type)


def family_dir(family: str) -> Path:
    """registry/_families/<family>/ — the GOVERNED canonical template for a document
    family (Lessons Learned, Change Note, …). This is the default home for a template
    in the family model; per-client dirs are the exception, not the rule."""
    return REGISTRY / "_families" / slugify(family)


def find_template(client: str, doc_type: str) -> tuple[Path, dict]:
    """Resolve a registered template. Returns (template_path, manifest_dict)."""
    d = template_dir(client, doc_type)
    man_path = d / "manifest.json"
    if not man_path.exists():
        raise FileNotFoundError(
            f"No template registered for client='{client}', doc_type='{doc_type}'. "
            f"Expected {man_path}. Run registry.py list to see what exists."
        )
    manifest = load_manifest(man_path)
    tpl_path = d / manifest["template_file"]
    if not tpl_path.exists():
        raise FileNotFoundError(f"Manifest references missing template file: {tpl_path}")
    return tpl_path, manifest


# ── Format-agnostic paragraph iteration ─────────────────────────────────────--
# Both python-docx and python-pptx paragraphs expose `.runs` where each run has a
# settable `.text`. We yield paragraph objects so scan/replace/render is written once.

def _iter_docx_elem(element, parent):
    """Yield paragraphs under an lxml element (body or table cell), in document
    order, recursing into nested tables and content controls (w:sdt). Cover-page
    building blocks wrap the title/subtitle in block-level content controls, so
    without the sdt recursion that text is invisible to the fill."""
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.ns import qn

    for child in element.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        elif child.tag == qn("w:tbl"):
            table = Table(child, parent)
            for row in table.rows:
                for cell in row.cells:
                    yield from _iter_docx_elem(cell._element, cell)
        elif child.tag == qn("w:sdt"):
            content = child.find(qn("w:sdtContent"))
            if content is not None:
                yield from _iter_docx_elem(content, parent)


def _iter_textbox_paras(root, parent):
    """Yield paragraphs inside text boxes / drawing canvases (w:txbxContent) anywhere
    under `root`. Cover subtitles, callouts and sidebars often live in floating text
    boxes, which the normal body walk never reaches. A DrawingML text box is usually
    duplicated in an mc:AlternateContent (a modern Choice + a VML Fallback), so the
    same text appears in two w:txbxContent — we yield both so a fill stays consistent."""
    from docx.text.paragraph import Paragraph
    from docx.oxml.ns import qn
    for txbx in root.iter(qn("w:txbxContent")):
        for p in txbx.iterchildren(qn("w:p")):
            yield Paragraph(p, parent)


def iter_docx_paragraphs(doc):
    """Every paragraph in a .docx: body, tables (recursive), content controls, text
    boxes, and headers/footers."""
    yield from _iter_docx_elem(doc.element.body, doc._body)
    yield from _iter_textbox_paras(doc.element.body, doc._body)
    for section in doc.sections:
        for hf in (section.header, section.footer,
                   section.first_page_header, section.first_page_footer,
                   section.even_page_header, section.even_page_footer):
            for para in hf.paragraphs:
                yield para
            for table in hf.tables:
                for row in table.rows:
                    for cell in row.cells:
                        yield from _iter_docx_elem(cell._element, cell)
            yield from _iter_textbox_paras(hf._element, hf)


def iter_pptx_paragraphs(prs):
    """Every text paragraph in a .pptx across slides, shapes (incl. groups) & tables."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk_shapes(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk_shapes(shape.shapes)
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    yield para
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            yield para

    for slide in prs.slides:
        yield from walk_shapes(slide.shapes)
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                yield para


def para_text(paragraph) -> str:
    """Concatenated text of a paragraph's runs (works for docx and pptx)."""
    return "".join(run.text for run in paragraph.runs)


def all_docx_tables(doc):
    """Every table in the body, in document order, recursing into nested tables.

    Row-groups (repeating rows) live in body content tables. This is the ONE
    enumeration `templatize.propose` (which assigns a table_index), `templatize.build`
    (which tags that index's template row), and `fill.expand_row_groups` all share, so
    a row-group's table_index means the same thing at propose, build and fill time.
    Header/footer tables are intentionally excluded — they are not row-group carriers."""
    from docx.table import Table
    out = []

    def walk(tables):
        for t in tables:
            out.append(t)
            for row in t.rows:
                for cell in row.cells:
                    walk(cell.tables)
    walk(doc.tables)
    return out


# ── Document properties (cover pages & data-bound content controls) ──────────--
# Word "cover page" building blocks render their title / subtitle / date / author
# from CONTENT CONTROLS that are *data-bound* to package parts, not from body runs.
# So the value shown on the cover lives in docProps/core.xml (Dublin Core) or in a
# customXml CoverPageProperties store — NOT in word/document.xml. The run-based
# scan/replace above never reaches it, which is why a fill can update the body but
# leave the cover stale (and validate.py, which only reads runs, misses it).
#
# These helpers add the *same* {{ tag }} machinery to those parts: propose can offer
# the property leaves as candidates, build injects tags into them, fill replaces the
# tags, and validate scans them for leftovers. Editing the bound value updates the
# cover the next time the file is rendered/opened.

# core.xml leaves that carry human content (skip revision/dates/lastModifiedBy/etc.)
_CORE_CONTENT_LOCALNAMES = {
    "title", "subject", "creator", "description", "keywords", "category",
}
_COVERPROPS_NS = "http://schemas.microsoft.com/office/2006/coverPageProps"

# A clean field-name + label hint per known property leaf (best-effort; the assisted
# review step can rename). Keyed by local name.
_PROPERTY_HINTS = {
    "title": ("doc_title", "Cover title (document property)"),
    "subject": ("doc_subtitle", "Cover subtitle (document property)"),
    "creator": ("author", "Author (document property)"),
    "description": ("doc_description", "Description (document property)"),
    "keywords": ("keywords", "Keywords (document property)"),
    "category": ("category", "Category (document property)"),
    "PublishDate": ("publish_date", "Cover date (CoverPageProperties.PublishDate)"),
    "Abstract": ("abstract", "Cover abstract (CoverPageProperties.Abstract)"),
    "CompanyAddress": ("company_address", "Cover company address"),
    "CompanyPhone": ("company_phone", "Cover company phone"),
    "CompanyEmail": ("company_email", "Cover company email"),
}


def _local(tag) -> str:
    """Local name of an lxml qualified tag (strips the {ns})."""
    from lxml import etree
    return etree.QName(tag).localname if isinstance(tag, str) else ""


def _property_kind(part_name: str, root) -> str | None:
    """Classify a package part as a content-bearing property store, or None."""
    if part_name == "docProps/core.xml":
        return "core"
    if part_name.startswith("customXml/item") and part_name.endswith(".xml"):
        if root is not None and _local(root.tag) == "CoverPageProperties":
            return "cover"
    return None


def _content_leaves(root, kind):
    """Yield the (element) leaves of a property store that carry fillable content."""
    if kind == "core":
        for el in root:
            if _local(el.tag) in _CORE_CONTENT_LOCALNAMES:
                yield el
    else:  # cover: every leaf element (no element children)
        for el in root.iter():
            if len(el) == 0 and el is not root:
                yield el


def iter_property_leaves(path: str | Path) -> list[dict]:
    """Read a file's content-bearing property leaves (for `propose`).

    Returns a list of dicts: current_text, part, localname, suggest_name, label.
    Only docx/pptx packages that actually contain these parts yield anything.
    """
    import zipfile
    from lxml import etree

    out: list[dict] = []
    try:
        zin = zipfile.ZipFile(str(path))
    except (zipfile.BadZipFile, FileNotFoundError):
        return out
    with zin:
        for name in zin.namelist():
            if not (name == "docProps/core.xml"
                    or (name.startswith("customXml/item") and name.endswith(".xml"))):
                continue
            try:
                root = etree.fromstring(zin.read(name))
            except etree.XMLSyntaxError:
                continue
            kind = _property_kind(name, root)
            if not kind:
                continue
            for el in _content_leaves(root, kind):
                text = (el.text or "").strip()
                if not text:
                    continue
                ln = _local(el.tag)
                sug_name, label = _PROPERTY_HINTS.get(ln, (slugify(ln), f"Document property: {ln}"))
                out.append({
                    "current_text": el.text,   # keep raw (may hold newlines)
                    "part": name,
                    "localname": ln,
                    "suggest_name": sug_name,
                    "label": label,
                })
    return out


def ordered_replacer(pairs: list[tuple[str, str]], stats: dict | None = None):
    """Build a text->text transform applying substring replacements, longest `old`
    first (so a longer value is tagged before a shorter substring clobbers it).
    If `stats` is given, counts replacements per `new` string."""
    pairs = sorted((p for p in pairs if p[0]), key=lambda kv: -len(kv[0]))

    def transform(text: str) -> str:
        for old, new in pairs:
            if old and old in text:
                if stats is not None:
                    stats[new] = stats.get(new, 0) + text.count(old)
                text = text.replace(old, new)
        return text

    return transform


def patch_property_parts(path: str | Path, out_path: str | Path, transform) -> int:
    """Apply `transform` to every content-bearing property leaf and rewrite the
    package. Returns the number of leaves changed. Safe when out_path == path.
    Non-property parts are copied byte-for-byte (compression preserved)."""
    import os
    import tempfile
    import zipfile
    from lxml import etree

    path = Path(path)
    with zipfile.ZipFile(str(path)) as zin:
        entries = [(i, zin.read(i.filename)) for i in zin.infolist()]

    changed = 0
    patched: dict[str, bytes] = {}
    for info, data in entries:
        name = info.filename
        if not (name == "docProps/core.xml"
                or (name.startswith("customXml/item") and name.endswith(".xml"))):
            continue
        try:
            root = etree.fromstring(data)
        except etree.XMLSyntaxError:
            continue
        kind = _property_kind(name, root)
        if not kind:
            continue
        part_changed = False
        for el in _content_leaves(root, kind):
            if el.text:
                new = transform(el.text)
                if new != el.text:
                    el.text = new
                    part_changed = True
                    changed += 1
        if part_changed:
            patched[name] = etree.tostring(root, xml_declaration=True,
                                           encoding="UTF-8", standalone=True)

    if not patched:
        if str(out_path) != str(path):
            import shutil
            shutil.copyfile(str(path), str(out_path))
        return 0

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".zip",
                                      dir=str(Path(out_path).parent))
    tmp.close()
    with zipfile.ZipFile(tmp.name, "w") as zout:
        for info, data in entries:
            zout.writestr(info, patched.get(info.filename, data))
    os.replace(tmp.name, str(out_path))
    return changed


def property_texts(path: str | Path) -> list[str]:
    """Every content-bearing property leaf's text (for the validator's tag scan)."""
    return [d["current_text"] for d in iter_property_leaves(path)]


# ── Image / logo slots ───────────────────────────────────────────────────────--
# A client's cover background and the logo repeated on every slide/page are PICTURES,
# not text. To make them swappable per client we treat each embedded *media part*
# (word/media/imageN, ppt/media/imageN) as a named image SLOT: one media part may be
# referenced by many shapes (a logo on 30 slides), so replacing that one part swaps it
# everywhere at once — the same "preserve the frame, inject the content" idea applied
# to pictures. The shape geometry (position/size/crop) lives in the drawing XML, not
# the media, so it is untouched by a media swap.

_MEDIA_PREFIXES = ("word/media/", "ppt/media/", "xl/media/")
_PIL_FMT = {"png": "PNG", "jpg": "JPEG", "jpeg": "JPEG", "gif": "GIF",
            "bmp": "BMP", "tif": "TIFF", "tiff": "TIFF"}  # emf/wmf: vector, copied raw


def _slot_name(width, height, ext, idx) -> str:
    """Heuristic name: wide+short small raster -> logo; big -> background; else imageN."""
    if width and height:
        ar = width / height if height else 0
        if ar >= 2.0 and max(width, height) <= 900:
            return f"logo_{idx}"
        if width >= 1200 or height >= 900:
            return f"background_{idx}"
    return f"image_{idx}"


def iter_image_slots(path: str | Path) -> list[dict]:
    """Every embedded image part (for `propose`): media_part, ext, bytes, dims, refs."""
    import zipfile

    out: list[dict] = []
    try:
        z = zipfile.ZipFile(str(path))
    except (zipfile.BadZipFile, FileNotFoundError):
        return out
    with z:
        names = z.namelist()
        rels = [n for n in names if n.endswith(".rels")]
        idx = 0
        for n in names:
            if not (n.startswith(_MEDIA_PREFIXES) and not n.endswith("/")):
                continue
            idx += 1
            data = z.read(n)
            width = height = None
            try:
                import io
                from PIL import Image
                width, height = Image.open(io.BytesIO(data)).size
            except Exception:
                pass
            base = n.split("/")[-1]
            refs = sum(z.read(r).decode("utf-8", "ignore").count(base) for r in rels)
            out.append({
                "media_part": n,
                "ext": n.rsplit(".", 1)[-1].lower(),
                "bytes": len(data),
                "width": width,
                "height": height,
                "refs": refs,
                "suggest_name": _slot_name(width, height, n.rsplit(".", 1)[-1].lower(), idx),
            })
    return out


def swap_media_parts(path: str | Path, out_path: str | Path,
                     mapping: dict[str, str]) -> int:
    """Replace image SLOTS: mapping is {media_part_name: new_image_file}. The new image
    is re-encoded to the ORIGINAL part's format so every reference (and the file's
    content-type) stays valid; the shape frame scales it, so geometry is preserved.
    Returns the number of parts swapped. Safe when out_path == path."""
    import io
    import os
    import tempfile
    import zipfile

    mapping = {k: v for k, v in (mapping or {}).items() if v}
    with zipfile.ZipFile(str(path)) as zin:
        entries = [(i, zin.read(i.filename)) for i in zin.infolist()]

    new_bytes: dict[str, bytes] = {}
    for part, new_file in mapping.items():
        ext = part.rsplit(".", 1)[-1].lower()
        pil_fmt = _PIL_FMT.get(ext)
        if pil_fmt is None:                       # vector (emf/wmf) — copy raw
            with open(new_file, "rb") as f:
                new_bytes[part] = f.read()
            continue
        from PIL import Image
        img = Image.open(new_file)
        if pil_fmt == "JPEG" and img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format=pil_fmt)
        new_bytes[part] = buf.getvalue()

    if not new_bytes:
        if str(out_path) != str(path):
            import shutil
            shutil.copyfile(str(path), str(out_path))
        return 0

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".zip",
                                      dir=str(Path(out_path).parent))
    tmp.close()
    with zipfile.ZipFile(tmp.name, "w") as zout:
        for info, data in entries:
            zout.writestr(info, new_bytes.get(info.filename, data))
    os.replace(tmp.name, str(out_path))
    return len(new_bytes)


def iter_unsupported_objects(path: str | Path) -> list[dict]:
    """Find content the engine CANNOT fill: SmartArt/diagram text and native charts.

    Text inside a SmartArt/diagram (``diagrams/data*.xml``) or a chart
    (``charts/chart*.xml``) is not a normal text run — python-docx/pptx never see it,
    so the fill leaves the source's content in place. Rather than silently ship stale
    text, `propose` surfaces these and QA flags them. Returns dicts:
    {kind, part, chars, sample} — kind in {'smartart','chart'}."""
    import re as _re
    import zipfile

    out: list[dict] = []
    try:
        z = zipfile.ZipFile(str(path))
    except (zipfile.BadZipFile, FileNotFoundError):
        return out
    with z:
        for n in z.namelist():
            kind = None
            if _re.search(r"diagrams/data\d*\.xml$", n):
                kind = "smartart"
            elif _re.search(r"charts/chart\d*\.xml$", n):
                kind = "chart"
            if not kind:
                continue
            xml = z.read(n).decode("utf-8", "ignore")
            texts = [t for t in _re.findall(r"<a:t>(.*?)</a:t>", xml) if t.strip()]
            # A SmartArt diagram that has been abstracted to an image (or fully blanked)
            # carries no text — it is no longer an unfilled object, so don't flag it.
            if kind == "smartart" and not texts:
                continue
            out.append({
                "kind": kind,
                "part": n,
                "chars": sum(len(t) for t in texts),
                "sample": "; ".join(texts[:5])[:120],
            })
    return out


# ── SmartArt (diagram) text — fillable for FIXED-structure diagrams ──────────────
# SmartArt text is NOT a normal run (python-docx/pptx can't see it): it lives in
# ``ppt|word/diagrams/data#.xml`` as <a:t> runs, mirrored in a cached ``drawing#.xml``
# that LibreOffice/PowerPoint actually RENDER. So to fill SmartArt text we rewrite the
# <a:t> content in BOTH parts (data keeps PowerPoint correct; drawing keeps the render
# correct). This works when the node COUNT is fixed across projects — it changes text,
# not structure. Variable-count SmartArt (add/remove team cards) needs the data model +
# connections regenerated and can't be done here — rebuild those as a native table
# (row-groups) or abstract them to an image placeholder (see smartart_to_placeholder).

_SMARTART_PART_RE = re.compile(r"diagrams/(?:data|drawing)\d*\.xml$")
_SMARTART_DATA_RE = re.compile(r"diagrams/data\d*\.xml$")
_AT_RE = re.compile(r"(<a:t>)(.*?)(</a:t>)", re.S)


def _xml_unescape(s: str) -> str:
    from xml.sax.saxutils import unescape
    return unescape(s)


def _xml_escape(s: str) -> str:
    from xml.sax.saxutils import escape
    return escape(s)


def smartart_texts(path: str | Path) -> list[dict]:
    """Every <a:t> text in SmartArt DATA parts: {text, part}. Used by `propose` to
    offer SmartArt text as candidates and by `validate` to scan it for leftover tags
    and source residue (python-pptx never sees this text)."""
    import zipfile
    out: list[dict] = []
    try:
        z = zipfile.ZipFile(str(path))
    except (zipfile.BadZipFile, FileNotFoundError):
        return out
    with z:
        for n in z.namelist():
            if not _SMARTART_DATA_RE.search(n):
                continue
            xml = z.read(n).decode("utf-8", "ignore")
            for m in _AT_RE.finditer(xml):
                if m.group(2).strip():
                    # Return UNESCAPED text so candidates/residue terms are natural
                    # ("A & B", not "A &amp; B"); patch re-escapes on write.
                    out.append({"text": _xml_unescape(m.group(2)), "part": n})
    return out


def patch_smartart_parts(path: str | Path, out_path: str | Path, transform,
                         only_parts: set | None = None) -> int:
    """Apply `transform` to the text inside every <a:t> in SmartArt data+drawing parts
    (optionally restricted to `only_parts` — e.g. one diagram's data+drawing), keeping
    the data model and cached drawing in sync. Returns the number of <a:t> changed.
    Safe when out_path == path."""
    import os
    import shutil
    import tempfile
    import zipfile

    path = Path(path)
    with zipfile.ZipFile(str(path)) as zin:
        entries = [(i, zin.read(i.filename)) for i in zin.infolist()]

    patched: dict[str, bytes] = {}
    changed = 0
    for info, data in entries:
        n = info.filename
        if not _SMARTART_PART_RE.search(n):
            continue
        if only_parts is not None and _diagram_index(n) not in only_parts:
            continue
        txt = data.decode("utf-8", "ignore")
        cnt = [0]

        def repl(m, _cnt=cnt):
            # Work in UNESCAPED space so a tag/value that contains &, <, > matches and
            # is re-escaped correctly — raw substitution into XML would break the part.
            inner = _xml_unescape(m.group(2))
            new = transform(inner)
            if new != inner:
                _cnt[0] += 1
            return m.group(1) + _xml_escape(new) + m.group(3)

        new_txt = _AT_RE.sub(repl, txt)
        if new_txt != txt:
            patched[n] = new_txt.encode("utf-8")
            changed += cnt[0]

    if not patched:
        if str(out_path) != str(path):
            shutil.copyfile(str(path), str(out_path))
        return 0

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".zip",
                                      dir=str(Path(out_path).parent))
    tmp.close()
    with zipfile.ZipFile(tmp.name, "w", zipfile.ZIP_DEFLATED) as zout:
        for info, data in entries:
            zout.writestr(info, patched.get(info.filename, data))
    _replace_with_retry(tmp.name, str(out_path))
    return changed


def _replace_with_retry(src: str, dst: str, attempts: int = 5) -> None:
    """os.replace, resilient to transient Windows locks (AV/search indexer briefly
    holding a just-written file). Falls back to copy+remove if replace keeps failing."""
    import os
    import shutil
    import time
    for i in range(attempts):
        try:
            os.replace(src, dst)
            return
        except PermissionError:
            if i == attempts - 1:
                shutil.copyfile(src, dst)      # last resort: overwrite in place
                try:
                    os.remove(src)
                except OSError:
                    pass
                return
            time.sleep(0.2 * (i + 1))


def _diagram_index(part_name: str) -> str:
    """'ppt/diagrams/data2.xml' -> '2'; 'drawing2.xml' -> '2'. Pairs a data part with
    its drawing cache so `only_parts={'2'}` touches both."""
    m = re.search(r"(?:data|drawing)(\d*)\.xml$", part_name)
    return m.group(1) if m else ""


def make_placeholder_image(out_png: str | Path, label: str, w_px: int, h_px: int) -> None:
    """A neutral labelled placeholder box (for abstracting a project-specific SmartArt/
    figure to a swappable image slot)."""
    from PIL import Image, ImageDraw
    w = max(240, min(int(w_px) or 800, 2400))
    h = max(140, min(int(h_px) or 450, 2400))
    img = Image.new("RGB", (w, h), (238, 240, 244))
    d = ImageDraw.Draw(img)
    d.rectangle([3, 3, w - 4, h - 4], outline=(150, 160, 175), width=3)
    msg = f"[ {label} ]"
    try:
        d.text((16, max(10, h // 2 - 8)), msg, fill=(90, 100, 115))
    except Exception:
        pass
    img.save(str(out_png), "PNG")


def smartart_to_placeholder(prs, target_indices: set, png_for) -> list[dict]:
    """Replace each SmartArt graphicFrame whose diagram index is in `target_indices`
    with a placeholder PICTURE of the same geometry, and remove the graphicFrame.
    `png_for(index, w_px, h_px)` returns a path to a placeholder image for that frame.
    Returns [{index, media_hint, name}] describing what was swapped, so the caller can
    register image fields + blank the orphaned diagram text (for the residue check).
    Operates on a python-pptx Presentation before it is saved."""
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    EMU_PER_PX = 9525
    done: list[dict] = []
    for slide in prs.slides:
        rels = slide.part.rels
        for shape in list(slide.shapes):
            el = shape._element
            if not el.tag.endswith("}graphicFrame"):
                continue
            gd = el.find(f"{{{A_NS}}}graphic/{{{A_NS}}}graphicData")
            if gd is None or gd.get("uri") != DGM_NS:
                continue
            relids = gd.find(f"{{{DGM_NS}}}relIds")
            if relids is None:
                continue
            dm = relids.get(f"{{{R_NS}}}dm")
            try:
                data_part = rels[dm].target_partname       # '/ppt/diagrams/data2.xml'
            except KeyError:
                continue
            idx = _diagram_index(str(data_part))
            if idx not in target_indices:
                continue
            L, T, W, H = shape.left, shape.top, shape.width, shape.height
            png = png_for(idx, (W or 0) // EMU_PER_PX, (H or 0) // EMU_PER_PX)
            pic = slide.shapes.add_picture(png, L, T, W, H)
            blip = pic._element.find(f".//{{{A_NS}}}blip")
            media_part = ""
            if blip is not None:
                embed = blip.get(f"{{{R_NS}}}embed")
                try:
                    media_part = str(rels[embed].target_partname).lstrip("/")
                except KeyError:
                    media_part = ""
            el.getparent().remove(el)
            done.append({"index": idx, "media_part": media_part})
    return done


def replace_in_paragraph(paragraph, old: str, new: str) -> int:
    """Replace every occurrence of `old` with `new` inside one paragraph.

    Preserves run formatting. If `old` sits within a single run we edit that run in
    place. If it spans runs we consolidate the paragraph's runs into the first
    (keeping the first run's formatting) so `new` — critically, a ``{{ tag }}`` —
    lands in exactly one run. Returns the number of replacements made.
    """
    if not old:
        return 0
    runs = paragraph.runs
    if not runs:
        return 0

    # Fast path: value contained within a single run.
    n = 0
    for run in runs:
        if old in run.text:
            count = run.text.count(old)
            run.text = run.text.replace(old, new)
            n += count
    if n:
        return n

    # Slow path: value spans runs -> consolidate into the first run, then replace.
    full = "".join(r.text for r in runs)
    if old not in full:
        return 0
    count = full.count(old)
    runs[0].text = full.replace(old, new)
    for r in runs[1:]:
        r.text = ""
    return count
